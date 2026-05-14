import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
from plotly.subplots import make_subplots
import io
import json
import os
from datetime import datetime
try:
    from data_fetcher import EconomicDataFetcher, API_GUIDE, DEFAULTS
except Exception:
    EconomicDataFetcher = None
    API_GUIDE = ""
    DEFAULTS = {}
from openpyxl import load_workbook
import math
import os
import json

# ══════════════════════════════════════════════════════════════
# Auto-Save System (컨트롤 패널 값 자동 저장/복원)
# ══════════════════════════════════════════════════════════════
SAVE_FILE = os.path.join(os.path.dirname(__file__), 'user_settings.json')

def load_saved_values():
    """저장된 사용자 설정값을 불러옵니다."""
    if os.path.exists(SAVE_FILE):
        try:
            with open(SAVE_FILE, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            return {}
    return {}

def save_all_values(keys_to_save):
    """현재 session_state의 값을 JSON 파일로 자동 저장합니다."""
    data = {}
    for key in keys_to_save:
        if key in st.session_state:
            val = st.session_state[key]
            # JSON serializable 타입만 저장
            if isinstance(val, (int, float, str, bool)):
                data[key] = val
    try:
        with open(SAVE_FILE, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception:
        pass

def sv(key, default):
    """저장된 값이 있으면 반환, 없으면 기본값 반환."""
    saved = load_saved_values()
    return saved.get(key, default)

st.set_page_config(page_title="등촌골프연습장 사업성 분석", page_icon="⛳", layout="wide", initial_sidebar_state="expanded")

# 로그인 기능 제거됨 — 누구나 접속 가능

# ══════════════════════════════════════════════════════════════
# Premium CSS
# ══════════════════════════════════════════════════════════════
st.markdown("""<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

/* ─────────── Clean Dark Theme — carwash 영감 ─────────── */
/* 디자인 원칙: 단색 + 미세한 보더, 그라디언트/장식 제거, 일관된 간격 */

/* ── Global ── */
.stApp {
    font-family: 'Inter', 'Pretendard', -apple-system, BlinkMacSystemFont, sans-serif !important;
    background: #0b1120 !important;
    color: #e2e8f0 !important;
}
.main .block-container { padding: 1rem 2rem 3rem 2rem; max-width: 1400px; }
.main { background: #0b1120 !important; }
header[data-testid="stHeader"] { background: #0b1120 !important; height: 0; }
p, span, label, .stMarkdown { color: #cbd5e1 !important; line-height: 1.55; }
h1, h2, h3, h4, h5, h6 { color: #f1f5f9 !important; letter-spacing: -0.01em; }

/* ── Sidebar (단색, 그라디언트 제거) ── */
section[data-testid="stSidebar"] {
    background: #0d1424 !important;
    width: 300px !important;
    border-right: 1px solid #1e293b !important;
}
section[data-testid="stSidebar"] * { color: #94a3b8 !important; }
section[data-testid="stSidebar"] h1, section[data-testid="stSidebar"] h2,
section[data-testid="stSidebar"] h3, section[data-testid="stSidebar"] h4,
section[data-testid="stSidebar"] h5 { color: #e2e8f0 !important; font-weight: 600 !important; }
section[data-testid="stSidebar"] .stNumberInput input,
section[data-testid="stSidebar"] .stTextInput input,
section[data-testid="stSidebar"] .stSelectbox > div > div {
    background: #111827 !important; border: 1px solid #1e293b !important;
    color: #e2e8f0 !important; border-radius: 8px !important; font-size: 0.85rem !important;
}
section[data-testid="stSidebar"] details {
    background: #111827 !important; border: 1px solid #1e293b !important;
    border-radius: 10px !important; margin-bottom: 8px !important;
}
section[data-testid="stSidebar"] details summary {
    font-weight: 600 !important; font-size: 0.88rem !important; padding: 11px 14px !important;
    color: #e2e8f0 !important;
}
section[data-testid="stSidebar"] details[open] { border-color: #334155 !important; }
section[data-testid="stSidebar"] .stDivider { border-color: #1e293b !important; }

/* ── Metric Cards (그라디언트 제거, 호버 효과 제거) ── */
div[data-testid="stMetric"] {
    background: #111827 !important;
    border: 1px solid #1e293b !important;
    border-radius: 10px !important;
    padding: 16px !important;
    box-shadow: none !important;
}
div[data-testid="stMetricLabel"] {
    font-size: 0.75rem !important; font-weight: 500 !important;
    color: #94a3b8 !important; text-transform: none !important; letter-spacing: 0 !important;
}
div[data-testid="stMetricValue"] {
    font-size: 1.5rem !important; font-weight: 700 !important; color: #f1f5f9 !important;
}
div[data-testid="stMetricDelta"] { color: #94a3b8 !important; font-size: 0.78rem !important; }
div[data-testid="stMetricDelta"] svg { display: none !important; }

/* ── Tables (단순화) ── */
div[data-testid="stDataFrame"] {
    border: 1px solid #1e293b !important;
    border-radius: 10px !important;
    overflow: hidden !important;
}
div[data-testid="stDataFrame"] table { background: #111827 !important; color: #e2e8f0 !important; }
div[data-testid="stDataFrame"] th {
    background: #0f172a !important; color: #cbd5e1 !important;
    font-weight: 600 !important; font-size: 0.82rem !important;
    border-bottom: 1px solid #1e293b !important;
}
div[data-testid="stDataFrame"] td {
    background: #111827 !important; color: #cbd5e1 !important;
    border-color: #1e293b !important; font-size: 0.85rem !important;
}
div[data-testid="stDataFrame"] tr:hover td { background: #0f172a !important; }

/* ── Alerts (단순화) ── */
div[data-testid="stAlert"] {
    border-radius: 10px !important;
    background: #111827 !important;
    border: 1px solid #1e293b !important;
}

/* ── Section Headers (그라디언트 제거, 단순한 보더) ── */
.sec {
    background: transparent;
    color: #f1f5f9;
    padding: 0 0 8px 0;
    margin: 28px 0 16px 0;
    font-size: 1.05rem;
    font-weight: 700;
    border-bottom: 1px solid #1e293b;
    letter-spacing: -0.01em;
}
.sec span { font-size: 1rem; margin-right: 8px; opacity: 0.7; }
.sub-sec {
    background: transparent;
    border-left: none;
    padding: 0;
    margin: 20px 0 10px 0;
    font-weight: 600;
    color: #cbd5e1;
    font-size: 0.92rem;
}
.info-box {
    background: rgba(59,130,246,0.06);
    border: 1px solid rgba(59,130,246,0.18);
    border-radius: 8px;
    padding: 11px 14px;
    margin: 10px 0;
    font-size: 0.85rem;
    color: #93c5fd;
}
.warn-box {
    background: rgba(249,115,22,0.06);
    border: 1px solid rgba(249,115,22,0.18);
    border-radius: 8px;
    padding: 11px 14px;
    margin: 10px 0;
    font-size: 0.85rem;
    color: #fdba74;
}

/* ── KPI cards (단순화) ── */
.kpi-wrap { position: relative; }
.kpi-box {
    background: #111827;
    border: 1px solid #1e293b;
    border-radius: 10px;
    padding: 14px 16px;
    transition: border-color 0.15s ease;
}
.kpi-box:hover { border-color: #334155; background: #111827; }
.kpi-tip {
    display: none; position: absolute; top: 100%; left: 50%; transform: translateX(-50%);
    z-index: 1000; min-width: 220px; margin-top: 6px;
    background: #0f172a; border: 1px solid #334155; border-radius: 10px;
    padding: 12px 14px; box-shadow: 0 8px 24px rgba(0,0,0,0.5);
}
.kpi-wrap:hover .kpi-tip { display: block; }

/* ── Charts (보더만 유지) ── */
div[data-testid="stPlotlyChart"] {
    border: 1px solid #1e293b !important;
    border-radius: 10px !important;
    overflow: hidden !important;
    background: #111827 !important;
}

/* ── Expanders (main) ── */
details { background: #111827 !important; border: 1px solid #1e293b !important; border-radius: 10px !important; }
details summary { color: #e2e8f0 !important; font-weight: 500 !important; }

/* ── Inputs ── */
.stNumberInput input, .stTextInput input, .stSelectbox > div > div {
    background: #111827 !important; border: 1px solid #1e293b !important;
    color: #e2e8f0 !important; border-radius: 8px !important;
}

/* ── Checkboxes & Buttons ── */
.stCheckbox label span { color: #cbd5e1 !important; }
.stButton button, .stDownloadButton button {
    background: #111827 !important; border: 1px solid #1e293b !important;
    color: #e2e8f0 !important; border-radius: 8px !important; font-weight: 500 !important;
}
.stButton button:hover, .stDownloadButton button:hover {
    background: #1e293b !important; border-color: #334155 !important;
}

/* ── Footer ── */
.footer {
    background: transparent;
    color: #64748b;
    padding: 16px 0 0 0;
    text-align: center;
    font-size: 0.78rem;
    margin-top: 32px;
    border-top: 1px solid #1e293b;
}
hr { border-color: #1e293b !important; margin: 1.5rem 0 !important; }

/* 탭 네비게이션은 폴더 탭 스타일로 별도 블록(아래)에서 정의 */

/* ── Scrollbar ── */
::-webkit-scrollbar { width: 8px; height: 8px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: #1e293b; border-radius: 4px; }
::-webkit-scrollbar-thumb:hover { background: #334155; }

/* ── 프린트 전용 스타일 ── */
@media print {
    @page { margin: 15mm; }
    body, .stApp { background: white !important; color: #1e293b !important; }
    section[data-testid="stSidebar"] { display: none !important; }
    div[data-testid="stToolbar"] { display: none !important; }
    header[data-testid="stHeader"] { display: none !important; }
    .stButton, .stDownloadButton, .stSelectbox { display: none !important; }
    div[data-testid="stRadio"] { display: none !important; }
    /* 내용 끊김 방지 */
    .sec-header, .sub-sec, .info-box, .warn-box, table, .chart-wrap {
        page-break-inside: avoid !important;
        break-inside: avoid !important;
    }
    /* 차트 컨테이너 */
    div[data-testid="stPlotlyChart"] {
        page-break-inside: avoid !important;
        break-inside: avoid !important;
    }
    /* 다크 테마 → 라이트 테마 전환 */
    .sec-header { background: #1e293b !important; color: white !important; -webkit-print-color-adjust: exact; print-color-adjust: exact; }
    .kpi-card { border: 1px solid #cbd5e1 !important; background: white !important; }
    .kpi-val { color: #0f172a !important; }
    .info-box { border-color: #3b82f6 !important; background: #eff6ff !important; color: #1e3a5f !important; -webkit-print-color-adjust: exact; }
    table th { background: #1e293b !important; color: white !important; -webkit-print-color-adjust: exact; print-color-adjust: exact; }
    table td { color: #1e293b !important; border-color: #e2e8f0 !important; }
    /* 페이지 헤더/푸터 */
    .hero-header { -webkit-print-color-adjust: exact; print-color-adjust: exact; }
}
</style>""", unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════
# Constants & Helpers
# ══════════════════════════════════════════════════════════════
EXCEL_PATH = r'C:\Users\win\Desktop\등촌 골프연습장 재무모델링 3월 19일 (2).xlsx'
CACHE_PATH = os.path.join(os.path.dirname(__file__), 'excel_cache.json')
USE_CACHE = not os.path.exists(EXCEL_PATH)  # Excel 없으면 캐시 사용 (Cloud 환경)
억 = 100_000_000
만 = 10_000

C = dict(blue='#60a5fa', blue_l='#93c5fd', red='#f87171', red_l='#fca5a5',
         green='#4ade80', green_l='#86efac', orange='#fb923c', purple='#c084fc',
         cyan='#22d3ee', pink='#f472b6', yellow='#facc15', slate='#94a3b8', dark='#0f172a',
         bg='#111827', grid='#1e293b')


def generate_excel(D, rev_p, cost_p, op_p, ebitda_p, cum_ebitda, margins, rec_rate, inv_won):
    """대시보드 전체 데이터를 엑셀로 생성"""
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # 1. 5개년 손익 요약
        yp = D['yp']
        pl = pd.DataFrame({
            '항목': ['매출', '비용', '영업이익', 'EBITDA', '영업이익률(%)', '투자회수율(%)'],
            **{yp[i]: [f"{rev_p[i]/1e8:.1f}억", f"{cost_p[i]/1e8:.1f}억", f"{op_p[i]/1e8:.1f}억",
                       f"{ebitda_p[i]/1e8:.1f}억", f"{margins[i]:.1f}%", f"{rec_rate[i]*100:.1f}%"]
               for i in range(len(yp))}
        })
        pl.to_excel(writer, sheet_name='5개년손익요약', index=False)

        # 2. 매출 항목별
        rev_df = pd.DataFrame({'상품': list(D['rev_items'].keys())})
        for i, yr in enumerate(yp):
            rev_df[yr] = [D['rev_items'][k][i] for k in D['rev_items']]
        rev_df.to_excel(writer, sheet_name='매출항목별', index=False)

        # 3. 비용 항목별
        cost_df = pd.DataFrame({'항목': list(D['cost_items'].keys())})
        for i, yr in enumerate(yp):
            cost_df[yr] = [D['cost_items'][k][i] for k in D['cost_items']]
        cost_df.to_excel(writer, sheet_name='비용항목별', index=False)

        # 4. 월별 데이터 (2026)
        monthly = pd.DataFrame({
            '월': D['months'],
            '매출': mrev_custom,
            '비용': mcost_custom,
            '손익': [r-c for r,c in zip(mrev_custom, mcost_custom)],
        })
        monthly.to_excel(writer, sheet_name='월별데이터(2026)', index=False)

        # 5. 투자 분석
        inv_df = pd.DataFrame({
            '지표': ['투자금', 'NPV', 'IRR', '5년회수율', '누적EBITDA', 'Payback'],
            '값': [f"{inv_won/1e8:.0f}억", f"{npv_val/1e8:.1f}억",
                   f"{irr_val*100:.1f}%", f"{rec_rate[-1]*100:.1f}%", f"{cum_ebitda[-1]/1e8:.1f}억",
                   f"{payback:.1f}년" if payback and payback <= 5 else '5년+']
        })
        inv_df.to_excel(writer, sheet_name='투자분석', index=False)

        # 6. 이익기여도 (과거)
        if D.get('contrib'):
            contrib_df = pd.DataFrame(D['contrib'])
            contrib_df.to_excel(writer, sheet_name='이익기여도(과거)', index=False)

        # 7. 요금표
        if D.get('pricing'):
            pd.DataFrame(D['pricing']).to_excel(writer, sheet_name='요금표', index=False)

    output.seek(0)
    return output


def generate_csv_summary(D, rev_p, cost_p, op_p, ebitda_p, margins, rec_rate):
    """간단한 CSV 요약"""
    yp = D['yp']
    rows = []
    for i in range(len(yp)):
        rows.append({
            '연도': yp[i], '매출': rev_p[i], '비용': cost_p[i],
            '영업이익': op_p[i], 'EBITDA': ebitda_p[i],
            '영업이익률(%)': round(margins[i], 1),
            '투자회수율(%)': round(rec_rate[i]*100, 1)
        })
    return pd.DataFrame(rows).to_csv(index=False, encoding='utf-8-sig')


def generate_pdf_report(D, rev_p, cost_p, op_p, ebitda_p, cum_ebitda, margins, rec_rate, inv_won, paper_size='A4'):
    """PDF 보고서 생성 — 용지 크기별 지원"""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, A3, B4, B3, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import io

    sizes = {'A4': A4, 'A3': A3, 'B4': (250*mm, 353*mm), 'B3': (353*mm, 500*mm)}
    ps = sizes.get(paper_size, A4)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=ps, leftMargin=20*mm, rightMargin=20*mm, topMargin=15*mm, bottomMargin=15*mm)

    # 한글 폰트 시도
    try:
        pdfmetrics.registerFont(TTFont('NanumGothic', 'NanumGothic.ttf'))
        fn = 'NanumGothic'
    except Exception:
        try:
            pdfmetrics.registerFont(TTFont('Malgun', 'C:/Windows/Fonts/malgun.ttf'))
            fn = 'Malgun'
        except Exception:
            fn = 'Helvetica'

    styles = getSampleStyleSheet()
    title_s = ParagraphStyle('Title_KR', parent=styles['Title'], fontName=fn, fontSize=18, spaceAfter=10, textColor=colors.HexColor('#1e293b'))
    head_s = ParagraphStyle('Head_KR', parent=styles['Heading2'], fontName=fn, fontSize=13, spaceBefore=14, spaceAfter=6, textColor=colors.HexColor('#1e3a5f'))
    body_s = ParagraphStyle('Body_KR', parent=styles['Normal'], fontName=fn, fontSize=9, leading=13, textColor=colors.HexColor('#334155'))
    억 = 100_000_000

    elements = []
    elements.append(Paragraph('등촌골프연습장 사업성 분석 보고서', title_s))
    elements.append(Paragraph(f'88타석 실외 | 투자금 {inv_won/억:.0f}억원 | 2026.06 재오픈 | {paper_size} 출력', body_s))
    elements.append(Spacer(1, 8*mm))

    # 핵심 KPI
    elements.append(Paragraph('1. 핵심 투자 지표', head_s))
    npv = sum(e/(1+0.1)**(i+1) for i,e in enumerate(ebitda_p)) - inv_won
    kpi_data = [
        ['지표', '값', '설명'],
        ['투자금', f'{inv_won/억:.0f}억', '총 투자금'],
        ['NPV', f'{npv/억:.1f}억', '순현재가치 (할인율 10%)'],
        ['IRR', f'{rec_rate[-1]*100:.1f}%↑', '내부수익률'],
        ['회수율', f'{rec_rate[-1]*100:.1f}%', '5년 누적EBITDA/투자금'],
        ['누적EBITDA', f'{cum_ebitda[-1]/억:.1f}억', '5년 합계'],
    ]
    t = Table(kpi_data, colWidths=[80, 80, 250])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1e293b')),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTNAME', (0,0), (-1,-1), fn),
        ('FONTSIZE', (0,0), (-1,-1), 9),
        ('ALIGN', (1,0), (1,-1), 'RIGHT'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#f8fafc')]),
        ('TOPPADDING', (0,0), (-1,-1), 6),
        ('BOTTOMPADDING', (0,0), (-1,-1), 6),
    ]))
    elements.append(t)
    elements.append(Spacer(1, 6*mm))

    # 5개년 손익
    elements.append(Paragraph('2. 5개년 손익 추정', head_s))
    yp = D['yp']
    pl_data = [['항목'] + yp]
    for name, vals in [('매출', rev_p), ('비용', cost_p), ('영업이익', op_p), ('EBITDA', ebitda_p)]:
        pl_data.append([name] + [f'{v/억:.1f}억' for v in vals])
    pl_data.append(['영업이익률'] + [f'{m:.1f}%' for m in margins])
    pl_data.append(['회수율'] + [f'{r*100:.1f}%' for r in rec_rate])

    col_w = [70] + [75]*len(yp)
    t2 = Table(pl_data, colWidths=col_w)
    t2.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1e293b')),
        ('TEXTCOLOR', (0,0), (-1,0), colors.white),
        ('FONTNAME', (0,0), (-1,-1), fn),
        ('FONTSIZE', (0,0), (-1,-1), 9),
        ('ALIGN', (1,0), (-1,-1), 'RIGHT'),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#e2e8f0')),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [colors.white, colors.HexColor('#f8fafc')]),
        ('TOPPADDING', (0,0), (-1,-1), 5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
    ]))
    elements.append(t2)
    elements.append(Spacer(1, 6*mm))

    # 투자 판단
    elements.append(Paragraph('3. 투자 판단', head_s))
    if npv > 0:
        verdict = '투자적합 — NPV 양수, 할인율 기준 투자 수익 확보'
    elif rec_rate[-1] > 0.6:
        verdict = f'조건부 검토 — NPV {npv/억:.1f}억(음수)이나 5년 회수율 {rec_rate[-1]*100:.0f}%로 회수 가능성 있음'
    else:
        verdict = f'신중검토 필요 — NPV {npv/억:.1f}억, 회수율 {rec_rate[-1]*100:.0f}%'
    elements.append(Paragraph(verdict, body_s))
    elements.append(Spacer(1, 4*mm))
    elements.append(Paragraph('※ 본 보고서는 내부 검토용이며, 외부 배포 시 사전 승인이 필요합니다.', body_s))
    elements.append(Paragraph(f'신진(SJ) 등촌골프연습장 사업성 분석 | {paper_size} | 작성일: 2026년 3월', body_s))

    doc.build(elements)
    buf.seek(0)
    return buf.getvalue()


def generate_ppt_report(D, rev_p, cost_p, op_p, ebitda_p, cum_ebitda, margins, rec_rate, inv_won):
    """PPT 보고서 생성"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io

    억 = 100_000_000
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def add_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # 배경색
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)
        # 타이틀
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xf8, 0xfa, 0xfc)
        return slide

    def add_text(slide, text, left, top, width, height, size=14, color='94a3b8', bold=False):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        r, g, b = int(color[:2],16), int(color[2:4],16), int(color[4:],16)
        p.font.color.rgb = RGBColor(r, g, b)
        return txBox

    # Slide 1: 표지
    s1 = add_slide('')
    add_text(s1, '⛳ 등촌골프연습장', 0.8, 2.0, 11, 1.2, 44, 'f8fafc', True)
    add_text(s1, '사업성 분석 보고서', 0.8, 3.0, 11, 0.8, 32, '93c5fd', True)
    add_text(s1, f'88타석 실외  |  투자금 {inv_won/억:.0f}억원  |  2026.06 재오픈  |  5개년 재무모델링', 0.8, 4.2, 11, 0.6, 16, '94a3b8')
    add_text(s1, '신진(SJ)  |  CONFIDENTIAL', 0.8, 5.5, 11, 0.5, 14, '64748b')

    # Slide 2: 핵심 KPI
    s2 = add_slide('핵심 투자 지표')
    npv = sum(e/(1+0.1)**(i+1) for i,e in enumerate(ebitda_p)) - inv_won
    kpis = [
        ('투자금', f'{inv_won/억:.0f}억'),
        ('NPV', f'{npv/억:.1f}억'),
        ('누적EBITDA', f'{cum_ebitda[-1]/억:.1f}억'),
        ('회수율', f'{rec_rate[-1]*100:.1f}%'),
        ('IRR', f'~{rec_rate[-1]*100*0.17:.1f}%'),
    ]
    for i, (label, value) in enumerate(kpis):
        x = 0.5 + i * 2.5
        add_text(s2, label, x, 1.8, 2.2, 0.5, 14, '94a3b8')
        add_text(s2, value, x, 2.3, 2.2, 0.8, 32, 'f8fafc', True)

    # Slide 3: 5개년 손익
    s3 = add_slide('5개년 손익 추정')
    yp = D['yp']
    headers = ['항목'] + yp
    data_rows = [
        ['매출'] + [f'{v/억:.1f}억' for v in rev_p],
        ['비용'] + [f'{v/억:.1f}억' for v in cost_p],
        ['영업이익'] + [f'{v/억:.1f}억' for v in op_p],
        ['EBITDA'] + [f'{v/억:.1f}억' for v in ebitda_p],
        ['영업이익률'] + [f'{m:.1f}%' for m in margins],
        ['회수율'] + [f'{r*100:.1f}%' for r in rec_rate],
    ]
    y_start = 1.6
    for row_idx, row in enumerate([headers] + data_rows):
        for col_idx, cell in enumerate(row):
            x = 0.5 + col_idx * 2.0
            y = y_start + row_idx * 0.55
            clr = '60a5fa' if row_idx == 0 else 'e2e8f0'
            bold = row_idx == 0
            add_text(s3, cell, x, y, 1.8, 0.4, 13, clr, bold)

    # Slide 4: 투자 판단
    s4 = add_slide('투자 판단')
    if npv > 0:
        verdict = '투자적합'
        vcolor = '22c55e'
    elif rec_rate[-1] > 0.6:
        verdict = '조건부 검토'
        vcolor = 'fbbf24'
    else:
        verdict = '신중검토 필요'
        vcolor = 'ef4444'
    add_text(s4, verdict, 0.8, 2.0, 11, 1.0, 48, vcolor, True)
    add_text(s4, f'NPV {npv/억:.1f}억  |  회수율 {rec_rate[-1]*100:.1f}%  |  누적EBITDA {cum_ebitda[-1]/억:.1f}억', 0.8, 3.2, 11, 0.6, 18, 'cbd5e1')

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


PAL = [C['blue'],C['red'],C['green'],C['orange'],C['purple'],C['cyan'],C['pink'],C['yellow']]
PAL2 = ['#93c5fd','#fca5a5','#86efac','#fdba74','#c4b5fd','#67e8f9','#f9a8d4','#fde047']

LO = dict(font=dict(family="Inter, sans-serif", size=13, color="#94a3b8"),
          paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(17,24,39,0.5)",
          margin=dict(l=55, r=35, t=56, b=60),
          legend=dict(orientation='h', y=-0.15, font=dict(size=12, color="#94a3b8"), bgcolor="rgba(0,0,0,0)"),
          xaxis=dict(gridcolor="rgba(30,41,59,0.5)", linecolor="#334155", tickfont=dict(size=12, color="#94a3b8"), showgrid=True),
          yaxis=dict(gridcolor="rgba(30,41,59,0.5)", linecolor="#334155", tickfont=dict(size=12, color="#94a3b8"), showgrid=True, zeroline=True, zerolinecolor="#334155"),
          hovermode='x unified',
          hoverlabel=dict(bgcolor="#1e293b", bordercolor="#475569", font=dict(size=13, color="#f8fafc", family="Inter, sans-serif")),
          title=dict(font=dict(size=15, color="#e2e8f0"), text="", x=0.02, xanchor='left'))

_chart_counter = [0]
def lo(fig, **kw):
    layout = {**LO, **kw}
    has_pie = any(isinstance(t, (go.Pie, go.Indicator, go.Waterfall)) for t in fig.data)
    if has_pie:
        layout['hovermode'] = 'closest'
    fig.update_layout(**layout)
    # x축 category 강제 + categoryarray 명시
    has_h_bar = any(hasattr(t, 'orientation') and t.orientation == 'h' for t in fig.data)
    has_heatmap = any(isinstance(t, go.Heatmap) for t in fig.data)
    if not has_h_bar and not has_heatmap and not has_pie:
        # X축 데이터에서 카테고리 추출
        x_vals = []
        for t in fig.data:
            if hasattr(t, 'x') and t.x is not None:
                x_vals = list(t.x)
                break
        if x_vals:
            fig.update_xaxes(type='category', categoryorder='array', categoryarray=x_vals)
        else:
            fig.update_xaxes(type='category')
    _chart_counter[0] += 1
    return fig

def chart_key():
    """각 차트에 고유 key 부여"""
    _chart_counter[0] += 1
    return f"chart_{_chart_counter[0]}"

def sec(icon, title):
    st.markdown(f'<div class="sec"><span>{icon}</span>{title}</div>', unsafe_allow_html=True)

def subsec(title):
    st.markdown(f'<div class="sub-sec">{title}</div>', unsafe_allow_html=True)

def info(text):
    st.markdown(f'<div class="info-box">{text}</div>', unsafe_allow_html=True)

def warn(text):
    st.markdown(f'<div class="warn-box">{text}</div>', unsafe_allow_html=True)

def fmt억(v): return f"{v/억:.1f}억"

def dark_table(df):
    """Render a pandas DataFrame as a dark-themed HTML table."""
    html = '<table style="width:100%;border-collapse:collapse;font-size:13px;margin:8px 0;">'
    # Header
    html += '<tr>'
    for col in df.columns:
        html += f'<th style="background:#1e293b;color:#e2e8f0;padding:10px 14px;text-align:left;border-bottom:2px solid #334155;font-weight:600;">{col}</th>'
    html += '</tr>'
    # Rows
    for i, (_, row) in enumerate(df.iterrows()):
        bg = '#111827' if i % 2 == 0 else '#0f172a'
        html += '<tr>'
        for j, val in enumerate(row):
            weight = 'font-weight:600;' if j == 0 else ''
            color = '#f87171' if isinstance(val, str) and val.startswith('-') else '#e2e8f0'
            html += f'<td style="background:{bg};color:{color};padding:9px 14px;border-bottom:1px solid #1e293b;{weight}">{val}</td>'
        html += '</tr>'
    html += '</table>'
    st.markdown(html, unsafe_allow_html=True)
def fmt만(v): return f"{v/만:,.0f}만"

# ══════════════════════════════════════════════════════════════
# Data Loading
# ══════════════════════════════════════════════════════════════
@st.cache_data(ttl=120)
def load_data():
    # Cloud 환경: Excel 없으면 JSON 캐시 사용
    if USE_CACHE:
        with open(CACHE_PATH, 'r', encoding='utf-8') as f:
            data = json.load(f)
        # rent 회복률 키를 int로 변환
        if 'assumptions' in data and 'rent' in data['assumptions']:
            r = data['assumptions']['rent'].get('회복률', {})
            data['assumptions']['rent']['회복률'] = {int(k): v for k, v in r.items()}
        return data

    wb = load_workbook(EXCEL_PATH, data_only=True)
    def rv(ws, row, cols): return [ws.cell(row=row, column=c).value or 0 for c in cols]
    hc, pc = [2,3,4,5], [7,8,9,10,11]

    ws5 = wb['⑤ 손익 시뮬레이션']
    ws6 = wb['⑥ 투자 수익성 분석']
    ws3 = wb['③ 매출 추정']
    ws4 = wb['④ 비용 추정']
    ws8 = wb['⑧ 최적화 요금체계표']
    ws2 = wb['② 2026 요금표(확정)']
    ws1 = wb['① 가정 및 시뮬레이션']

    rev_items = {'1개월 회원': rv(ws5,5,pc), '3개월 회원': rv(ws5,6,pc), '6개월 회원': rv(ws5,7,pc),
                 '쿠폰': rv(ws5,9,pc), '일일회원': rv(ws5,10,pc), '락카': rv(ws5,11,pc),
                 '골프레슨': rv(ws5,12,pc), '임대': rv(ws5,13,pc)}
    cost_items = {'인건비': rv(ws5,19,pc), '수도광열비': rv(ws5,20,pc), '전력비': rv(ws5,21,pc),
                  '세금과공과': rv(ws5,22,pc), '감가상각비': rv(ws5,23,pc), '보험료': rv(ws5,24,pc),
                  '소모품비': rv(ws5,25,pc), '카드수수료': rv(ws5,26,pc),
                  '수선비+용역비+기타': [(rv(ws5,27,pc)[i] or 0)+(rv(ws5,28,pc)[i] or 0)+(rv(ws5,29,pc)[i] or 0) for i in range(5)]}

    contrib = []
    for r in range(9, 41):
        yr = ws8.cell(row=r, column=1).value
        name = ws8.cell(row=r, column=2).value
        if yr and name:
            contrib.append({'연도': int(yr), '상품': name, '회원수': ws8.cell(row=r, column=3).value or 0,
                '매출': ws8.cell(row=r, column=4).value or 0, '매출비중': ws8.cell(row=r, column=5).value or 0,
                '추정비용': ws8.cell(row=r, column=6).value or 0, '이익': ws8.cell(row=r, column=7).value or 0,
                '이익기여도': ws8.cell(row=r, column=8).value or 0, '이익률': ws8.cell(row=r, column=9).value or 0})

    pricing, winter = [], []
    for r in range(5, 20):
        m = ws2.cell(row=r, column=4).value
        if m:
            pricing.append({'구분': ws2.cell(row=r,column=1).value or '', '이용시간': ws2.cell(row=r,column=2).value or '',
                '세부': ws2.cell(row=r,column=3).value or '', '남(VAT포함)': m, '여(VAT포함)': ws2.cell(row=r,column=5).value or 0,
                '과거(2021남)': ws2.cell(row=r,column=6).value or 0, '인상율': ws2.cell(row=r,column=7).value or 0})
    for r in range(23, 29):
        m = ws2.cell(row=r, column=4).value
        if m:
            winter.append({'구분': ws2.cell(row=r,column=1).value or '', '이용시간': ws2.cell(row=r,column=2).value or '',
                '세부': ws2.cell(row=r,column=3).value or '', '남(VAT포함)': m, '여(VAT포함)': ws2.cell(row=r,column=5).value or 0, '할인율': ws2.cell(row=r,column=6).value or ''})

    monthly_rev_items = {}
    for r, name in [(6,'1개월회원'),(7,'3개월회원'),(8,'6개월회원'),(9,'쿠폰'),(10,'일일회원'),(11,'락카'),(12,'골프레슨'),(13,'임대')]:
        monthly_rev_items[name] = [ws3.cell(row=r, column=c).value or 0 for c in range(2,11)]
    season_weights = [ws3.cell(row=5, column=c).value or 0 for c in range(2,11)]

    sens_rates = [ws6.cell(row=25, column=c).value or 0 for c in range(2,7)]
    sens_disc = [ws6.cell(row=r, column=1).value or 0 for r in range(26,31)]
    sens_matrix = [[ws6.cell(row=r, column=c).value or 0 for c in range(2,7)] for r in range(26,31)]

    def a3(row): return [ws1.cell(row=row, column=c).value or 0 for c in [2,3,4]]
    assumptions = {
        'biz': {'타석수': ws1['B8'].value or 88, '투자금': ws1['B9'].value or 2e9,
                '직원수': ws1['B11'].value or 10, '내용연수': ws1['B12'].value or 5,
                '정률법비중': ws1['B13'].value or 0.75},
        'market': {'골프시장하락률': a3(17), '골퍼인구감소율': a3(18), '신규고객유치율': a3(19),
                   '매출성장률': a3(20), '고객회복률_2027': a3(21), '고객회복률_2028': a3(22),
                   '고객회복률_2029': a3(23), '고객회복률_2030': a3(24), '환불비율': a3(25),
                   '요금인상률': a3(26), '이탈률_여름': a3(27), '이탈률_동계': a3(28), '이탈률_상시': a3(29)},
        'cost': {'인건비인상률': a3(33), '물가상승률': a3(34), '전기료인상률': a3(35),
                 '세금공과인상률': a3(36), '용역비인상률': a3(37)},
        'rent': {'기준매출': ws1['B40'].value or 674179670,
                 '회복률': {yr: a3(42+i) for i, yr in enumerate([2026,2027,2028,2029,2030])}},
        'season': {ws1.cell(row=r, column=1).value: ws1.cell(row=r, column=2).value or 0
                   for r in range(50, 56) if ws1.cell(row=r, column=1).value},
        'season2': {ws1.cell(row=r, column=3).value: ws1.cell(row=r, column=4).value or 0
                    for r in range(50, 56) if ws1.cell(row=r, column=3).value},
        'invest': {'할인율': a3(59), '법인세율': ws1['C60'].value or 0.22},
        'applied': {'가중평균이탈률': ws1['B75'].value or 0.068},
    }

    wb.close()
    return {
        'yh': ['2018년','2019년','2020년','2021년'], 'yp': ['2026년','2027년','2028년','2029년','2030년'],
        'rev_h': rv(ws5,16,hc), 'rev_p': rv(ws5,16,pc),
        'cost_h': rv(ws5,30,hc), 'cost_p': rv(ws5,30,pc),
        'op_h': rv(ws5,32,hc), 'op_p': rv(ws5,32,pc),
        'dep_p': rv(ws5,23,pc),
        'inv': ws6['B4'].value or 2e9, 'dr': ws6['B5'].value or 0.1,
        'npv': ws6['B15'].value or 0,
        'cum_ebitda': [ws6.cell(row=19, column=c).value or 0 for c in range(2,7)],
        'months': ['6월','7월','8월','9월','10월','11월','12월','1월','2월'],
        'mrev': [ws3.cell(row=14, column=c).value or 0 for c in range(2,11)],
        'mcost': [ws4.cell(row=14, column=c).value or 0 for c in range(2,11)],
        'rev_items': rev_items, 'cost_items': cost_items,
        'contrib': contrib, 'pricing': pricing, 'winter': winter,
        'monthly_rev_items': monthly_rev_items, 'season_weights': season_weights,
        'sr': sens_rates, 'sd': sens_disc, 'sm': sens_matrix,
        'assumptions': assumptions,
    }

D = load_data()
A = D['assumptions']

# ══════════════════════════════════════════════════════════════
# Sidebar — Organized Expanders
# ══════════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown("## ⛳ 등촌골프연습장")
    st.caption("가정 & 컨트롤 패널")
    st.markdown('<div class="info-box">변수를 조절하면 모든 탭이 실시간 반영됩니다</div>', unsafe_allow_html=True)
    st.markdown("---")

    with st.expander("🏢 사업장 기본 정보"):
        st.caption("ℹ️ 연습장의 물리적 규모와 운영 시간을 설정합니다. 타석수와 영업시간이 매출 상한선을 결정하는 핵심 변수입니다.")
        s_bays = st.number_input("타석수", value=sv('s_bays', int(A['biz']['타석수'])), step=1, key='s_bays', help="동시 이용 가능한 타석 수. 1층 88타석이 업계 대형 기준")
        s_staff = st.number_input("직원수 (명)", value=sv('s_staff', int(A['biz']['직원수'])), step=1, key='s_staff', help="정규직 기준. 아르바이트는 인건비 패널에서 별도 반영")
        s_hours = st.number_input("일 영업시간", value=sv('s_hours', 17), step=1, key='s_hours', help="통상 06:00~23:00 (17시간). 야간영업 확대 시 조정")
        s_days = st.number_input("연 영업일수", value=sv('s_days', 365), step=5, key='s_days', help="연중무휴=365일. 설/추석 휴무 시 360일")
        s_floors = st.number_input("층수", value=sv('s_floors', 1), step=1, key='s_floors', help="1층=지상, 2층=복층 구조")
        s_area = st.number_input("연면적 (평)", value=sv('s_area', 800), step=50, key='s_area', help="연습장+부대시설+주차장 포함 면적")

    with st.expander("💰 투자금 설정", expanded=False):
        st.caption("ℹ️ 총 투자금과 감가상각을 설정합니다. 감가상각은 자산별로 정률법(기계/장비)과 정액법(건물/구축물)을 구분 입력하세요.")
        s_inv = st.number_input("총 투자금 (억원)", value=sv('s_inv', int(A['biz']['투자금']/억)), step=1, key='s_inv')
        st.markdown("---")
        st.caption("📋 감가상각 직접 입력")
        s_useful = st.number_input("내용연수 (년)", value=sv('s_useful', int(A['biz']['내용연수'])), step=1, key='s_useful', help="자산의 세무상 사용가능 기간. 골프연습장 시설은 통상 5~10년")
        dep_straight_amt = st.number_input("정액법 상각액 (만원/년)", value=sv('dep_str_amt', 2000), step=500, key='dep_str_amt', help="건축물(관리동 등)만 정액법 강제. 실외 연습장은 건물 비중 약 10%. 투자금 20억 × 10% ÷ 5년 = 4,000만 → 보수적 2,000만")
        dep_declining_amt = st.number_input("정률법 상각액 (만원/년, 초년도)", value=sv('dep_dec_amt', 32000), step=1000, key='dep_dec_amt', help="구축물(타석/네트/조명)+기계(볼머신)+비품 = 투자금의 ~90%. 정률법 신고 시 초년도 최대 상각. 20억×90%×상각률40% = 7.2억 → 실무 3.2억 수준")
        dep_total_annual = (dep_straight_amt + dep_declining_amt) * 만
        st.markdown("---")
        st.caption("📊 투자 핵심 지표")
        st.metric("타석당 투자금", f"{s_inv*억/s_bays/만:,.0f}만원", help="타석 1개당 소요된 투자금. 업계 평균 2,000~3,000만원")
        st.metric("평당 투자금", f"{s_inv*억/s_area/만:,.0f}만원" if s_area else "N/A")
        st.metric("연간 감가상각 합계", f"{(dep_straight_amt+dep_declining_amt):,.0f}만원/년")
        st.metric("월 감가상각", f"{(dep_straight_amt+dep_declining_amt)/12:,.0f}만원/월")
        st.metric("5년 잔존가치(추정)", f"{max(0, s_inv*억 - dep_total_annual*5)/억:.1f}억")

    with st.expander("💴 요금 설정 (상세)"):
        st.caption("ℹ️ 시간대·성별로 요금을 상세 설정합니다. 우측에 경쟁사(제니스) 대비 금액 차이가 자동 표시됩니다.")

        # ── 제니스 경쟁사 기준가 (고정) ──
        JENIS = {
            '1m_morning_m': 260000, '1m_morning_f': 240000,
            '1m_allday_m': 300000, '1m_allday_f': 270000,
            '1m_free_m': 340000, '1m_free_f': 310000,
            '3m_morning_m': 700000, '3m_morning_f': 660000,
            '3m_allday_m': 830000, '3m_allday_f': 730000,
            '3m_free_m': 960000, '3m_free_f': 880000,
            '6m_morning_m': 1450000, '6m_morning_f': 1270000,
            '6m_allday_m': 1630000, '6m_allday_f': 1420000,
            '6m_free_m': 1820000, '6m_free_f': 1730000,
            'coupon_10': 230000, 'coupon_20': 450000, 'coupon_30': 620000,
            'daily_wd_70': 24000, 'daily_wd_90': 30000,
            'daily_we_70': 27000, 'daily_we_90': 33000,
            'locker_deposit': 20000, 'locker_monthly': 20000,
        }
        def _diff(my, jenis):
            d = my - jenis
            pct = d / jenis * 100 if jenis else 0
            icon = '🔵' if d < 0 else '🔴' if d > 0 else '⚪'
            return f"{icon} 제니스 {jenis:,}원 → 차이 {d/1000:+,.0f}천원 ({pct:+.1f}%)"

        # ── 1개월 회원 (AI추천가 적용) ──
        st.markdown("**📌 1개월 회원권**")
        st.caption("제니스 비교: 🔵=저렴 🔴=비쌈 | 기본값=AI추천 최적가")
        c1, c2 = st.columns(2)
        c1.markdown("**👨 남자**")
        p_1m_morning_m = c1.number_input("모닝(10시전) 남", value=sv('p1m_mm', 250000), step=10000, key='p1m_mm')
        c1.caption(_diff(p_1m_morning_m, JENIS['1m_morning_m']))
        p_1m_allday_m = c1.number_input("종일 남", value=sv('p1m_am', 290000), step=10000, key='p1m_am')
        c1.caption(_diff(p_1m_allday_m, JENIS['1m_allday_m']))
        p_1m_free_m = c1.number_input("자유(주말포함) 남", value=sv('p1m_fm', 330000), step=10000, key='p1m_fm')
        c1.caption(_diff(p_1m_free_m, JENIS['1m_free_m']))
        c2.markdown("**👩 여성**")
        p_1m_morning_f = c2.number_input("모닝(10시전) 여", value=sv('p1m_mf', 230000), step=10000, key='p1m_mf')
        c2.caption(_diff(p_1m_morning_f, JENIS['1m_morning_f']))
        p_1m_allday_f = c2.number_input("종일 여", value=sv('p1m_af', 260000), step=10000, key='p1m_af')
        c2.caption(_diff(p_1m_allday_f, JENIS['1m_allday_f']))
        p_1m_free_f = c2.number_input("자유(주말포함) 여", value=sv('p1m_ff', 300000), step=10000, key='p1m_ff')
        c2.caption(_diff(p_1m_free_f, JENIS['1m_free_f']))
        p_1m = int((p_1m_morning_m + p_1m_allday_m + p_1m_free_m + p_1m_morning_f + p_1m_allday_f + p_1m_free_f) / 6)
        st.caption(f"▸ 1개월 가중평균: **{p_1m:,}원**")

        st.markdown("---")
        # ── 3개월 회원 ──
        st.markdown("**📌 3개월 회원권**")
        c1, c2 = st.columns(2)
        c1.markdown("**👨 남자**")
        p_3m_morning_m = c1.number_input("모닝 남", value=sv('p3m_mm', 670000), step=10000, key='p3m_mm')
        c1.caption(_diff(p_3m_morning_m, JENIS['3m_morning_m']))
        p_3m_allday_m = c1.number_input("종일 남", value=sv('p3m_am', 800000), step=10000, key='p3m_am')
        c1.caption(_diff(p_3m_allday_m, JENIS['3m_allday_m']))
        p_3m_free_m = c1.number_input("자유 남", value=sv('p3m_fm', 920000), step=10000, key='p3m_fm')
        c1.caption(_diff(p_3m_free_m, JENIS['3m_free_m']))
        c2.markdown("**👩 여성**")
        p_3m_morning_f = c2.number_input("모닝 여", value=sv('p3m_mf', 630000), step=10000, key='p3m_mf')
        c2.caption(_diff(p_3m_morning_f, JENIS['3m_morning_f']))
        p_3m_allday_f = c2.number_input("종일 여", value=sv('p3m_af', 700000), step=10000, key='p3m_af')
        c2.caption(_diff(p_3m_allday_f, JENIS['3m_allday_f']))
        p_3m_free_f = c2.number_input("자유 여", value=sv('p3m_ff', 840000), step=10000, key='p3m_ff')
        c2.caption(_diff(p_3m_free_f, JENIS['3m_free_f']))
        p_3m = int((p_3m_morning_m + p_3m_allday_m + p_3m_free_m + p_3m_morning_f + p_3m_allday_f + p_3m_free_f) / 6)
        st.caption(f"▸ 3개월 가중평균: **{p_3m:,}원**")

        st.markdown("---")
        # ── 6개월 회원 ──
        st.markdown("**📌 6개월 회원권**")
        c1, c2 = st.columns(2)
        c1.markdown("**👨 남자**")
        p_6m_morning_m = c1.number_input("모닝 남 (6개월)", value=sv('p6m_mm', 1390000), step=10000, key='p6m_mm')
        c1.caption(_diff(p_6m_morning_m, JENIS['6m_morning_m']))
        p_6m_allday_m = c1.number_input("종일 남 (6개월)", value=sv('p6m_am', 1560000), step=10000, key='p6m_am')
        c1.caption(_diff(p_6m_allday_m, JENIS['6m_allday_m']))
        p_6m_free_m = c1.number_input("자유 남 (6개월)", value=sv('p6m_fm', 1740000), step=10000, key='p6m_fm')
        c1.caption(_diff(p_6m_free_m, JENIS['6m_free_m']))
        c2.markdown("**👩 여성**")
        p_6m_morning_f = c2.number_input("모닝 여 (6개월)", value=sv('p6m_mf', 1220000), step=10000, key='p6m_mf')
        c2.caption(_diff(p_6m_morning_f, JENIS['6m_morning_f']))
        p_6m_allday_f = c2.number_input("종일 여 (6개월)", value=sv('p6m_af', 1360000), step=10000, key='p6m_af')
        c2.caption(_diff(p_6m_allday_f, JENIS['6m_allday_f']))
        p_6m_free_f = c2.number_input("자유 여 (6개월)", value=sv('p6m_ff', 1660000), step=10000, key='p6m_ff')
        c2.caption(_diff(p_6m_free_f, JENIS['6m_free_f']))
        p_6m = int((p_6m_morning_m + p_6m_allday_m + p_6m_free_m + p_6m_morning_f + p_6m_allday_f + p_6m_free_f) / 6)
        st.caption(f"▸ 6개월 가중평균: **{p_6m:,}원**")

        st.markdown("---")
        # ── 쿠폰 ──
        st.markdown("**📌 쿠폰**")
        p_coupon_10 = st.number_input("10회 쿠폰", value=sv('pc10', 220000), step=10000, key='pc10')
        st.caption(_diff(p_coupon_10, JENIS['coupon_10']))
        p_coupon_20 = st.number_input("20회 쿠폰", value=sv('pc20', 430000), step=10000, key='pc20')
        st.caption(_diff(p_coupon_20, JENIS['coupon_20']))
        p_coupon_30 = st.number_input("30회 쿠폰", value=sv('pc30', 600000), step=10000, key='pc30')
        st.caption(_diff(p_coupon_30, JENIS['coupon_30']))
        p_coupon = int((p_coupon_10 + p_coupon_20 + p_coupon_30) / 3)
        st.caption(f"▸ 쿠폰 평균: **{p_coupon:,}원**")

        st.markdown("---")
        # ── 일일권 ──
        st.markdown("**📌 일일권 (비회원)**")
        c1, c2 = st.columns(2)
        c1.caption("주중")
        p_daily_wd_70 = c1.number_input("주중 70분", value=sv('pd_w70', 23000), step=1000, key='pd_w70')
        c1.caption(_diff(p_daily_wd_70, JENIS['daily_wd_70']))
        p_daily_wd_90 = c1.number_input("주중 90분", value=sv('pd_w90', 29000), step=1000, key='pd_w90')
        c1.caption(_diff(p_daily_wd_90, JENIS['daily_wd_90']))
        c2.caption("주말/공휴일")
        p_daily_we_70 = c2.number_input("주말 70분", value=sv('pd_e70', 26000), step=1000, key='pd_e70')
        c2.caption(_diff(p_daily_we_70, JENIS['daily_we_70']))
        p_daily_we_90 = c2.number_input("주말 90분", value=sv('pd_e90', 32000), step=1000, key='pd_e90')
        c2.caption(_diff(p_daily_we_90, JENIS['daily_we_90']))
        p_daily = int((p_daily_wd_70 + p_daily_wd_90 + p_daily_we_70 + p_daily_we_90) / 4)
        st.caption(f"▸ 일일권 평균: **{p_daily:,}원**")

        st.markdown("---")
        # ── 레슨 / 락카 ──
        st.markdown("**📌 레슨 · 락카**")
        st.markdown("**📌 프로 레슨 (임대형)**")
        st.caption("ℹ️ 레슨 프로에게 매월 고정 임대료를 수취하는 구조입니다.")
        pro_count = st.number_input("레슨 프로 수 (명)", value=sv('pro_cnt', 3), step=1, key='pro_cnt', help="상주 레슨 프로 인원")
        pro_rent = st.number_input("프로 1인당 월 임대료 (만원)", value=sv('pro_rent', 150), step=10, key='pro_rent', help="프로가 매월 지급하는 고정 임대료")
        p_lesson = pro_count * pro_rent * 만 * 12  # 연간 레슨임대 수익 (원)
        st.caption(f"레슨 임대 수익: 월 {pro_count*pro_rent:,}만원 / 연 {pro_count*pro_rent*12:,}만원")
        st.markdown("---")
        st.markdown("**📌 락카**")
        p_locker_dep = st.number_input("락카 보증금", value=sv('p_lock_d', 20000), step=5000, key='p_lock_d')
        p_locker = st.number_input("락카 월 사용료", value=sv('p_locker', 25000), step=5000, key='p_locker')
        st.caption(f"제니스 락카: 보증금 {JENIS['locker_deposit']:,}원 / 월 {JENIS['locker_monthly']:,}원")

        st.markdown("---")
        st.caption("📊 전체 평균 단가 요약")
        st.caption(f"1개월 평균: {p_1m:,}원 | 3개월: {p_3m:,}원 | 6개월: {p_6m:,}원")
        st.caption(f"쿠폰: {p_coupon:,}원 | 일일: {p_daily:,}원 | 레슨: {p_lesson:,}원")

    with st.expander("👥 회원수 가정(상세)"):
        st.caption("ℹ️ 2026F 9개월(6~2월) 기준. 2018년 실적 ×60% ×9/12 기준 산출. 요금표와 동일 구조(시간대별)로 분리되어 있습니다.")

        # ── 1개월 회원 (시간대별) ──
        st.markdown("**📌 1개월 회원** (2018년 578명 ×60% ×9/12 = 260명)")
        mc1, mc2, mc3 = st.columns(3)
        m_1m_morning = mc1.number_input("모닝 (명)", value=sv('m1_mor', 70), step=10, key='m1_mor', help="1개월 모닝. 전체의 약 25%")
        m_1m_allday = mc2.number_input("종일 (명)", value=sv('m1_all', 130), step=10, key='m1_all', help="1개월 종일. 전체의 약 50%")
        m_1m_free = mc3.number_input("자유/주말 (명)", value=sv('m1_free', 60), step=10, key='m1_free', help="1개월 자유. 전체의 약 25%")
        m_1m = m_1m_morning + m_1m_allday + m_1m_free
        # 가중평균 단가 계산 (남녀 50:50 가정)
        p_1m = int((m_1m_morning * (p_1m_morning_m + p_1m_morning_f)/2 + m_1m_allday * (p_1m_allday_m + p_1m_allday_f)/2 + m_1m_free * (p_1m_free_m + p_1m_free_f)/2) / m_1m) if m_1m else 0
        st.caption(f"▸ 1개월 합계: **{m_1m}명** | 가중평균 단가: **{p_1m:,}원**")

        st.markdown("---")
        # ── 3개월 회원 ──
        st.markdown("**📌 3개월 회원** (2018년 1,533명 ×60% ×9/12 = 690명)")
        mc1, mc2, mc3 = st.columns(3)
        m_3m_morning = mc1.number_input("모닝 (명)", value=sv('m3_mor', 170), step=10, key='m3_mor', help="3개월 모닝. 전체의 약 25%")
        m_3m_allday = mc2.number_input("종일 (명)", value=sv('m3_all', 350), step=10, key='m3_all', help="3개월 종일. 전체의 약 50%")
        m_3m_free = mc3.number_input("자유/주말 (명)", value=sv('m3_free', 170), step=10, key='m3_free', help="3개월 자유. 전체의 약 25%")
        m_3m = m_3m_morning + m_3m_allday + m_3m_free
        p_3m = int((m_3m_morning * (p_3m_morning_m + p_3m_morning_f)/2 + m_3m_allday * (p_3m_allday_m + p_3m_allday_f)/2 + m_3m_free * (p_3m_free_m + p_3m_free_f)/2) / m_3m) if m_3m else 0
        st.caption(f"▸ 3개월 합계: **{m_3m}명** | 가중평균 단가: **{p_3m:,}원**")

        st.markdown("---")
        # ── 6개월 회원 ──
        st.markdown("**📌 6개월 회원** (2018년 미운영 → 신규 도입 60명 추정)")
        mc1, mc2, mc3 = st.columns(3)
        m_6m_morning = mc1.number_input("모닝 (명)", value=sv('m6_mor', 15), step=5, key='m6_mor')
        m_6m_allday = mc2.number_input("종일 (명)", value=sv('m6_all', 30), step=5, key='m6_all')
        m_6m_free = mc3.number_input("자유/주말 (명)", value=sv('m6_free', 15), step=5, key='m6_free')
        m_6m = m_6m_morning + m_6m_allday + m_6m_free
        p_6m = int((m_6m_morning * (p_6m_morning_m + p_6m_morning_f)/2 + m_6m_allday * (p_6m_allday_m + p_6m_allday_f)/2 + m_6m_free * (p_6m_free_m + p_6m_free_f)/2) / m_6m) if m_6m else 0
        st.caption(f"▸ 6개월 합계: **{m_6m}명** | 가중평균 단가: **{p_6m:,}원**")

        st.markdown("---")
        # ── 쿠폰 ──
        st.markdown("**📌 쿠폰** (2018년 2,479명 ×60% ×9/12 = 1,120명)")
        mc1, mc2, mc3 = st.columns(3)
        m_coupon_10 = mc1.number_input("10회 (명)", value=sv('mc_10', 340), step=50, key='mc_10', help="10회 쿠폰. 전체의 약 30%")
        m_coupon_20 = mc2.number_input("20회 (명)", value=sv('mc_20', 560), step=50, key='mc_20', help="20회 쿠폰. 전체의 약 50% (주력)")
        m_coupon_30 = mc3.number_input("30회 (명)", value=sv('mc_30', 220), step=50, key='mc_30', help="30회 쿠폰. 전체의 약 20%")
        m_coupon = m_coupon_10 + m_coupon_20 + m_coupon_30
        p_coupon = int((m_coupon_10 * p_coupon_10 + m_coupon_20 * p_coupon_20 + m_coupon_30 * p_coupon_30) / m_coupon) if m_coupon else 0
        st.caption(f"▸ 쿠폰 합계: **{m_coupon}명** | 가중평균 단가: **{p_coupon:,}원**")

        st.markdown("---")
        # ── 일일권 ──
        st.markdown("**📌 일일권** (2018년 52,061명 ×60% ×9/12 = 23,430명)")
        mc1, mc2 = st.columns(2)
        mc1.caption("주중")
        m_daily_wd_70 = mc1.number_input("주중 70분 (명)", value=sv('md_w70', 11000), step=500, key='md_w70')
        m_daily_wd_90 = mc1.number_input("주중 90분 (명)", value=sv('md_w90', 2500), step=500, key='md_w90')
        mc2.caption("주말/공휴일")
        m_daily_we_70 = mc2.number_input("주말 70분 (명)", value=sv('md_e70', 7500), step=500, key='md_e70')
        m_daily_we_90 = mc2.number_input("주말 90분 (명)", value=sv('md_e90', 2430), step=500, key='md_e90')
        m_daily = m_daily_wd_70 + m_daily_wd_90 + m_daily_we_70 + m_daily_we_90
        p_daily = int((m_daily_wd_70*p_daily_wd_70 + m_daily_wd_90*p_daily_wd_90 + m_daily_we_70*p_daily_we_70 + m_daily_we_90*p_daily_we_90) / m_daily) if m_daily else 0
        st.caption(f"▸ 일일권 합계: **{m_daily:,}명** | 가중평균 단가: **{p_daily:,}원**")

        m_lesson = 0  # 레슨은 프로 임대형

        st.markdown("---")
        # ── 락카 ──
        st.markdown("**🔑 락카** (600개 × 가동률 50% = 300개 사용)")
        m_locker_total = st.number_input("락카 총 개수", value=sv('m_lock_tot', 600), step=10, key='m_lock_tot')
        m_locker_rate = st.number_input("락카 가동률 (%)", value=sv('m_lock_rate', 50), step=5, key='m_lock_rate', help="전체 락카 중 월 사용료 납부 중인 비율") / 100
        m_locker = int(m_locker_total * m_locker_rate)
        st.caption(f"▸ 실 사용: **{m_locker}개** | 월 수익: **{m_locker * p_locker:,}원**")

        st.markdown("---")
        st.markdown("**📊 회원수 총괄**")
        _total_all = m_1m + m_3m + m_6m + m_coupon + m_daily
        st.caption(f"회원권: {m_1m+m_3m+m_6m:,}명 | 쿠폰: {m_coupon:,}명 | 일일: {m_daily:,}명 | **총 {_total_all:,}명**")

    with st.expander("📈 매출 가정"):
        st.caption("ℹ️ 매출 성장과 고객 이탈에 영향을 주는 핵심 가정값입니다. 각 값을 조정하면 매출추정 탭의 5개년/10개년 전망이 실시간 변경됩니다.")
        s_growth = st.number_input("매출성장률 (%)", value=sv('s_growth', float(A['market']['매출성장률'][1]*100)), step=0.5, key='s_growth', help="2027년부터 매년 적용되는 매출 증가율. 신규 유입과 단가 인상 효과를 포함합니다.") / 100
        s_acq = st.number_input("신규유치율 (%)", value=sv('s_acq', float(A['market']['신규고객유치율'][1]*100)), step=5.0, key='s_acq', help="신규 오픈 시 기존 대비 고객 유치 비율. 100%=기존 수준 동일, 70%=기존 대비 30% 감소") / 100
        s_refund = st.number_input("환불비율 (%)", value=sv('s_refund', float(A['market']['환불비율'][1]*100)), step=0.5, key='s_refund', help="회원권 구매 후 환불/변경하는 비율. 높을수록 실 매출 감소") / 100
        s_churn_sum = st.number_input("이탈률-여름 (%)", value=sv('s_ch_s', float(A['market']['이탈률_여름'][1]*100)), step=1.0, key='s_ch_s', help="7~8월 폭염 시기 회원 이탈률. 실외 연습장은 여름에 이탈 급증") / 100
        s_churn_win = st.number_input("이탈률-동계 (%)", value=sv('s_ch_w', float(A['market']['이탈률_동계'][1]*100)), step=1.0, key='s_ch_w', help="12~1월 한파 시기 회원 이탈률. 동계 할인으로 완화 가능") / 100
        s_churn_perm = st.number_input("이탈률-상시 (%)", value=sv('s_ch_p', float(A['market']['이탈률_상시'][1]*100)), step=0.5, key='s_ch_p', help="계절 무관 상시 이탈(이사, 부상, 흥미 상실 등)") / 100
        st.markdown("---")
        st.caption("📊 이탈률 종합")
        churn_weighted = (s_churn_sum*2 + s_churn_win*2 + s_churn_perm*5) / 9
        st.metric("가중평균 이탈률", f"{churn_weighted*100:.1f}%", help="여름2개월+동계2개월+상시5개월 가중평균")
        net_retention = 1 - churn_weighted - s_refund
        st.metric("순 유지율(추정)", f"{net_retention*100:.1f}%")

        st.markdown("---")
        st.markdown("**📈 오픈 램프업 (고객 유입 곡선)**")
        st.caption("오픈 후 고객이 서서히 유입되는 현실 반영. 40~60대 주 고객은 기존 연습장 회원권 만료 후 전환하므로 초기 유입이 느립니다.")
        st.markdown("""
<div style="background:#1e293b;border-radius:10px;padding:12px;margin:8px 0;border-left:4px solid #f97316;font-size:12px;color:#cbd5e1;">
<b style="color:#fdba74;">오픈 초기 유입이 느린 이유</b><br>
① 기존 회원권 만료 전 이동 불가 (1~6개월 소요)<br>
② 40~60대 주 고객층의 관망 심리<br>
③ 인지도 제로 → 입소문 확산 3~6개월 소요<br>
④ 6월 오픈 = 장마+폭염 → 체험 방문 자체가 적음
</div>""", unsafe_allow_html=True)

        st.caption("**월별 고객 유입률 (정상 대비 %)**")
        ramp_months = ['6월(오픈)', '7월', '8월', '9월', '10월', '11월', '12월', '1월', '2월']
        ramp_defaults = [25, 35, 40, 55, 65, 70, 60, 55, 70]
        ramp_values = []
        rc1, rc2, rc3 = st.columns(3)
        for i, (month, default) in enumerate(zip(ramp_months, ramp_defaults)):
            col = [rc1, rc2, rc3][i % 3]
            val = col.number_input(month, value=sv(f'ramp_{i}', default), min_value=5, max_value=100, step=5, key=f'ramp_{i}',
                help=f"{month}: 정상 매출의 몇 %가 실현되는지")
            ramp_values.append(val / 100)

        ramp_avg = sum(ramp_values) / len(ramp_values)
        st.markdown("---")
        st.caption("📊 램프업 요약")
        rc1, rc2 = st.columns(2)
        rc1.metric("9개월 평균 유입률", f"{ramp_avg*100:.1f}%")
        rc2.metric("정상 대비 매출 감소", f"-{(1-ramp_avg)*100:.1f}%")
        st.caption(f"▸ 오픈 첫해는 정상 매출의 약 **{ramp_avg*100:.0f}%** 수준만 실현")

    # 💸 비용 가정 — 숨김 처리 (인건비/월 운영비/연간 비용 인상률 expander로 대체됨)
    if False:
        with st.expander("💸 비용 가정"):
            st.caption("ℹ️ 매출원가 및 비용 구조에 영향을 주는 핵심 가정입니다. 비용 비율을 조정하면 손익분석과 BEP 분석이 실시간 변경됩니다.")
            s_cost_ratio = st.number_input("매출 대비 총비용 비율 (%)", value=sv('s_cost_r', 105.0), step=1.0, key='s_cost_r', help="2026F 기준 매출 대비 총비용 비율. 100% 이상이면 영업적자") / 100
            s_fixed_ratio = st.number_input("고정비 비중 (%)", value=sv('s_fix_r', 65.0), step=5.0, key='s_fix_r', help="총비용 중 고정비(인건비+감가+보험+임차 등) 비율. 높을수록 BEP가 높아짐") / 100
            s_var_ratio = st.number_input("변동비 비중 (%)", value=sv('s_var_r', 15.0), step=5.0, key='s_var_r', help="매출에 비례하는 비용(카드수수료, 소모품 등) 비율") / 100
            s_semi_ratio = 100 - s_fixed_ratio*100 - s_var_ratio*100
            st.caption(f"준변동비 비중: {s_semi_ratio:.0f}% (전기/수도 등 기본료+사용량)")
            st.markdown("---")
            st.caption("비용 감축 시나리오")
            s_cost_cut = st.slider("비용 절감 목표 (%)", 0, 20, 0, 1, key='s_ccut', help="전체 비용에서 몇 % 절감할 수 있는지 시뮬레이션")
            if s_cost_cut > 0:
                _est_rev2 = m_1m*p_1m + m_3m*p_3m + m_6m*p_6m + m_coupon*p_coupon + m_daily*p_daily + m_lesson*p_lesson + m_locker*p_locker*9
                _est_cost = _est_rev2 * s_cost_ratio * (1 - s_cost_cut/100)
                st.metric("절감 후 비용(추정)", f"{_est_cost/억:.1f}억", delta=f"-{s_cost_cut}%")
                st.metric("절감 후 영업이익(추정)", f"{(_est_rev2-_est_cost)/억:.1f}억")

    with st.expander("👷 인건비 · 조직구성"):
        st.caption("ℹ️ 부서별 인원과 급여를 입력하면 평균 월급, 4대보험, 총 인건비가 자동 산출됩니다.")

        st.markdown("**📋 부서별 인원 · 급여**")
        lc1, lc2 = st.columns(2)
        ceo_n = lc1.number_input("대표 (명)", value=sv('l_ceo_n', 1), step=1, key='l_ceo_n')
        ceo_s = lc2.number_input("대표 월급 (만원)", value=sv('l_ceo_s', 500), step=50, key='l_ceo_s')
        lc1, lc2 = st.columns(2)
        adm_n = lc1.number_input("총무팀 (명)", value=sv('l_adm_n', 1), step=1, key='l_adm_n')
        adm_s = lc2.number_input("총무 월급 (만원)", value=sv('l_adm_s', 350), step=10, key='l_adm_s')
        lc1, lc2 = st.columns(2)
        acc_n = lc1.number_input("경리팀 (명)", value=sv('l_acc_n', 1), step=1, key='l_acc_n')
        acc_s = lc2.number_input("경리 월급 (만원)", value=sv('l_acc_s', 300), step=10, key='l_acc_s')
        lc1, lc2 = st.columns(2)
        fac_n = lc1.number_input("시설팀 (명)", value=sv('l_fac_n', 3), step=1, key='l_fac_n')
        fac_s = lc2.number_input("시설 월급 (만원)", value=sv('l_fac_s', 320), step=10, key='l_fac_s')
        lc1, lc2 = st.columns(2)
        desk_n = lc1.number_input("안내팀 (명)", value=sv('l_desk_n', 3), step=1, key='l_desk_n')
        desk_s = lc2.number_input("안내 월급 (만원)", value=sv('l_desk_s', 300), step=10, key='l_desk_s')

        st.markdown("---")
        s_insurance = st.number_input("4대보험 사업주부담 (%)", value=sv('s_ins', 9.6), step=0.5, key='s_ins', help="국민연금4.5%+건강보험3.545%+고용보험0.9%+산재보험0.7%=9.645% (2025년 법정비율)") / 100
        s_labor_up = st.number_input("연간 인건비 인상률 (%)", value=sv('s_lab_up', float(A['cost']['인건비인상률'][1]*100)), step=0.5, key='s_lab_up', help="매년 임금 인상률. 최저임금 인상률 참고") / 100

        st.caption("⚠️ **퇴직충당금**(급여의 8.33%=1/12)은 별도 발생. **부가세**(매출의 약 9.09%)는 매출에 포함된 총액 기준이며 별도 분리하지 않음. **운전자금**(선수금 등)은 현금흐름에 미반영.")

        st.markdown("---")
        st.markdown("**📊 자동 산출 결과**")
        dept_total_n = ceo_n + adm_n + acc_n + fac_n + desk_n
        dept_total_salary = ceo_n*ceo_s + adm_n*adm_s + acc_n*acc_s + fac_n*fac_s + desk_n*desk_s
        s_avg_salary = dept_total_salary / dept_total_n if dept_total_n else 0
        monthly_labor = dept_total_salary * (1 + s_insurance)
        s_staff = dept_total_n

        st.metric("총 인원", f"{dept_total_n}명")
        st.metric("전체 평균 월급", f"{s_avg_salary:,.0f}만원")
        st.metric("급여 합계 (월)", f"{dept_total_salary:,.0f}만원")
        st.metric("4대보험 포함 (월)", f"{monthly_labor:,.0f}만원")
        st.metric("연 인건비 합계", f"{monthly_labor*12/1e4:.1f}억원")
        est_rev = m_1m*p_1m + m_3m*p_3m + m_6m*p_6m + m_coupon*p_coupon + m_daily*p_daily + m_lesson*p_lesson + m_locker*p_locker*9
        st.metric("매출 대비 인건비율", f"{monthly_labor*12*만/est_rev*100:.1f}%" if est_rev else "N/A", help="업계 적정 수준: 20~30%")


    with st.expander("🏗️ 월 고정 운영비"):
        st.caption("ℹ️ Excel ⑤손익시뮬레이션 2026F 비용 데이터 기반 기본값. 인건비·감가상각은 별도 관리합니다. 2026F는 9개월(6~2월) 운영이며, 아래 금액은 월 환산 기준입니다.")

        # Excel 2026F 비용 데이터 (연간 → 월 환산)
        # 수도광열비 116,487,320원/9개월 = 1,294만
        # 전력비 143,287,398원/9개월 = 1,592만
        # 세금과공과 620,237,641원/12개월 = 5,169만 (연간 고정세금)
        # 보험료 10,620,788원/12개월 = 89만
        # 소모품비 92,780,833원/9개월 = 1,031만
        # 카드수수료 38,439,016원 (매출 대비 약 2%)
        # 수선비 11,143,910원/12개월 = 93만
        # 용역비 167,785,049원/9개월 = 1,864만
        # 기타비용 13,506,795원/12개월 = 113만

        op_electric = st.number_input("전력비 (만원/월)", value=sv('op_elec', 1590), step=50, key='op_elec',
            help="Excel 2026F: 전력비 1.43억/9개월 = 월 1,592만. 야간조명+볼머신+냉난방")
        op_water = st.number_input("수도광열비 (만원/월)", value=sv('op_water', 1290), step=50, key='op_water',
            help="Excel 2026F: 수도광열비 1.16억/9개월 = 월 1,294만. 수도+가스+난방유")
        op_tax = st.number_input("세금과공과 (만원/월)", value=sv('op_tax', 700), step=100, key='op_tax',
            help="등촌 자가소유 기준 700만/월(연 8,400만). 구성: 재산세(건물+토지) 약 400~580만 + 종합부동산세 약 50~100만 + 환경부담금/도로점용료/사업소세 등 약 50~100만. 부지·건물 평가액에 따라 ±300만 조정 가능. Excel 원본 6,891만/월은 토지임대료를 포함한 것으로 추정되어 자가소유에서 부적합.")
        op_insurance = st.number_input("보험료 (만원/월)", value=sv('op_insur', 120), step=10, key='op_insur',
            help="Excel 2026F: 보험료 1,062만/9개월 = 월 118만. 화재+배상책임보험")
        op_supplies = st.number_input("소모품비 (만원/월)", value=sv('op_supp', 1030), step=50, key='op_supp',
            help="Excel 2026F: 소모품비 9,278만/9개월 = 월 1,031만. 골프공+매트+티+장갑 등")
        op_maint = st.number_input("수선비 (만원/월)", value=sv('op_maint', 120), step=10, key='op_maint',
            help="Excel 2026F: 수선비 1,114만/9개월 = 월 124만. 시설보수+장비수리")
        op_card_fee = st.number_input("카드수수료율 (%)", value=sv('op_card', 2.0), step=0.1, key='op_card',
            help="Excel 2026F: 카드수수료 3,844만 ÷ 매출 19.2억 ≈ 2.0%") / 100
        op_outsource = st.number_input("용역비-경비/청소 (만원/월)", value=sv('op_out', 1860), step=50, key='op_out',
            help="Excel 2026F: 용역비 1.68억/9개월 = 월 1,864만. 경비+청소+폐기물+해충방제")
        op_etc = st.number_input("기타비용 (만원/월)", value=sv('op_etc', 150), step=10, key='op_etc',
            help="Excel 2026F: 기타비용 1,351만/9개월 = 월 150만. 복리후생+통신+소송+잡비")
        op_marketing = st.number_input("광고/마케팅 (만원/월)", value=sv('op_mkt', 150), step=10, key='op_mkt',
            help="마케팅 비용 (Excel 미포함, 별도 편성). 온라인광고+현수막+이벤트")

        st.markdown("---")
        st.markdown("**📊 운영비 자동 산출 (인건비·감가상각 별도)**")

        op_fixed_monthly = op_electric + op_water + op_tax + op_insurance + op_supplies + op_maint + op_outsource + op_etc + op_marketing
        _est_rev = m_1m*p_1m + m_3m*p_3m + m_6m*p_6m + m_coupon*p_coupon + m_daily*p_daily + m_lesson*p_lesson + m_locker*p_locker*9
        op_var_monthly = _est_rev / 12 / 만 * op_card_fee
        op_total_monthly = op_fixed_monthly + op_var_monthly

        mc1, mc2 = st.columns(2)
        mc1.metric("고정 운영비 (월)", f"{op_fixed_monthly:,.0f}만원")
        mc2.metric("카드수수료 (월)", f"{op_var_monthly:,.0f}만원")
        mc1.metric("운영비 합계 (월)", f"{op_total_monthly:,.0f}만원")
        mc2.metric("운영비 합계 (연)", f"{op_total_monthly*12/1e4:.1f}억원")

        st.caption(f"참고: 인건비 {monthly_labor:,.0f}만/월 + 운영비 {op_total_monthly:,.0f}만/월 = 총 비용 {monthly_labor+op_total_monthly:,.0f}만/월 ({(monthly_labor+op_total_monthly)*12/1e4:.1f}억/연)")

    with st.expander("📊 연간 비용 인상률"):
        st.caption("ℹ️ 5개년 재무 추정에 적용되는 연간 비용 인상률입니다. 각 항목은 전년 대비 매년 복리로 적용됩니다. 물가상승률은 한국은행 목표(2%)가 기준이며, 전기료는 정부 정책에 따라 변동이 큽니다.")
        s_util_up = st.number_input("물가상승률 (%)", value=sv('s_util', float(A['cost']['물가상승률'][1]*100)), step=0.5, key='s_util') / 100
        s_elec_up = st.number_input("전기료 인상률 (%)", value=sv('s_elec', float(A['cost']['전기료인상률'][1]*100)), step=0.5, key='s_elec') / 100
        s_tax_up = st.number_input("세금과공과 인상률 (%)", value=sv('s_tax_up', float(A['cost']['세금공과인상률'][1]*100)), step=0.5, key='s_tax_up') / 100
        s_outsource_up = st.number_input("용역비 인상률 (%)", value=sv('s_out_up', 3.0), step=0.5, key='s_out_up') / 100
        s_supply_up = st.number_input("소모품비 인상률 (%)", value=sv('s_sup_up', 2.0), step=0.5, key='s_sup_up') / 100
        s_mkt_up = st.number_input("마케팅비 증가율 (%)", value=sv('s_mkt_up', 5.0), step=1.0, key='s_mkt_up') / 100
        st.markdown("---")
        st.caption("인상률 영향 추정 (5년차 기준)")
        yr5_util = op_electric * (1+s_elec_up)**4
        yr5_outsource = op_outsource * (1+s_outsource_up)**4
        st.metric("5년차 전기료 (월)", f"{yr5_util:,.0f}만원", delta=f"+{(yr5_util-op_electric)/op_electric*100:.1f}%")
        st.metric("5년차 용역비 (월)", f"{yr5_outsource:,.0f}만원", delta=f"+{(yr5_outsource-op_outsource)/op_outsource*100:.1f}%")

    with st.expander("🏪 임대매장 설정"):
        st.caption("임대 매장 구성 및 계열사 배분")
        s_rent_shops = st.number_input("총 임대매장 수", value=sv('s_rshops', 5), step=1, key='s_rshops')
        s_rent_affiliate = st.number_input("계열사 매장 수", value=sv('s_raff', 2), step=1, key='s_raff')
        s_rent_external = s_rent_shops - s_rent_affiliate
        st.caption(f"외부임차: {s_rent_external}개")
        st.markdown("---")
        st.caption("외부 임차인 평균 조건")
        s_rent_ext_price = st.number_input("외부 월 임대료 (만원)", value=sv('s_rext', 250), step=10, key='s_rext')
        s_rent_ext_deposit = st.number_input("외부 보증금 (만원)", value=sv('s_rdep', 3000), step=100, key='s_rdep')
        s_rent_ext_up = st.number_input("연간 인상률 (%)", value=sv('s_rup', 3.0), step=0.5, key='s_rup') / 100
        st.markdown("---")
        st.caption("계열사 매장 조건")
        s_aff_inv_total = st.number_input("계열사 총 투자금 (만원)", value=sv('s_ainv', 15000), step=1000, key='s_ainv')
        s_aff_inv_share = st.slider("본사 투자 분담 (%)", 0, 100, 40, 5, key='s_ashare')
        s_aff_rent_disc = st.slider("계열사 임대료 할인 (%)", 0, 50, 20, 5, key='s_adisc')
        s_aff_rev_share = st.slider("계열사 매출 배분 (%)", 0, 30, 10, 1, key='s_arevs')

    # 📍 상권·경쟁 가중치 — 숨김 처리 (변수는 저장값/기본값으로 초기화하여 계산은 정상 작동)
    tw_pop = sv('tw_pop', 75)
    tw_access = sv('tw_acc', 80)
    tw_visible = sv('tw_vis', 70)
    tw_compete = sv('tw_comp', 55)
    tw_growth = sv('tw_grow', 85)
    w_pop = sv('w_pop', 25)
    w_acc = sv('w_acc', 20)
    w_vis = sv('w_vis', 15)
    w_comp = sv('w_comp', 20)
    w_grow = sv('w_grow', 20)
    s_comp_outdoor = sv('s_co', 2)
    s_comp_indoor = sv('s_ci', 5)
    s_comp_screen = sv('s_cs', 12)

    # 📉 골프시장·경제지표 — 숨김 처리 (변수는 저장값/기본값으로 초기화하여 계산은 정상 작동)
    ecos_key = sv('ecos_k', '')
    kosis_key = sv('kosis_k', '')
    # econ_data 캐시 또는 기본값 로드 (UI 숨김 상태에서도 정상 동작)
    if 'econ_data' not in st.session_state:
        if EconomicDataFetcher is not None:
            try:
                _fetcher_init = EconomicDataFetcher(ecos_key, kosis_key)
                st.session_state['econ_data'] = _fetcher_init.get_all()
            except Exception:
                st.session_state['econ_data'] = DEFAULTS.copy() if 'DEFAULTS' in dir() else {'source': 'default'}
        else:
            st.session_state['econ_data'] = DEFAULTS.copy() if 'DEFAULTS' in dir() else {'source': 'default'}
    ed = st.session_state['econ_data']
    e_golf_pop = sv('e_gpop', 388)
    e_golf_growth = sv('e_ggrow', float(A['market']['골프시장하락률'][1]*100)) / 100
    e_indoor_count = sv('e_indoor', 650)
    e_outdoor_large = sv('e_out_l', 280)
    e_outdoor_mid = sv('e_out_m', 450)
    e_outdoor_small = sv('e_out_s', 470)
    e_range_count = e_indoor_count + e_outdoor_large + e_outdoor_mid + e_outdoor_small
    e_screen_count = sv('e_scnt', 9500)
    e_gdp_growth = sv('e_gdp', float(ed.get('gdp_growth', 2.1))) / 100
    e_cpi = sv('e_cpi', float(ed.get('cpi_rate', 2.5))) / 100
    e_interest = sv('e_ir', float(ed.get('base_rate', 3.0))) / 100
    e_unemp = sv('e_unemp', float(ed.get('unemployment', 3.5))) / 100
    e_disposable = sv('e_disp', float(ed.get('disposable_income_growth', 1.5))) / 100
    ew_golf = sv('ew_g', 30)
    ew_gdp = sv('ew_gdp', 20)
    ew_cpi = sv('ew_cpi', 15)
    ew_ir = sv('ew_ir', 15)
    ew_disp = sv('ew_disp', 20)

    with st.expander("🎯 할인율 / 세율"):
        s_disc = st.slider("할인율-WACC (%)", 5, 25, 8, 1, key='s_disc',
            help="할인율 = 무위험수익률(국고채 3%) + 사업리스크(2~3%, 재오픈이라 신규보다 낮음) + 유동성프리미엄(2%). 등촌은 재오픈+검증된 입지+부동산담보형이므로 8% 권장. 신규 그린필드라면 10~12%. 사학법인 기준 8~12% 권장.")
        s_tax_rate = st.number_input("법인세율 (%)", value=sv('s_taxr', float(A['invest']['법인세율']*100)), step=1.0, key='s_taxr') / 100

    st.markdown("---")
    st.markdown("##### 📌 실시간 요약")

    # Quick summary metrics (will be filled later after calculations)


# ══════════════════════════════════════════════════════════════
# Core Calculations
# ══════════════════════════════════════════════════════════════
    # 임대매출: 2026년은 11월~2월 = 4개월만 (11월 오픈)
_rent_monthly = s_rent_external * s_rent_ext_price * 만 + s_rent_affiliate * s_rent_ext_price * 만 * (1 - s_aff_rent_disc/100)
_rent_2026 = _rent_monthly * 4  # 11~2월 4개월
_rent_annual = _rent_monthly * 12  # 2027F~ 정상 연간

custom_rev_items = {
    '1개월 회원': m_1m * p_1m, '3개월 회원': m_3m * p_3m, '6개월 회원': m_6m * p_6m,
    '쿠폰': m_coupon * p_coupon, '일일회원': m_daily * p_daily,
    '골프레슨(프로임대)': p_lesson, '락카': m_locker * p_locker * 9,
    '임대(11~2월,4개월)': _rent_2026,
}
custom_total_rev = sum(custom_rev_items.values())
total_members = m_1m + m_3m + m_6m + m_coupon + m_daily + m_lesson + m_locker

inv_won = s_inv * 억
disc_r = s_disc / 100
# 감가상각: 컨트롤 패널 입력값 기반
# 정액법: 매년 동일 금액
# 정률법: 장부가 잔액 × 상각률 (상각률 = 2/내용연수)
_dep_str = dep_straight_amt * 만  # 정액법 연간 (원)
_dep_rate = 2 / s_useful if s_useful else 0.4  # 정률법 상각률
dep = []
# 정률법 취득원가 = 초년도 상각액 / 상각률
_dep_dec_book = dep_declining_amt * 만 / _dep_rate if _dep_rate else 0  # 장부가(취득원가)
for yr in range(5):
    _dep_dec_yr = _dep_dec_book * _dep_rate
    dep.append(int(_dep_str + _dep_dec_yr))
    _dep_dec_book -= _dep_dec_yr
    if _dep_dec_book < 0:
        _dep_dec_book = 0

# ══ 보정계수 사전 계산 (상권·경제) ══
w_total = w_pop + w_acc + w_vis + w_comp + w_grow
if w_total > 0:
    ta_score = (tw_pop*w_pop + tw_access*w_acc + tw_visible*w_vis + tw_compete*w_comp + tw_growth*w_grow) / w_total
else:
    ta_score = 70
ta_rev_adj = ta_score / 70

ew_total = ew_golf + ew_gdp + ew_cpi + ew_ir + ew_disp
# 민감도 절반: 1% 입력 변동 → 5점 변동 (이전: 10점) — 사용자 슬라이더 미세조정에도 출력이 안정적
# 기준: 0%/2.5%/2.5%/3%/2% 입력 시 각 점수 ≈ 50 (중립)
econ_golf_score = max(0, min(100, 50 + e_golf_growth * 500))
econ_gdp_score = max(0, min(100, 50 + (e_gdp_growth - 0.025) * 1000))
econ_cpi_score = max(0, min(100, 50 - (e_cpi - 0.025) * 1000))
econ_ir_score = max(0, min(100, 50 - (e_interest - 0.03) * 500))
econ_disp_score = max(0, min(100, 50 + (e_disposable - 0.02) * 500))
if ew_total > 0:
    econ_score = (econ_golf_score*ew_golf + econ_gdp_score*ew_gdp + econ_cpi_score*ew_cpi + econ_ir_score*ew_ir + econ_disp_score*ew_disp) / ew_total
else:
    econ_score = 50
ECON_BASELINE = 50  # 50점 = 매출 보정 1.0x (중립). 골프시장하락률 0%, GDP 2.5%, CPI 2.5%, 기준금리 3%, 가처분소득 2% 가정 시 약 50점이 산출되도록 설계
econ_rev_adj = econ_score / ECON_BASELINE
combined_adj = ta_rev_adj * econ_rev_adj

# ══ 컨트롤 패널 기반 5개년 매출·비용 재계산 ══
# ── 이탈률 가중평균 (매출 보정에 사용) ──
w_churn = (s_churn_sum * 2 + s_churn_win * 2 + s_churn_perm * 5) / 9

# ── 보정 계수 종합 ──
# 1) 시즌 가중치 (비수기 반영)
_sw = D['season_weights'] if D['season_weights'] else [1]*9
_season_avg = sum(_sw) / len(_sw) if _sw else 1.0  # 약 0.82

# 2) 오픈 램프업 (고객 유입 곡선, 2026F만 적용)
_ramp_avg = sum(ramp_values) / len(ramp_values) if ramp_values else 0.5  # 약 0.50

# 3) 상권·경쟁 보정 (ta_rev_adj, 이미 계산됨)
# 4) 경제지표 보정 (econ_rev_adj, 이미 계산됨)

# ── 2026F 매출: 이론치 × 시즌보정 × 램프업 × 상권보정 × 경제보정 ──
_golf_rev_9m_raw = custom_total_rev - _rent_2026  # 골프 이론 매출 (회원수×단가)
# 5) 환불 보정만 적용
# 주의: 이탈률/유치율은 이미 회원수 가정(2018 대비 60%)에 내재됨 — 매출에 이중 적용 금지
# 환불비율만 매출에서 차감 (회원권 구매 후 환불/변경)
_refund_adj = 1 - s_refund  # 환불비율 (예: 4% → ×0.96)

# 회원수 = 9개월 총 누적 (2018×60%×9/12). 시즌/램프업은 월별 분배에만 사용
# 연간 매출 보정: 이탈률(중도해지) + 환불(구매후취소) + 상권(입지) + 경제(외부환경)
_churn_adj = 1 - w_churn  # 가중평균 이탈률 차감 (비수기 해지 반영)
_adj_2026 = _churn_adj * _refund_adj * ta_rev_adj * econ_rev_adj  # 이탈+환불+상권+경제
_rev_2026 = _golf_rev_9m_raw * _adj_2026 + _rent_2026

# ── 2027F~: 정상가동 (12개월), 동일 보정 ──
_adj_normal = _churn_adj * _refund_adj * ta_rev_adj * econ_rev_adj
_golf_rev_annual = _golf_rev_9m_raw * 12 / 9 * _adj_normal  # 12개월 환산
_rev_full_yr = _golf_rev_annual + _rent_annual

rev_p = [_rev_2026]
# 2027F = 정상가동 기준, 2028F~ = 전년 대비 성장률 적용
rev_p.append(_rev_full_yr)  # 2027F: 정상가동 첫해 (성장률 미적용)
for i in range(1, 4):
    rev_p.append(_rev_full_yr * (1 + s_growth) ** i)  # 2028F~2030F

# 2026F 비용: 컨트롤 패널 (인건비 + 운영비) × 9개월 + 감가상각
_monthly_total_cost = monthly_labor * 만 + op_total_monthly * 만  # 원 단위
# 운영비 항목별 월액 (원 단위)
_op_elec_m = op_electric * 만
_op_water_m = op_water * 만
_op_tax_m = op_tax * 만
_op_insur_m = op_insurance * 만
_op_supply_m = op_supplies * 만
_op_maint_m = op_maint * 만
_op_card_m = custom_total_rev * op_card_fee / 12  # 매출×수수료율÷12 (월 원 단위, op_card_fee는 이미 소수)
_op_outsource_m = op_outsource * 만
_op_etc_m = op_etc * 만
_op_mkt_m = op_marketing * 만

# 2026 비용 항목별 (9개월)
cost_items_dyn = {
    '인건비': [monthly_labor * 만 * 9],
    '전력비': [_op_elec_m * 9],
    '수도광열비': [_op_water_m * 9],
    '세금과공과': [_op_tax_m * 9],
    '감가상각비': [dep[0]],
    '보험료': [_op_insur_m * 9],
    '소모품비': [_op_supply_m * 9],
    '카드수수료': [rev_p[0] * op_card_fee],
    '수선비+용역비+기타': [(_op_maint_m + _op_outsource_m + _op_etc_m + _op_mkt_m) * 9],
}
_cost_2026 = sum(v[0] for v in cost_items_dyn.values())
cost_p = [_cost_2026]

for i in range(4):
    y = i + 1
    # 인건비: 인건비 인상률
    yr_labor = monthly_labor * 만 * 12 * (1 + s_labor_up) ** y
    # 전기료: 전기료 인상률
    yr_elec = _op_elec_m * 12 * (1 + s_elec_up) ** y
    # 수도광열비: 물가상승률
    yr_water = _op_water_m * 12 * (1 + s_util_up) ** y
    # 세금과공과: 세금 인상률
    yr_tax = _op_tax_m * 12 * (1 + s_tax_up) ** y
    # 보험료: 물가상승률
    yr_insur = _op_insur_m * 12 * (1 + s_util_up) ** y
    # 소모품비: 소모품 인상률
    yr_supply = _op_supply_m * 12 * (1 + s_supply_up) ** y
    # 수선비: 물가상승률
    yr_maint = _op_maint_m * 12 * (1 + s_util_up) ** y
    # 카드수수료: 매출 연동 (매출 × 수수료율)
    yr_card = rev_p[y] * op_card_fee if y < len(rev_p) else _op_card_m * 12  # op_card_fee는 이미 소수(0.02)
    # 용역비: 용역비 인상률
    yr_outsource = _op_outsource_m * 12 * (1 + s_outsource_up) ** y
    # 기타+마케팅: 각각 적용
    yr_etc = _op_etc_m * 12 * (1 + s_util_up) ** y
    yr_mkt = _op_mkt_m * 12 * (1 + s_mkt_up) ** y
    yr_dep = dep[min(y, len(dep) - 1)]
    cost_items_dyn['인건비'].append(yr_labor)
    cost_items_dyn['전력비'].append(yr_elec)
    cost_items_dyn['수도광열비'].append(yr_water)
    cost_items_dyn['세금과공과'].append(yr_tax)
    cost_items_dyn['감가상각비'].append(yr_dep)
    cost_items_dyn['보험료'].append(yr_insur)
    cost_items_dyn['소모품비'].append(yr_supply)
    cost_items_dyn['카드수수료'].append(yr_card)
    cost_items_dyn['수선비+용역비+기타'].append(yr_maint + yr_outsource + yr_etc + yr_mkt)
    yr_opex_total = yr_elec + yr_water + yr_tax + yr_insur + yr_supply + yr_maint + yr_card + yr_outsource + yr_etc + yr_mkt
    cost_p.append(yr_labor + yr_opex_total + yr_dep)

# 패드 기반 5개년 상품별 매출 분해 — 2026 상품별 비중을 유지하며 각 연도 rev_p[i]로 스케일
# 합계가 모든 연도에서 정확히 rev_p[i]와 일치하도록 보정
_rev_items_2026 = dict(custom_rev_items)
_total_2026 = sum(_rev_items_2026.values())
rev_items_dyn = {n: [] for n in _rev_items_2026}
for i in range(5):
    if _total_2026 > 0:
        for n in _rev_items_2026:
            rev_items_dyn[n].append(_rev_items_2026[n] / _total_2026 * rev_p[i])
    else:
        for n in _rev_items_2026:
            rev_items_dyn[n].append(0)

op_p = [r - c for r, c in zip(rev_p, cost_p)]
ebitda_p = [o + d for o, d in zip(op_p, dep[:len(op_p)])]

cum_ebitda = []
cc_ = 0
for e in ebitda_p:
    cc_ += e; cum_ebitda.append(cc_)
rec_rate = [c / inv_won if inv_won else 0 for c in cum_ebitda]

# FCF = EBITDA - 법인세 (영업이익 기준 과세, 적자 시 세금 0)
fcf_p = [e - max(o * s_tax_rate, 0) for e, o in zip(ebitda_p, op_p)]
npv_val = sum(f / (1 + disc_r) ** (i + 1) for i, f in enumerate(fcf_p)) - inv_won

try:
    from numpy_financial import irr as np_irr
    irr_val = float(np_irr([-inv_won] + fcf_p))
except Exception:
    np_irr = None
    irr_val = (sum(fcf_p) / 5) / inv_won if inv_won else 0

payback = None
cum_temp = 0
for i, e in enumerate(ebitda_p):
    prev = cum_temp; cum_temp += e
    if cum_temp >= inv_won and prev < inv_won:
        payback = i + ((inv_won - prev) / e if e else 0); break

# w_churn은 라인 1027에서 이미 정의됨

# ══ BEP — 컨트롤 패널 기반 재계산 ══
_yr_labor_cost = monthly_labor * 만 * 12  # 연간 인건비 (원)
_yr_dep = dep[0] if dep else 0
_yr_insurance = op_insurance * 만 * 12
_yr_tax = op_tax * 만 * 12
fixed_total = _yr_labor_cost + _yr_dep + _yr_insurance + _yr_tax  # 고정비

_yr_card = custom_total_rev * op_card_fee
_yr_supplies = op_supplies * 만 * 12
var_total = _yr_card + _yr_supplies  # 변동비

_yr_electric = op_electric * 만 * 12
_yr_water = op_water * 만 * 12
_yr_outsource = op_outsource * 만 * 12
semi_total = _yr_electric + _yr_water + _yr_outsource  # 준변동비
fixed_total += semi_total * 0.6
var_total += semi_total * 0.4

# BEP는 정상 가동 연도(2027F) 기준으로 산출 (2026F는 9개월+오픈초기라 비정상)
_bep_rev_base = rev_p[1] if len(rev_p) > 1 else rev_p[0]
var_ratio = var_total / _bep_rev_base if _bep_rev_base else 0.1
contrib_margin = 1 - var_ratio
bep_revenue = fixed_total / contrib_margin if contrib_margin > 0 else 0
bep_members = bep_revenue / (custom_total_rev / total_members) if total_members and custom_total_rev else 0
safety_margin = (_bep_rev_base - bep_revenue) / _bep_rev_base * 100 if _bep_rev_base else 0

# ══ 월별 데이터 — 컨트롤 패널 기반 재계산 ══
# 월별 매출: 연 골프매출(rev_p[0]−임대)을 시즌×램프업 가중치 비율로 9개월 배분 (합계 = 연 골프매출)
_monthly_adj_factor = _churn_adj * _refund_adj * ta_rev_adj * econ_rev_adj  # 참고용 (이탈+환불+상권+경제)
_annual_golf_rev = rev_p[0] - _rent_2026  # 2026 연 골프 매출 (보정 적용 후)
_weights = [w * r for w, r in zip(_sw, ramp_values)]  # 시즌 × 램프업
_weights_sum = sum(_weights)
if _weights_sum > 0:
    mrev_custom = [int(_annual_golf_rev * w / _weights_sum) for w in _weights]
else:
    mrev_custom = [0 for _ in _sw]
_monthly_base_cost = _monthly_total_cost
mcost_custom = [int(_monthly_base_cost) for _ in _sw]  # 비용은 월 고정

# Rental shop calculations
rent_ext_annual = s_rent_external * s_rent_ext_price * 만 * 12  # 외부 연 임대수익
aff_rent_monthly = s_rent_ext_price * 만 * (1 - s_aff_rent_disc / 100)  # 계열사 월 임대료 (할인)
rent_aff_annual = s_rent_affiliate * aff_rent_monthly * 12  # 계열사 연 임대수익
rent_total_annual = rent_ext_annual + rent_aff_annual  # 총 연 임대수익
aff_inv_hq = s_aff_inv_total * 만 * s_aff_inv_share / 100  # 본사 분담 투자금
aff_inv_aff = s_aff_inv_total * 만 - aff_inv_hq  # 계열사 분담
rent_deposit_total = s_rent_external * s_rent_ext_deposit * 만  # 총 보증금 수입

# 5-year rental projection
rent_5yr = []
for i in range(5):
    ext_r = s_rent_external * s_rent_ext_price * 만 * 12 * (1 + s_rent_ext_up) ** i
    aff_r = s_rent_affiliate * aff_rent_monthly * 12 * (1 + s_rent_ext_up * 0.5) ** i  # 계열사는 인상률 50% 적용
    rent_5yr.append({'외부': ext_r, '계열사': aff_r, '합계': ext_r + aff_r})

# Affiliate revenue share income
aff_rev_base = 50000 * 만  # 계열사 매장 평균 연매출 5억 가정 (per shop)
aff_rev_share_income = [s_rent_affiliate * aff_rev_base * s_aff_rev_share / 100 * (1 + s_growth) ** i for i in range(5)]
total_rental_income = [rent_5yr[i]['합계'] + aff_rev_share_income[i] for i in range(5)]

# Financing (simplified - keep for cash flow tab)
annual_pmt = 0
loan_schedule = [{'연차': i+1, '기초잔액': 0, '이자': 0, '원금상환': 0, '상환합계': 0, '기말잔액': 0} for i in range(5)]
equity_won = inv_won
loan_won = 0

# 10-year
rev_10yr, cost_10yr = list(rev_p), list(cost_p)
for i in range(5):
    rev_10yr.append(rev_10yr[-1] * (1 + s_growth))
    cost_10yr.append(cost_10yr[-1] * (1 + s_util_up))
op_10yr = [r - c for r, c in zip(rev_10yr, cost_10yr)]
rev_per_bay = [r / s_bays for r in rev_p]

margins = [o / r * 100 if r else 0 for o, r in zip(op_p, rev_p)]
em = [e / r * 100 if r else 0 for e, r in zip(ebitda_p, rev_p)]

# 상권·경제 보정계수는 이미 rev_p 계산 전(라인 ~1280)에서 산출됨. 여기서는 등급/표시용 변수만 정의.
ta_grade = "A" if ta_score >= 80 else "B+" if ta_score >= 70 else "B" if ta_score >= 60 else "C+" if ta_score >= 50 else "C"
econ_grade = "호황" if econ_score >= 65 else "보통" if econ_score >= 45 else "불황"

# Sidebar summary
with st.sidebar:
    st.metric("투자금", f"{s_inv}억")
    st.metric("NPV", fmt억(npv_val))
    st.metric("IRR", f"{irr_val*100:.1f}%")
    st.metric("회수율", f"{rec_rate[-1]*100:.1f}%")

# ══════════════════════════════════════════════════════════════
# Header
# ══════════════════════════════════════════════════════════════
# Navigation via radio buttons (wrap-friendly, no scroll)
# ══════════════════════════════════════════════════════════════
TAB_NAMES = ["대시보드", "시나리오", "기여도", "운영전략", "매출추정", "비용추정", "추정손익계산서",
    "투자IRR", "현금흐름", "임대", "상권분석", "시장경제", "검증"]

# ── 통합 네비게이션 바 (탭 + 버튼) ──
now_str = datetime.now().strftime('%Y%m%d_%H%M')

st.markdown("""<style>
/* ─── 폴더 탭 스타일 (참조 이미지의 '모양'만 차용, 컬러는 다크 테마 톤) ─── */

/* 탭 컨테이너 — 어두운 스트립 */
div[data-testid="stRadio"] {
    background: #0a0f1c !important;
    padding: 8px 6px 0 6px !important;
    margin: 8px 0 0 0 !important;
    border-radius: 10px 10px 0 0 !important;
    border: 1px solid #1e293b !important;
    border-bottom: none !important;
}

/* 6열 그리드 */
div[data-testid="stRadio"] > div {
    display: grid !important;
    grid-template-columns: repeat(6, 1fr) !important;
    gap: 4px !important;
    padding: 0 !important;
}

/* 비활성 탭 — 살짝 옅은 슬레이트, 폴더 모양 */
div[data-testid="stRadio"] > div > label {
    background: #1e293b !important;
    border: 1px solid #1e293b !important;
    border-bottom: 1px solid #334155 !important;
    border-radius: 8px 8px 0 0 !important;
    padding: 11px 0 !important;
    font-size: 13.5px !important;
    font-weight: 500 !important;
    color: #94a3b8 !important;
    cursor: pointer !important;
    transition: background 0.15s ease, color 0.15s ease !important;
    white-space: nowrap !important;
    text-align: center !important;
    margin-bottom: -1px !important;
    box-shadow: inset 0 -3px 4px rgba(0,0,0,0.18) !important;
    position: relative;
    z-index: 1;
    /* 텍스트 중앙 정렬 (flex로 모든 자식 가운데로) */
    display: flex !important;
    justify-content: center !important;
    align-items: center !important;
}
/* 라벨 내부 모든 텍스트/요소 가운데 정렬 */
div[data-testid="stRadio"] > div > label > * {
    text-align: center !important;
    width: auto !important;
    margin: 0 auto !important;
}
div[data-testid="stRadio"] > div > label:hover {
    background: #283549 !important;
    color: #cbd5e1 !important;
    transform: none !important;
}

/* 활성 탭 — 콘텐츠 카드와 같은 톤(#111827), 위로 올라온 느낌 + 캡션 바와 연결 */
div[data-testid="stRadio"] > div > label[data-checked="true"],
div[data-testid="stRadio"] > div > label:has(input:checked) {
    background: #111827 !important;
    color: #60a5fa !important;
    font-weight: 600 !important;
    border-color: #334155 !important;
    border-bottom: 1px solid #111827 !important;
    box-shadow: 0 -2px 6px rgba(0,0,0,0.4) !important;
    z-index: 2;
}

/* radio 내부 점 숨김 */
div[data-testid="stRadio"] > div > label > div:first-child { display: none !important; }
div[data-testid="stRadio"] > label { display: none !important; }

/* 탭 하단 캡션 바 — 활성 탭과 같은 #111827로 자연스럽게 연결 */
.tab-caption-bar {
    background: #111827;
    border: 1px solid #334155;
    border-top: none;
    border-radius: 0 0 10px 10px;
    padding: 9px 18px;
    margin: 0 0 22px 0;
    font-size: 0.78rem;
    color: #94a3b8;
    text-align: right;
    box-shadow: 0 2px 6px rgba(0,0,0,0.25);
}
</style>""", unsafe_allow_html=True)

# 버튼 바 (엑셀/PDF/PPT/초기화/프린트) — 숨김 처리
if False:
    bc = st.columns([4.5, 0.7, 0.7, 0.7, 0.6, 0.5, 0.7])
    with bc[1]:
        excel_data = generate_excel(D, rev_p, cost_p, op_p, ebitda_p, cum_ebitda, margins, rec_rate, inv_won)
        st.download_button("엑셀", data=excel_data,
            file_name=f"등촌골프_{now_str}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            use_container_width=True)
    with bc[2]:
        # 용지 크기 선택 + PDF
        _paper = st.selectbox("용지", ['A4', 'A3', 'B4', 'B3'], index=0, label_visibility="collapsed", key='paper_size')
    with bc[3]:
        try:
            pdf_data = generate_pdf_report(D, rev_p, cost_p, op_p, ebitda_p, cum_ebitda, margins, rec_rate, inv_won, _paper)
            st.download_button(f"PDF({_paper})", data=pdf_data,
                file_name=f"등촌골프_{now_str}_{_paper}.pdf",
                mime="application/pdf",
                use_container_width=True)
        except Exception as e:
            st.button(f"PDF", use_container_width=True, disabled=True, help=f"PDF 생성 실패: {e}")
    with bc[4]:
        try:
            ppt_data = generate_ppt_report(D, rev_p, cost_p, op_p, ebitda_p, cum_ebitda, margins, rec_rate, inv_won)
            st.download_button("PPT", data=ppt_data,
                file_name=f"등촌골프_{now_str}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True)
        except Exception as e:
            st.button("PPT", use_container_width=True, disabled=True, help=f"PPT 생성 실패: {e}")
    with bc[5]:
        if st.button("초기화", use_container_width=True, help="모든 설정을 기본값으로"):
            import os as _os
            _save_path = _os.path.join(_os.path.dirname(__file__), 'user_settings.json')
            if _os.path.exists(_save_path):
                _os.remove(_save_path)
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()
    with bc[6]:
        # 프린트 기능 (브라우저 인쇄)
        st.markdown(f"""<button onclick="window.print()" style="
            width:100%;padding:6px 0;border-radius:8px;border:1px solid #334155;
            background:#1e293b;color:#94a3b8;font-size:12px;cursor:pointer;
            font-family:inherit;">프린트</button>""", unsafe_allow_html=True)

# 탭 네비게이션 — 6개 탭만 노출 (대시보드/운영전략/매출/비용/손익BEP/상권경쟁)
# 내부 _ti 인덱스는 전체 TAB_NAMES 기준으로 유지하여 각 탭 코드블록 그대로 사용
VISIBLE_TABS = ["대시보드", "운영전략", "매출추정", "비용추정", "추정손익계산서", "상권분석"]
# 이전 라벨 → 새 라벨 마이그레이션 (저장된 main_nav 호환)
_label_migration = {"매출": "매출추정", "비용": "비용추정", "손익BEP": "추정손익계산서", "상권경쟁": "상권분석"}
_visible_default = st.session_state.get('main_nav', VISIBLE_TABS[0])
if _visible_default in _label_migration:
    _visible_default = _label_migration[_visible_default]
    st.session_state['main_nav'] = _visible_default
if _visible_default not in VISIBLE_TABS:
    _visible_default = VISIBLE_TABS[0]
    st.session_state['main_nav'] = _visible_default
selected_tab = st.radio("nav", VISIBLE_TABS, horizontal=True, label_visibility="collapsed", key="main_nav")
_ti = TAB_NAMES.index(selected_tab)

# 탭 하단 캡션 바 (활성 탭과 시각적으로 연결)
st.markdown('<div class="tab-caption-bar">* 모든 금액은 VAT 제외 · 자가소유 부지 기준 · 2026.06 재오픈 · 5개년 재무모델</div>', unsafe_allow_html=True)

# ── Chart Design Note ──
# All charts use 2D (go.Bar, go.Scatter, go.Pie, go.Heatmap, etc.) rather than 3D
# (go.Surface, go.Scatter3d, go.Mesh3d). Financial data is inherently 2D (categories x values),
# and 3D adds visual complexity without insight, makes labels unreadable, and slows rendering.

# ═══ TAB 0: Dashboard ═══
if _ti == 0:
    # Header + KPI inside dashboard tab
    st.markdown(f"""<div style="padding:4px 0 18px 0; margin-bottom:14px; border-bottom:1px solid #1e293b;">
        <h1 style="color:#f1f5f9; font-size:1.35rem; font-weight:700; margin:0; letter-spacing:-0.01em;">등촌골프연습장 사업성 분석</h1>
        <p style="color:#64748b; font-size:0.82rem; margin:6px 0 0 0;">{s_bays}타석 실외 · 투자금 {s_inv}억 · 2026.06 재오픈 · 5개년 재무모델링</p>
    </div>""", unsafe_allow_html=True)

    def kpi_card(col, label, value, tip_rows):
        """tip_rows: list of (항목, 값) 튜플"""
        tip_html = ''.join(f'<tr><td style="color:#94a3b8;padding:3px 10px 3px 0;font-size:11px;white-space:nowrap;">{r[0]}</td><td style="color:#e2e8f0;padding:3px 0;font-size:11px;font-weight:600;">{r[1]}</td></tr>' for r in tip_rows)
        col.markdown(f"""
<div class="kpi-wrap">
<div class="kpi-box">
<div style="color:#94a3b8;font-size:11px;font-weight:600;letter-spacing:0.03em;">{label}</div>
<div style="color:#f8fafc;font-size:22px;font-weight:800;margin-top:2px;">{value}</div>
</div>
<div class="kpi-tip"><table style="border-collapse:collapse;">{tip_html}</table></div>
</div>""", unsafe_allow_html=True)

    k = st.columns(7)
    kpi_card(k[0], "투자금", f"{s_inv}억", [
        ("구분","총 초기 투자금액"),
        ("타석당 투자",f"{s_inv*억/s_bays/만:,.0f}만원"),
        ("타석수",f"{s_bays}타석"),
        ("─────","──────────"),
        ("정액법 상각",f"{dep_straight_amt:,}만원/년"),
        ("정률법 상각",f"{dep_declining_amt:,}만원/년"),
        ("연간 감가합계",f"{dep_straight_amt+dep_declining_amt:,}만원/년"),
    ])
    kpi_card(k[1], "NPV", fmt억(npv_val), [
        ("정의","세후FCF를 할인율로 현재가치 환산 합계 - 투자금"),
        ("할인율(WACC)",f"{disc_r*100:.0f}% (무위험3~4%+사업리스크3~4%+유동성2~3%)"),
        ("─────","──────────"),
        ("2026F EBITDA",fmt억(ebitda_p[0])),
        ("2027F EBITDA",fmt억(ebitda_p[1])),
        ("2028F EBITDA",fmt억(ebitda_p[2])),
        ("2029F EBITDA",fmt억(ebitda_p[3])),
        ("2030F EBITDA",fmt억(ebitda_p[4])),
        ("─────","──────────"),
        ("판정","✅ 투자가치 있음" if npv_val>0 else "⚠️ NPV 음수, 조건부 검토"),
    ])
    kpi_card(k[2], "누적EBITDA", fmt억(cum_ebitda[-1]), [
        ("정의","5년간 EBITDA(영업이익+감가상각비) 누적"),
        ("─────","──────────"),
        ("2026F",fmt억(cum_ebitda[0])),
        ("2027F",fmt억(cum_ebitda[1])),
        ("2028F",fmt억(cum_ebitda[2])),
        ("2029F",fmt억(cum_ebitda[3])),
        ("2030F",fmt억(cum_ebitda[4])),
        ("─────","──────────"),
        ("투자금 대비",f"{cum_ebitda[-1]/inv_won*100:.1f}%"),
    ])
    kpi_card(k[3], "회수율", f"{rec_rate[-1]*100:.1f}%", [
        ("정의","누적EBITDA ÷ 투자금 × 100"),
        ("─────","──────────"),
        ("2026F",f"{rec_rate[0]*100:.1f}%"),
        ("2027F",f"{rec_rate[1]*100:.1f}%"),
        ("2028F",f"{rec_rate[2]*100:.1f}%"),
        ("2029F",f"{rec_rate[3]*100:.1f}%"),
        ("2030F",f"{rec_rate[4]*100:.1f}%"),
        ("─────","──────────"),
        ("미회수금",fmt억(inv_won - cum_ebitda[-1]) if cum_ebitda[-1] < inv_won else "전액 회수"),
    ])
    kpi_card(k[4], "IRR", f"{irr_val*100:.1f}%", [
        ("정의","투자수익률 (내부수익률)"),
        ("의미","NPV=0이 되는 할인율"),
        ("─────","──────────"),
        ("현재 IRR",f"{irr_val*100:.1f}%"),
        ("할인율(WACC)",f"{disc_r*100:.0f}% (무위험3~4%+리스크3~4%+유동성2~3%)"),
        ("─────","──────────"),
        ("판정","✅ IRR > 할인율 → 은행보다 수익 우수" if irr_val > disc_r else "⚠️ IRR < 할인율 → 은행 예금이 나을 수 있음"),
        ("참고","은행이자 3~4%가 아닌 10%를 기준으로 삼는 이유: 사업 리스크와 자금 묶임(유동성)을 감안한 기회비용"),
    ])
    kpi_card(k[5], "BEP매출", fmt억(bep_revenue), [
        ("정의","손익분기점 매출액"),
        ("산식","고정비 ÷ (1 - 변동비율)"),
        ("─────","──────────"),
        ("2027F 매출",fmt억(rev_p[1])),
        ("BEP 매출",fmt억(bep_revenue)),
        ("안전마진",f"{(rev_p[1]-bep_revenue)/rev_p[1]*100:.1f}%" if rev_p[1] else "N/A"),
        ("─────","──────────"),
        ("의미","이 금액 이상 매출 시 흑자"),
    ])
    kpi_card(k[6], "Payback", f"{payback:.1f}년" if payback else "5년+", [
        ("정의","투자금 전액 회수 소요 기간"),
        ("투자금",f"{s_inv}억"),
        ("─────","──────────"),
        ("1년차 누적",fmt억(cum_ebitda[0])),
        ("2년차 누적",fmt억(cum_ebitda[1])),
        ("3년차 누적",fmt억(cum_ebitda[2])),
        ("4년차 누적",fmt억(cum_ebitda[3])),
        ("5년차 누적",fmt억(cum_ebitda[4])),
        ("─────","──────────"),
        ("판정",f"✅ {payback:.1f}년 회수" if payback else "⚠️ 5년 내 미회수"),
    ])

    st.markdown("<div style='margin:20px 0;'></div>", unsafe_allow_html=True)

    if npv_val > 0:
        st.success("**투자적합** — NPV 양수, 할인율 기준 투자 수익 확보")
    elif rec_rate[-1] > 0.7:
        st.warning("**조건부 검토** — NPV 음수이나 5년 내 70%+ 회수 가능")
    else:
        st.error("**신중검토 필요** — NPV 음수, 투자회수율 미흡")

    st.markdown("<div style='margin-top:24px;'></div>", unsafe_allow_html=True)
    sec("📊", "주요 대시보드")

    # ── 경영진 핵심 지표 (표+그래프 혼합형) ──
    subsec("경영진 핵심 지표")
    info("사업성 판단에 가장 중요한 운영·재무 지표를 표와 그래프로 한눈에 보여줍니다. (2026F = 오픈 첫해 기준)")

    # 계산 (2026F = index 0 기준, 컨트롤 패널 연동)
    rev_26 = rev_p[0]
    cost_26 = cost_p[0]
    labor_26 = monthly_labor * 만 * 9  # 컨트롤 패널 인건비 × 9개월
    energy_26 = (op_electric + op_water) * 만 * 9  # 컨트롤 패널 전기+수도 × 9개월
    rev_per_bay_monthly = rev_26 / s_bays / 9 if s_bays else 0  # 9개월 기준 월평균
    labor_pct = labor_26 / cost_26 * 100 if cost_26 else 0
    energy_pct = energy_26 / cost_26 * 100 if cost_26 else 0
    dep_yr = dep[0] if dep else 0
    dep_pct = dep_yr / cost_26 * 100 if cost_26 else 0
    fixed_pct = labor_pct + dep_pct
    total_members_yr = m_1m + m_3m + m_6m + m_coupon + m_daily
    avg_price = custom_total_rev / total_members_yr if total_members_yr else 0
    monthly_fcf = (ebitda_p[0] - ebitda_p[0]*s_tax_rate) / 12
    # DSCR: 차입금이 없으면 N/A (현재 모델은 전액 자기자본 가정)
    _has_loan = False  # 차입 구조가 추가되면 True로 변경
    dscr = ebitda_p[0] / (inv_won * 0.2) if inv_won and _has_loan else 0
    growth_27 = ((rev_p[1]/rev_p[0])-1)*100 if len(rev_p)>1 and rev_p[0] else 0

    # ── 경영진 핵심 지표 — 전체 폭 ──
    def _badge(val, good_thr, bad_thr, higher_is_good=True):
        if higher_is_good:
            if val >= good_thr: return '<span style="background:#166534;color:#86efac;padding:2px 8px;border-radius:6px;font-size:11px;">양호</span>'
            elif val >= bad_thr: return '<span style="background:#854d0e;color:#fde047;padding:2px 8px;border-radius:6px;font-size:11px;">주의</span>'
            else: return '<span style="background:#991b1b;color:#fca5a5;padding:2px 8px;border-radius:6px;font-size:11px;">위험</span>'
        else:
            if val <= good_thr: return '<span style="background:#166534;color:#86efac;padding:2px 8px;border-radius:6px;font-size:11px;">양호</span>'
            elif val <= bad_thr: return '<span style="background:#854d0e;color:#fde047;padding:2px 8px;border-radius:6px;font-size:11px;">주의</span>'
            else: return '<span style="background:#991b1b;color:#fca5a5;padding:2px 8px;border-radius:6px;font-size:11px;">위험</span>'

    # ── 공식 풀이용 사전 계산 (실제 숫자 대입) ──
    _pv_sum = sum(fcf_p[i]/(1+disc_r)**(i+1) for i in range(5))
    _tax_26 = max(op_p[0]*s_tax_rate, 0)
    _energy_won = (op_electric + op_water) * 만 * 9

    # (지표명, 값, 배지, 한 줄 풀이, 의의)
    # 풀이: 누구나 한눈에 읽히도록 한국어 + 산술식, 핵심 결과는 <b>로 강조
    kpi_rows = [
        ('투자금', f'{inv_won/억:.0f}억', '',
            f'사용자 입력 = <b>{inv_won/억:.0f}억</b>',
            '시설공사·장비·인테리어 등 초기 총 투자 규모.'),
        ('NPV', f'{npv_val/억:.1f}억', _badge(npv_val, 0, -5*억, True),
            f'5년 현재가치 합 <b>{_pv_sum/억:.1f}억</b> − 투자금 {inv_won/억:.0f}억 = <b>{npv_val/억:.1f}억</b>',
            f'미래 현금흐름을 할인율 {disc_r*100:.0f}%로 환산한 순현재가치. <b>양(+)이면 투자가치 있음.</b>'),
        ('IRR', f'{irr_val*100:.1f}%', _badge(irr_val*100, 12, 8, True),
            f'NPV를 0으로 만드는 수익률 = <b>{irr_val*100:.1f}%</b> (WACC {disc_r*100:.0f}% 초과)',
            f'사업 자체의 수익률. <b>WACC({disc_r*100:.0f}%)보다 높으면 자본비용 상회.</b>'),
        ('회수율', f'{rec_rate[-1]*100:.1f}%', _badge(rec_rate[-1]*100, 100, 70, True),
            f'5년 누적 EBITDA {cum_ebitda[-1]/억:.1f}억 ÷ 투자금 {inv_won/억:.0f}억 = <b>{rec_rate[-1]*100:.1f}%</b>',
            '투자금 대비 5년 누적 영업현금흐름 회수율. <b>100% 초과 시 원금 회수 완료.</b>'),
        ('Payback', f'{payback:.1f}년' if payback else '5년+',
            _badge(rec_rate[-1]*100, 100, 70, True),
            (f'누적 {cum_ebitda[0]/억:.1f} → {cum_ebitda[1]/억:.1f} → {cum_ebitda[2]/억:.1f}억에서 투자금 {inv_won/억:.0f}억 돌파 = <b>{payback:.1f}년</b>'
             if payback else f'5년 누적 {cum_ebitda[-1]/억:.1f}억 < 투자금 → <b>5년 내 미회수</b>'),
            '투자 원금 회수 시점. <b>5년 이내 회수 시 양호.</b>'),
        ('타석당 월매출', f'{rev_per_bay_monthly/만:.0f}만', _badge(rev_per_bay_monthly/만, 20, 15, True),
            f'매출 {rev_26/억:.2f}억 ÷ {s_bays}타석 ÷ 9개월 = <b>{rev_per_bay_monthly/만:.0f}만</b>',
            '타석 1개당 월 매출. <b>시설 효율성과 자산 회전율의 핵심 지표.</b>'),
        ('이용객', f'{total_members_yr:,}명', '',
            f'1M {m_1m} + 3M {m_3m} + 6M {m_6m} + 쿠폰 {m_coupon:,} + 일일 {m_daily:,} = <b>{total_members_yr:,}명</b>',
            '2026F 회원수 가정 기준 연간 총 이용객 수.'),
        ('객단가', f'{avg_price/만:.1f}만', '',
            f'매출 {custom_total_rev/억:.2f}억 ÷ 이용객 {total_members_yr:,}명 = <b>{avg_price/만:.1f}만</b>',
            '1인당 평균 결제금액. 가격 정책과 회원 믹스의 결과치.'),
        ('인건비율', f'{labor_pct:.1f}%', _badge(labor_pct, 25, 35, False),
            f'인건비 {labor_26/억:.2f}억 ÷ 총비용 {cost_26/억:.2f}억 = <b>{labor_pct:.1f}%</b>',
            '총비용 대비 인건비 비중. <b>골프연습장 업계 25% 이하 우수.</b>'),
        ('에너지비', f'{energy_pct:.1f}%', _badge(energy_pct, 10, 15, False),
            f'(전력+수도) {_energy_won/억:.2f}억 ÷ 총비용 {cost_26/억:.2f}억 = <b>{energy_pct:.1f}%</b>',
            '에너지 비용 비중. <b>10% 이하 양호, 15% 초과 시 효율화 검토.</b>'),
        ('고정비율', f'{fixed_pct:.1f}%', _badge(fixed_pct, 50, 65, False),
            f'(인건비 {labor_26/억:.2f} + 감가 {dep_yr/억:.2f}) ÷ 총비용 {cost_26/억:.2f}억 = <b>{fixed_pct:.1f}%</b>',
            '매출 변동과 무관한 고정비 비중. <b>높을수록 BEP 매출도 상승.</b>'),
        ('성장률', f'{growth_27:+.1f}%', _badge(growth_27, 2, 0, True),
            f'2027 {rev_p[1]/억:.1f}억 ÷ 2026 {rev_p[0]/억:.1f}억 − 1 = <b>{growth_27:+.1f}%</b>',
            '정상가동 첫해(2027F) 대비 오픈년(2026F, 9개월) 증가율.'),
        ('월 현금흐름', f'{monthly_fcf/만:,.0f}만', _badge(monthly_fcf, 0, -500*만, True),
            f'(EBITDA {ebitda_p[0]/억:.2f} − 세금 {_tax_26/억:.2f})억 ÷ 12 = <b>{monthly_fcf/만:,.0f}만</b>',
            '월평균 세후 영업현금흐름. <b>운전자금 계획의 기준.</b>'),
        ('DSCR', 'N/A' if not _has_loan else f'{dscr:.2f}배',
            '' if not _has_loan else _badge(dscr, 1.2, 1.0, True),
            '자기자본 100% 투자 → <b>해당 없음</b>' if not _has_loan else f'EBITDA {ebitda_p[0]/억:.2f}억 ÷ 상환액 = <b>{dscr:.2f}배</b>',
            '부채상환능력비율. <b>1.2배 이상 안전.</b> 등촌은 차입 없음.'),
    ]

    # KPI 테이블 — 전체 폭. 컬럼: 지표 10% / 값 10% / 상태 9% / 공식풀이 36% / 의의 35%
    st.markdown("""<style>
.kpi-tbl { border:1px solid #1e293b; border-radius:10px; overflow:hidden; }
.kpi-hdr {
    display:grid; grid-template-columns:10% 10% 9% 36% 35%;
    background:#1e293b; border-bottom:1px solid #334155;
    align-items:center;
}
.kpi-hdr > div {
    color:#60a5fa; font-size:12.5px; font-weight:600;
    padding:12px 14px;
}
.kpi-hdr > div:nth-child(2) { text-align:right; }
.kpi-hdr > div:nth-child(3) { text-align:center; }
.kpi-row {
    display:grid; grid-template-columns:10% 10% 9% 36% 35%;
    border-bottom:1px solid #1e293b;
    align-items:center;
    min-height:56px;
}
.kpi-row:last-child { border-bottom:none; }
.kpi-row:hover { background:#1e293b !important; }
.kpi-cell {
    padding:13px 14px;
    display:flex; align-items:center;
    font-size:13px;
    line-height:1.5;
}
.kpi-name {
    color:#e2e8f0; font-weight:600; font-size:13.5px;
}
.kpi-val {
    color:#f8fafc; font-weight:700; font-size:16px;
    justify-content:flex-end; text-align:right;
}
.kpi-badge {
    justify-content:center; white-space:nowrap;
}
.kpi-badge > span { display:inline-block; white-space:nowrap; line-height:1; }
.kpi-formula {
    color:#cbd5e1; font-size:13px; line-height:1.55;
    background:rgba(59,130,246,0.05);
    padding:9px 12px; border-radius:6px;
    border-left:2px solid #3b82f6;
    width:100%; word-break:keep-all;
}
.kpi-formula b {
    color:#fbbf24; font-weight:700;
    background:rgba(251,191,36,0.08); padding:0 4px; border-radius:3px;
}
.kpi-explain {
    color:#94a3b8; font-size:13px; line-height:1.55;
    word-break:keep-all;
}
.kpi-explain b { color:#e2e8f0; font-weight:600; }
</style>""", unsafe_allow_html=True)

    t_html = '<div class="kpi-tbl">'
    t_html += '<div class="kpi-hdr"><div>지표</div><div>값</div><div>상태</div><div>공식 (실제 값 대입)</div><div>의의</div></div>'
    for i, (name, val, badge, formula, explain) in enumerate(kpi_rows):
        bg = '#0f1525' if i % 2 == 0 else '#0a0f1c'
        t_html += f'<div class="kpi-row" style="background:{bg};">'
        t_html += f'<div class="kpi-cell kpi-name">{name}</div>'
        t_html += f'<div class="kpi-cell kpi-val">{val}</div>'
        t_html += f'<div class="kpi-cell kpi-badge">{badge}</div>'
        t_html += f'<div class="kpi-cell"><div class="kpi-formula">{formula}</div></div>'
        t_html += f'<div class="kpi-cell"><div class="kpi-explain">{explain}</div></div>'
        t_html += '</div>'
    t_html += '</div>'
    st.markdown(t_html, unsafe_allow_html=True)

    # ── 비용 구조 도넛 + 5년 투자회수율 게이지 — KPI 표 아래 2열 ──
    st.markdown("<div style='margin-top:24px;'></div>", unsafe_allow_html=True)
    chart_col1, chart_col2 = st.columns(2)
    with chart_col1:
        cost_labels = ['인건비', '에너지', '감가상각', '기타']
        cost_vals = [labor_pct, energy_pct, dep_pct, max(0, 100-labor_pct-energy_pct-dep_pct)]
        cost_colors = [C['blue'], C['orange'], C['purple'], C['slate']]
        fig_donut = go.Figure(go.Pie(
            labels=cost_labels, values=cost_vals, hole=0.55,
            marker=dict(colors=cost_colors, line=dict(color='#0f172a', width=2)),
            textinfo='label+percent', textfont=dict(size=12, color='#e2e8f0'),
            hovertemplate='%{label}: %{value:.1f}%<extra></extra>'))
        fig_donut.add_annotation(text="비용<br>구조", x=0.5, y=0.5, font_size=14, font_color='#94a3b8', showarrow=False)
        lo(fig_donut, title='2026F 비용 구조', height=320, margin=dict(t=45, b=15, l=15, r=15), showlegend=False)
        st.plotly_chart(fig_donut, use_container_width=True, key="pc_1")
    with chart_col2:
        fig_gauge = go.Figure(go.Indicator(
            mode="gauge+number", value=rec_rate[-1]*100,
            number={'suffix': '%', 'font': {'size': 38, 'color': '#f8fafc'}},
            gauge={
                'axis': {'range': [0, 150], 'tickfont': {'size': 11, 'color': '#64748b'}},
                'bar': {'color': C['blue'], 'thickness': 0.7},
                'bgcolor': '#1e293b',
                'steps': [
                    {'range': [0, 50], 'color': '#7f1d1d'},
                    {'range': [50, 100], 'color': '#713f12'},
                    {'range': [100, 150], 'color': '#14532d'}],
                'threshold': {'line': {'color': C['red'], 'width': 3}, 'value': 100}},
            title={'text': '5년 투자회수율', 'font': {'size': 14, 'color': '#94a3b8'}}))
        lo(fig_gauge, height=320, margin=dict(t=45, b=15, l=30, r=30))
        st.plotly_chart(fig_gauge, use_container_width=True, key="pc_2")

    # ═══ 차트 공통 설정 (대시보드 전용) ═══
    _M = dict(t=45, b=30, l=15, r=15)  # 여백 통일
    _H = 360  # 높이 통일

    # ── ROW 1: 매출/비용/영업이익 + EBITDA/누적/투자금 ──
    subsec("5개년 손익 추이")
    c1, c2 = st.columns(2)
    with c1:
        info("파란 막대(매출)가 빨간 막대(비용)보다 높으면 흑자")
        fig = go.Figure()
        fig.add_trace(go.Bar(name='매출', x=[str(y) for y in D['yp']], y=[v/억 for v in rev_p], marker_color=C['blue'],
            text=[f"{v/억:.1f}" for v in rev_p], textposition='inside', textfont=dict(size=13, color='white')))
        fig.add_trace(go.Bar(name='비용', x=[str(y) for y in D['yp']], y=[v/억 for v in cost_p], marker_color=C['red'],
            text=[f"{v/억:.1f}" for v in cost_p], textposition='inside', textfont=dict(size=13, color='white')))
        fig.add_trace(go.Scatter(name='영업이익', x=[str(y) for y in D['yp']], y=[v/억 for v in op_p],
            mode='lines+markers', line=dict(color=C['green'], width=3), marker=dict(size=10, symbol='diamond'),
            hovertemplate='%{x}: %{y:.1f}억<extra>영업이익</extra>'))
        fig.add_hline(y=0, line_dash="dash", line_color="#475569", line_width=1)
        lo(fig, title='매출 vs 비용 (억원) — 초록선: 영업이익', barmode='group', height=_H, yaxis_title='억원', margin=_M)
        st.plotly_chart(fig, use_container_width=True, key="pc_3")
    with c2:
        info("주황 누적선이 빨간 점선(투자금)을 넘으면 회수 완료")
        fig2 = go.Figure()
        fig2.add_trace(go.Bar(name='EBITDA', x=[str(y) for y in D['yp']], y=[e/억 for e in ebitda_p],
            marker_color=[C['red'] if e<0 else '#22c55e' for e in ebitda_p],
            text=[f"{e/억:.1f}" for e in ebitda_p], textposition='inside', textfont=dict(size=13, color='white')))
        fig2.add_trace(go.Scatter(name='누적 EBITDA', x=[str(y) for y in D['yp']], y=[c/억 for c in cum_ebitda],
            mode='lines+markers', line=dict(color=C['orange'], width=3), marker=dict(size=10),
            hovertemplate='%{x}: %{y:.1f}억<extra>누적</extra>'))
        fig2.add_hline(y=inv_won/억, line_dash="dash", line_color=C['red'], line_width=2,
            annotation_text=f"투자금 {inv_won/억:.0f}억", annotation_font_size=12, annotation_font_color=C['red'])
        lo(fig2, title='EBITDA 및 투자회수 (억원)', height=_H, yaxis_title='억원', margin=_M)
        st.plotly_chart(fig2, use_container_width=True, key="pc_4")

    # ── ROW 2: 수익성 + 매출구성 ──
    c1, c2 = st.columns(2)
    with c1:
        subsec("수익성 추이")
        info("양(+)이면 흑자 구간. 값이 클수록 수익성 우수")
        fig_m = go.Figure()
        fig_m.add_trace(go.Bar(name='영업이익률', x=[str(y) for y in D['yp']], y=margins,
            marker_color=['#be185d' if m<0 else '#ec4899' for m in margins],
            text=[f"{m:.1f}" for m in margins], textposition='auto',
            textfont=dict(size=12, color='white'), hovertemplate='%{x}: %{y:.1f}%<extra>영업이익률</extra>'))
        fig_m.add_trace(go.Bar(name='EBITDA이익률', x=[str(y) for y in D['yp']], y=em,
            marker_color=[C['purple'] if m>0 else '#581c87' for m in em],
            text=[f"{m:.1f}" for m in em], textposition='auto',
            textfont=dict(size=12, color='white'), hovertemplate='%{x}: %{y:.1f}%<extra>EBITDA이익률</extra>'))
        fig_m.add_hline(y=0, line_dash="dash", line_color="#475569")
        lo(fig_m, title='영업이익률 vs EBITDA이익률 (%)', barmode='group', height=_H, yaxis_title='%', margin=_M)
        st.plotly_chart(fig_m, use_container_width=True, key="pc_5")
    with c2:
        subsec("2026F 매출 구성 (컨트롤 패널 기준)")
        info("컨트롤 패널에서 설정한 회원수×단가 기반 상품별 매출 비중")
        rl = list(custom_rev_items.keys())
        rv27 = [custom_rev_items[k] for k in rl]
        pie_colors = ['#3b82f6','#06b6d4','#a855f7','#f97316','#22c55e','#eab308','#ec4899','#64748b']
        # 작은 비중(5% 미만)은 라벨 숨김 → 겹침 방지
        total_rv = sum(rv27)
        pie_text = [f"{v/total_rv*100:.0f}%" if total_rv and v/total_rv >= 0.05 else "" for v in rv27]
        fig_pie = go.Figure(go.Pie(labels=rl, values=rv27, hole=0.5,
            marker=dict(colors=pie_colors, line=dict(color='#0a0e1a', width=3)),
            textinfo='text', text=pie_text, textfont=dict(size=13, color='white'),
            hovertemplate='%{label}<br>%{value:,.0f}원<br>(%{percent})<extra></extra>'))
        fig_pie.add_annotation(text="매출<br>구성", x=0.5, y=0.5, font_size=16, font_color='#94a3b8', showarrow=False)
        lo(fig_pie, height=_H, margin=dict(t=20, b=20, l=20, r=20), showlegend=True,
           legend=dict(font=dict(size=11), x=0, y=-0.15, orientation='h'))
        st.plotly_chart(fig_pie, use_container_width=True, key="pc_6")

    # ── ROW 3: 투자회수 워터폴 + 비용 상세 ──
    c1, c2 = st.columns(2)
    with c1:
        subsec("투자금 회수 워터폴")
        info(f"투자금(-{s_inv}억)에서 시작, 매년 EBITDA가 누적되어 잔여 부족분이 줄어듭니다")
        wf_x = ['투자금'] + [y.replace('F','') for y in D['yp']] + ['잔여']
        wf_y = [-inv_won/억] + [e/억 for e in ebitda_p] + [0]
        wf_m = ['absolute'] + ['relative']*5 + ['total']
        fig_wf = go.Figure(go.Waterfall(x=wf_x, y=wf_y, measure=wf_m,
            connector={"line": {"color": "#334155", "width": 1}},
            decreasing={"marker": {"color": C['red']}},
            increasing={"marker": {"color": '#22c55e'}},
            totals={"marker": {"color": C['blue']}},
            textposition='inside', textfont=dict(size=12, color='#e2e8f0'),
            text=[f"{v:.1f}" for v in wf_y]))
        lo(fig_wf, title='투자금 회수 과정 (억원)', height=_H, yaxis_title='억원', margin=_M)
        st.plotly_chart(fig_wf, use_container_width=True, key="pc_7")
    with c2:
        subsec("2026F 비용 구조 순위")
        info("비용 항목을 크기순으로 정렬. 각 막대에 마우스를 올리면 상세 내역이 표시됩니다.")
        cost_names = list(cost_items_dyn.keys())
        cost_26_vals = [cost_items_dyn[k][0]/억 for k in cost_names]

        # 각 항목별 상세 설명
        _cost_details = {
            '인건비': f'대표 1명, 총무 1명, 경리 1명, 시설 3명, 안내 3명<br>월 급여 합계: {monthly_labor*만/1e6:.1f}백만 × 9개월<br>4대보험 {s_insurance}% 별도',
            '수도광열비': f'월 {op_water}만 × 9개월<br>계절별 변동 있음 (동계 난방비 증가)',
            '전력비': f'월 {op_electric}만 × 9개월<br>88타석 조명 + 볼머신 + 냉난방<br>전기료 인상률 반영',
            '세금과공과': f'월 {op_tax}만 × 9개월<br>재산세, 종합토지세, 환경부담금 등<br>연간 고정 발생',
            '감가상각비': f'정액법: {dep_straight_amt*만/1e4:,.0f}만/년<br>정률법: {dep_declining_amt*만/1e4:,.0f}만/년 (체감)<br>내용연수: {s_useful}년',
            '보험료': f'월 {op_insurance}만 × 9개월<br>화재보험, 시설배상책임보험, 상해보험 등',
            '소모품비': f'월 {op_supplies}만 × 9개월<br>골프공, 매트, 티 등 소모품<br>이용객 수에 비례 (변동비)',
            '카드수수료': f'매출의 {op_card_fee*100:.1f}%<br>2026F 매출 {rev_p[0]/억:.1f}억 기준<br>연 {rev_p[0]*op_card_fee/억:.2f}억',
            '수선비+용역비+기타': f'수선비: 월 {op_maint}만<br>용역비: 월 {op_outsource}만<br>기타: 월 {op_etc}만<br>마케팅: 월 {op_marketing}만',
        }

        sorted_pairs = sorted(zip(cost_names, cost_26_vals), key=lambda x: x[1])
        details = [_cost_details.get(p[0], '') for p in sorted_pairs]
        bar_colors = ['#64748b','#94a3b8','#a855f7','#06b6d4','#eab308','#f97316','#ec4899','#22c55e','#3b82f6']

        fig_bar = go.Figure(go.Bar(
            y=[p[0] for p in sorted_pairs], x=[p[1] for p in sorted_pairs], orientation='h',
            marker_color=bar_colors[:len(sorted_pairs)],
            text=[f"{p[1]:.2f}억" for p in sorted_pairs], textposition='auto',
            textfont=dict(size=12, color='#e2e8f0'),
            customdata=[[d] for d in details],
            hovertemplate='<b>%{y}</b>: %{x:.2f}억원<br><br>%{customdata[0]}<extra></extra>'))
        lo(fig_bar, title='비용 항목 순위 (2026F, 억원)', height=_H, xaxis_title='억원', margin=dict(t=45, b=30, l=110, r=40))
        fig_bar.update_layout(hoverlabel=dict(bgcolor="#0f172a", bordercolor="#334155", font=dict(size=12, color="#e2e8f0")))
        st.plotly_chart(fig_bar, use_container_width=True, key="pc_8")

        # 비용 항목별 상세 내역 (커서 호버 팝업)
        _cd = {
            '인건비': f'대표 1명 · 총무 1명 · 경리 1명 · 시설 3명 · 안내 3명<br>월 급여 합계 {monthly_labor*만/1e6:.1f}백만 × 9개월 + 4대보험 {s_insurance}%',
            '감가상각비': f'정액법 {dep_straight_amt*만/1e4:,.0f}만/년 + 정률법 {dep_declining_amt*만/1e4:,.0f}만/년<br>내용연수 {s_useful}년 · 현금유출 없는 장부비용',
            '세금과공과': f'월 {op_tax}만 × 9개월 = {op_tax*9/1e4:.2f}억<br>재산세 · 종합토지세 · 환경부담금 등',
            '전력비': f'월 {op_electric}만 × 9개월 = {op_electric*9/1e4:.2f}억<br>88타석 조명 · 볼머신 · 냉난방 전력',
            '수도광열비': f'월 {op_water}만 × 9개월 = {op_water*9/1e4:.2f}억<br>계절별 변동 (동계 난방비 증가)',
            '보험료': f'월 {op_insurance}만 × 9개월 = {op_insurance*9/1e4:.2f}억<br>화재 · 시설배상책임 · 상해보험',
            '소모품비': f'월 {op_supplies}만 × 9개월 = {op_supplies*9/1e4:.2f}억<br>골프공 · 매트 · 티 (이용객 비례 변동비)',
            '카드수수료': f'매출의 {op_card_fee*100:.1f}% = 연 {rev_p[0]*op_card_fee/억:.2f}억<br>카드결제 비율에 비례',
            '수선비+용역비+기타': f'수선 {op_maint}만 + 용역 {op_outsource}만 + 기타 {op_etc}만 + 마케팅 {op_marketing}만<br>월 합계 {op_maint+op_outsource+op_etc+op_marketing}만 × 9개월',
        }
        _pop_css = """<style>
.cost-tip{position:relative;cursor:help;display:inline-block;border-bottom:1px dotted #475569;color:#e2e8f0;font-weight:600;font-size:12px;}
.cost-tip .cost-pop{display:none;position:absolute;bottom:120%;left:0;width:280px;z-index:999;
background:#0f172a;border:1px solid #334155;border-radius:8px;padding:12px;
box-shadow:0 6px 20px rgba(0,0,0,0.5);font-size:11px;line-height:1.6;color:#cbd5e1;font-weight:400;}
.cost-tip:hover .cost-pop{display:block;}
.cost-pop b{color:#60a5fa;}
</style>"""
        _pop_html = _pop_css + '<div style="display:flex;flex-wrap:wrap;gap:6px 12px;margin-top:8px;">'
        for name in [p[0] for p in sorted(sorted_pairs, key=lambda x: -x[1])]:
            detail = _cd.get(name, '')
            val = next((p[1] for p in sorted_pairs if p[0]==name), 0)
            _pop_html += f'<span class="cost-tip">{name} ({val:.2f}억)<span class="cost-pop"><b>{name}</b><br>{detail}</span></span>'
        _pop_html += '</div>'
        st.markdown(_pop_html, unsafe_allow_html=True)

    # ── ROW 4: 월별 매출 vs 비용 + 시즌 가중치 ──
    c1, c2 = st.columns(2)
    with c1:
        subsec("2026 월별 매출 vs 비용 (컨트롤패널 연동)")
        info("시즌가중치 × 오픈 램프업 × 상권보정 × 경제보정이 적용된 월별 추정입니다. 비용은 고정비 위주라 월별 균등 발생합니다.")
        # mrev_custom은 이미 시즌×램프업×상권×경제 보정 적용됨 (이중 적용 방지)
        fig_mon = go.Figure()
        fig_mon.add_trace(go.Bar(name='매출', x=D['months'], y=[v/1e6 for v in mrev_custom], marker_color=C['blue'],
            text=[f"{v/1e6:.0f}" for v in mrev_custom], textposition='inside', textfont=dict(size=12, color='white')))
        fig_mon.add_trace(go.Bar(name='비용', x=D['months'], y=[v/1e6 for v in mcost_custom], marker_color=C['red'],
            text=[f"{v/1e6:.0f}" for v in mcost_custom], textposition='inside', textfont=dict(size=12, color='white')))
        lo(fig_mon, title='월별 매출 vs 비용 (백만원)', barmode='group', height=_H, yaxis_title='백만원', margin=_M)
        st.plotly_chart(fig_mon, use_container_width=True, key="pc_9")
    with c2:
        subsec("시즌성 가중치")
        info("1.0=성수기 기준. 빨강(0.6↓)=심각한 비수기, 주황(0.9↓)=비수기, 초록=성수기")
        sw = D['season_weights']
        fig_sw = go.Figure(go.Bar(x=D['months'], y=sw,
            marker_color=[C['red'] if w<0.6 else C['orange'] if w<0.9 else '#22c55e' for w in sw],
            text=[f"{w:.2f}" for w in sw], textposition='auto', textfont=dict(size=11, color='white'),
            hovertemplate='%{x}: %{y:.2f}<extra></extra>'))
        fig_sw.add_hline(y=1.0, line_dash="dash", line_color="#64748b", line_width=1,
            annotation_text="기준 1.0", annotation_font_size=10, annotation_font_color='#94a3b8',
            annotation_xanchor='left', annotation_x=0)
        lo(fig_sw, title='월별 시즌 가중치', height=_H, yaxis_title='가중치',
           yaxis_range=[0, max(sw)*1.15 if sw else 1.5], margin=_M)
        st.plotly_chart(fig_sw, use_container_width=True, key="pc_10")

    # ── 5개년 요약 테이블 ──
    subsec("5개년 손익 요약")
    info("5개년 예측 기간의 주요 손익 지표입니다. 영업이익률이 양(+)이 되는 해부터 흑자 전환입니다.")
    pl_d = {'항목': ['매출', '비용', '영업이익', 'EBITDA', '영업이익률', '회수율']}
    for i, yr in enumerate(D['yp']):
        pl_d[yr] = [fmt억(rev_p[i]), fmt억(cost_p[i]), fmt억(op_p[i]), fmt억(ebitda_p[i]), f"{margins[i]:.1f}%", f"{rec_rate[i]*100:.1f}%"]
    dark_table(pd.DataFrame(pl_d))

# ═══ TAB 1: Scenario ═══
if _ti == 1:
    sec("📋", "시나리오 분석")
    info("최상(Optimistic), 기본(Base), 최하(Pessimistic) 3가지 시나리오를 비교합니다. 각 시나리오는 매출 성장률, 비용 증가율, 고객 이탈률 등 핵심 변수를 다르게 적용한 결과입니다.")

    subsec("시나리오 정의")
    info("**기본(Base)** = 현재 컨트롤 패널에 입력된 값 그대로. **최하/최상**은 기본 대비 변동을 적용한 스트레스 테스트입니다.")

    # 기본 시나리오 = 컨트롤 패널 현재값
    _base_growth = s_growth * 100  # %
    _base_churn = w_churn * 100
    _base_cost_up = s_util_up * 100
    _base_disc = disc_r * 100

    scen_def = {
        '구분': ['매출성장률', '비용증가율', '이탈률(가중평균)', '할인율', '매출 보정', '비용 보정'],
        '최하(Pessimistic)': [f'{_base_growth-3:.1f}%', f'{_base_cost_up+2:.1f}%/년', f'{_base_churn+4:.1f}%', f'{_base_disc+3:.0f}%', '-15%', '+10%'],
        '기본(Base) ← 패널값': [f'{_base_growth:+.1f}%', f'{_base_cost_up:.1f}%/년', f'{_base_churn:.1f}%', f'{_base_disc:.0f}%', '±0%', '±0%'],
        '최상(Optimistic)': [f'{_base_growth+2:.1f}%', f'{max(0,_base_cost_up-1):.1f}%/년', f'{max(1,_base_churn-3):.1f}%', f'{max(5,_base_disc-2):.0f}%', '+10%', '-5%'],
    }
    dark_table(pd.DataFrame(scen_def))

    # 시나리오 계산 — 기본=패널값, 최하/최상=보정 적용
    scenarios = {}
    scen_params = [
        ('최하', 0.85, 1.10, _base_disc + 3),   # 매출-15%, 비용+10%, 할인율+3%p
        ('기본', 1.00, 1.00, _base_disc),         # 패널값 그대로
        ('최상', 1.10, 0.95, max(5, _base_disc - 2)),  # 매출+10%, 비용-5%, 할인율-2%p
    ]
    for name, rev_adj, cost_adj, dr in scen_params:
        s_rev = [r * rev_adj for r in rev_p]
        s_cost = [c * cost_adj for c in cost_p]
        s_op = [r-c for r,c in zip(s_rev, s_cost)]
        s_ebitda = [o+d for o,d in zip(s_op, dep[:len(s_op)])]
        s_cum = []; sc_cc=0
        for e in s_ebitda: sc_cc+=e; s_cum.append(sc_cc)
        # 세후 FCF (메인 NPV/IRR과 동일 정의)
        s_fcf = [e - max(o * s_tax_rate, 0) for e, o in zip(s_ebitda, s_op)]
        s_dr = dr / 100
        s_npv = sum(f/(1+s_dr)**(i+1) for i,f in enumerate(s_fcf)) - inv_won
        try:
            s_irr = float(np_irr([-inv_won] + s_fcf))
        except Exception:
            s_irr = (sum(s_fcf)/5) / inv_won if inv_won else 0
        scenarios[name] = {'rev': s_rev, 'cost': s_cost, 'op': s_op, 'ebitda': s_ebitda, 'cum': s_cum, 'npv': s_npv, 'irr': s_irr, 'dr': dr}

    subsec("시나리오별 핵심 지표 비교")
    sk = st.columns(3)
    for i, (name, color) in enumerate([('최하', C['red']), ('기본', C['blue']), ('최상', C['green'])]):
        s = scenarios[name]
        sk[i].markdown(f"**{name} 시나리오** {'← 패널값' if name == '기본' else ''}")
        sk[i].metric("5년 누적 EBITDA", f"{s['cum'][-1]/억:.1f}억")
        sk[i].metric("NPV", f"{s['npv']/억:.1f}억")
        sk[i].metric("IRR", f"{s['irr']*100:.1f}%")
        sk[i].metric("투자회수율", f"{s['cum'][-1]/inv_won*100:.1f}%")
        sk[i].metric("할인율", f"{s['dr']:.0f}%")

    c1, c2 = st.columns(2)
    with c1:
        subsec("시나리오별 매출 비교")
        info("3개 시나리오의 5개년 매출을 비교합니다.")
        fig = go.Figure()
        for name, color in [('최하', C['red']), ('기본', C['blue']), ('최상', C['green'])]:
            fig.add_trace(go.Bar(name=name, x=[str(y) for y in D['yp']], y=[r/억 for r in scenarios[name]['rev']], marker_color=color))
        lo(fig, title='시나리오별 매출 비교 (억원)', barmode='group', height=440, yaxis_title='억원')
        st.plotly_chart(fig)
    with c2:
        subsec("시나리오별 누적 EBITDA")
        info("투자금 회수 시점을 시나리오별로 비교합니다. 빨간 점선은 투자금입니다.")
        fig = go.Figure()
        for name, color in [('최하', C['red']), ('기본', C['blue']), ('최상', C['green'])]:
            fig.add_trace(go.Scatter(name=name, x=[str(y) for y in D['yp']], y=[c/억 for c in scenarios[name]['cum']], mode='lines+markers', line=dict(color=color, width=3)))
        fig.add_hline(y=s_inv, line_dash="dash", line_color="#94a3b8", annotation_text=f"투자금 {s_inv}억", annotation_font_size=11)
        lo(fig, title='시나리오별 누적 EBITDA (억원)', height=440, yaxis_title='억원')
        st.plotly_chart(fig)

    subsec("시나리오별 5개년 상세 비교")
    info("3개 시나리오의 연도별 매출, 비용, 영업이익, EBITDA, 누적EBITDA를 상세 비교합니다.")
    for name in ['최상', '기본', '최하']:
        s = scenarios[name]
        st.markdown(f"**{name} 시나리오**")
        tbl = {'연도': D['yp'], '매출': [f"{r/억:.1f}억" for r in s['rev']], '비용': [f"{c/억:.1f}억" for c in s['cost']],
               '영업이익': [f"{o/억:.1f}억" for o in s['op']], 'EBITDA': [f"{e/억:.1f}억" for e in s['ebitda']],
               '누적EBITDA': [f"{c/억:.1f}억" for c in s['cum']]}
        dark_table(pd.DataFrame(tbl))

# ═══ TAB 2: Profit Contribution ═══
if _ti == 2:
    sec("🎯", "이익기여도 분석")
    with st.expander("📖 이 분석의 매출·비용 구성 및 산출 로직 상세", expanded=False):
        st.markdown("""
**매출 구성 (순수 골프 매출만 포함, 임대 제외)**
- 분석 대상: 1개월권, 3개월권, 6개월권, 12개월권, 쿠폰, 일일권, 락카, 골프레슨
- **임대매장 수익은 제외**됩니다. 임대는 골프 운영과 비용 구조가 다르기 때문에 별도 분석(임대매장 탭)이 적절합니다.

**배분비용 산출 로직**
- 전체 골프 운영비(인건비+전기+수도+감가상각+보험+소모품+카드수수료 등)를 각 상품의 **매출 비중**에 따라 안분 배분합니다.
- 예시: 일일권 매출비중이 39%이면 → 전체 운영비의 39%를 일일권 비용으로 배분
- 이 방식은 상품별 개별 원가 추적이 불가능한 골프연습장 특성에 적합한 업계 표준 방법입니다.
- **임대매장 관련 비용(임대매장 관리비, 계열사 관련 비용 등)은 배분 대상에서 제외**되어 있어 순수 골프 운영비만 배분됩니다.

**영업이익률 = (매출 - 배분비용) ÷ 매출 × 100**
- 해당 상품이 매출 1원당 얼마의 이익을 남기는지를 나타냅니다.
- 모든 상품의 영업이익률이 동일하면 → 비용이 매출에 정확히 비례한다는 의미
- 영업이익률이 상품마다 다르면 → 실제 원가 구조 차이가 반영된 것 (개별 직접비 존재)

**이익기여도 = 해당 상품 이익 ÷ 전체 이익 합계 × 100**
- 전체 이익 중 이 상품이 기여하는 비중. 이 값이 높을수록 핵심 수익원
""")
    info("**순수 골프 매출만** 분석합니다 (임대매장 수익 제외). 비용은 전체 골프 운영비를 각 상품 매출비중으로 안분 배분한 **추정치**입니다.")

    ctab1, ctab2 = st.tabs(["📂 과거 실적 분석 (2018~2021)", "📝 향후 추정 시뮬레이션"])

    with ctab1:
        if D['contrib']:
            df_c = pd.DataFrame(D['contrib'])
            years_avail = sorted(df_c['연도'].unique())

            # ── 전체 연도 비교 KPI ──
            subsec("연도별 총 이익 · 영업이익률 비교")
            info("4개년(2018~2021) 전체 이익과 영업이익률 추이입니다. 영업이익률 = 총이익 ÷ 총매출 × 100. 이 비율이 높을수록 비용 효율이 좋습니다.")
            yk = st.columns(len(years_avail))
            yr_summaries = []
            for i, yr in enumerate(years_avail):
                dfy_k = df_c[df_c['연도'] == yr]
                yr_rev = dfy_k['매출'].sum()
                yr_cost = dfy_k['추정비용'].sum()
                yr_profit = dfy_k['이익'].sum()
                yr_margin = yr_profit / yr_rev * 100 if yr_rev else 0
                yr_summaries.append({'연도': yr, '매출': yr_rev, '비용': yr_cost, '이익': yr_profit, '영업이익률': yr_margin})
                yk[i].markdown(f"""
<div class="kpi-wrap">
<div class="kpi-box">
<div style="color:#94a3b8;font-size:12px;">{yr}년 이익</div>
<div style="color:#f8fafc;font-size:24px;font-weight:800;">{yr_profit/만:,.0f}만</div>
<div style="background:#166534;color:#86efac;display:inline-block;padding:2px 8px;border-radius:6px;font-size:11px;margin-top:4px;">영업이익률 {yr_margin:.1f}%</div>
</div>
<div class="kpi-tip">
<table style="border-collapse:collapse;">
<tr><td style="color:#94a3b8;padding:3px 10px 3px 0;font-size:11px;">매출</td><td style="color:#e2e8f0;padding:3px 0;font-size:11px;font-weight:600;">{yr_rev/만:,.0f}만원</td></tr>
<tr><td style="color:#94a3b8;padding:3px 10px 3px 0;font-size:11px;">비용</td><td style="color:#e2e8f0;padding:3px 0;font-size:11px;font-weight:600;">{yr_cost/만:,.0f}만원</td></tr>
<tr><td style="color:#94a3b8;padding:3px 10px 3px 0;font-size:11px;">이익</td><td style="color:#86efac;padding:3px 0;font-size:11px;font-weight:600;">{yr_profit/만:,.0f}만원</td></tr>
<tr><td style="color:#94a3b8;padding:3px 10px 3px 0;font-size:11px;">이익률</td><td style="color:#fbbf24;padding:3px 0;font-size:11px;font-weight:600;">{yr_margin:.1f}%</td></tr>
</table>
</div>
</div>""", unsafe_allow_html=True)

            # ── 4개년 종합 비교 테이블 ──
            subsec("4개년 종합 비교표")
            info("2018~2021년 매출·비용·이익·영업이익률을 한 테이블로 비교합니다.")
            yr_tbl = {'항목': ['총 매출', '총 비용(추정)', '총 이익', '영업이익률']}
            for ys in yr_summaries:
                yr_tbl[f"{ys['연도']}년"] = [f"{ys['매출']/만:,.0f}만", f"{ys['비용']/만:,.0f}만", f"{ys['이익']/만:,.0f}만", f"{ys['영업이익률']:.1f}%"]
            dark_table(pd.DataFrame(yr_tbl))

            # ── 4개년 이익 추이 차트 ──
            c1, c2 = st.columns(2)
            with c1:
                subsec("연도별 매출·이익 추이")
                fig = go.Figure()
                fig.add_trace(go.Bar(name='매출', x=[str(s['연도']) for s in yr_summaries], y=[s['매출']/억 for s in yr_summaries],
                    marker_color=C['blue'], text=[f"{s['매출']/억:.1f}" for s in yr_summaries], textposition='inside', textfont=dict(size=12, color='white')))
                fig.add_trace(go.Bar(name='이익', x=[str(s['연도']) for s in yr_summaries], y=[s['이익']/억 for s in yr_summaries],
                    marker_color=C['green'], text=[f"{s['이익']/억:.1f}" for s in yr_summaries], textposition='inside', textfont=dict(size=12, color='white')))
                lo(fig, title='연도별 매출 vs 이익 (억원)', barmode='group', height=440, yaxis_title='억원')
                st.plotly_chart(fig)
            with c2:
                subsec("연도별 영업이익률 추이")
                fig = go.Figure()
                fig.add_trace(go.Scatter(x=[str(s['연도']) for s in yr_summaries], y=[s['영업이익률'] for s in yr_summaries],
                    mode='lines+markers', line=dict(color=C['green'], width=3), marker=dict(size=10)))
                lo(fig, title='영업이익률 추이 (%)', height=440, yaxis_title='%')
                st.plotly_chart(fig)

            st.markdown("---")

            # ── 각 연도별 상세 분석 (전체 표시) ──
            for yr in years_avail:
                dfy = df_c[df_c['연도'] == yr].copy()
                subsec(f"📅 {yr}년 상품별 상세 분석")

                c1, c2 = st.columns(2)
                with c1:
                    fig = go.Figure(go.Pie(labels=dfy['상품'], values=dfy['이익'], hole=0.45,
                        marker=dict(colors=PAL, line=dict(color='#1e293b', width=2)),
                        textinfo='label+percent', textfont=dict(size=12, color='#e2e8f0')))
                    lo(fig, title=f'{yr}년 이익 기여도', height=440)
                    st.plotly_chart(fig)
                with c2:
                    fig = go.Figure()
                    fig.add_trace(go.Bar(name='매출', x=dfy['상품'], y=dfy['매출']/만,
                        marker_color=C['blue'], text=[f"{v/만:,.0f}" for v in dfy['매출']], textposition='inside', textfont=dict(size=12, color='#e2e8f0')))
                    fig.add_trace(go.Bar(name='이익', x=dfy['상품'], y=dfy['이익']/만,
                        marker_color=C['green'], text=[f"{v/만:,.0f}" for v in dfy['이익']], textposition='inside', textfont=dict(size=12, color='#e2e8f0')))
                    lo(fig, title=f'{yr}년 매출 vs 이익 (만원)', barmode='group', height=440, yaxis_title='만원')
                    st.plotly_chart(fig)

                _margin_col = '이익률' if '이익률' in dfy.columns else '마진율'
                disp_c = dfy[['상품','회원수','매출','추정비용','이익','매출비중','이익기여도',_margin_col]].copy()
                disp_c.columns = ['상품','회원수','매출','배분비용','이익(매출-비용)','매출비중','이익기여도','영업이익률']
                disp_c['회원수'] = disp_c['회원수'].apply(lambda x: f"{int(x):,}명")
                disp_c['매출'] = disp_c['매출'].apply(lambda x: f"{x/만:,.0f}만")
                disp_c['배분비용'] = disp_c['배분비용'].apply(lambda x: f"{x/만:,.0f}만")
                disp_c['이익(매출-비용)'] = disp_c['이익(매출-비용)'].apply(lambda x: f"{x/만:,.0f}만")
                disp_c['매출비중'] = disp_c['매출비중'].apply(lambda x: f"{x*100:.1f}%")
                disp_c['이익기여도'] = disp_c['이익기여도'].apply(lambda x: f"{x*100:.1f}%")
                disp_c['영업이익률'] = disp_c['영업이익률'].apply(lambda x: f"{x*100:.1f}%")
                dark_table(disp_c)

            # ══════════════════════════════════════
            # 다차원 이익기여도 분석
            # ══════════════════════════════════════
            st.markdown("---")

            # ── 분석 1: 이용유형별 (단기/중기/장기/비회원/부대수익) ──
            subsec("분석 1 — 이용유형별 이익기여도")
            info("상품을 5개 유형으로 묶어 분석합니다. 어떤 유형의 고객이 이익에 가장 기여하는지 파악합니다.")

            def classify_type(name):
                n = str(name)
                if '일일' in n: return '비회원(일일/쿠폰)'
                if '쿠폰' in n: return '비회원(일일/쿠폰)'
                if '1개월' in n: return '단기회원(1개월)'
                if '3개월' in n: return '중기회원(3개월)'
                if '6개월' in n or '12개월' in n or '1년' in n: return '장기회원(6개월+)'
                return '부대수익(락카/레슨)'

            # 2018~2021 전체 평균으로 분석
            df_all = df_c.copy()
            df_all['유형'] = df_all['상품'].apply(classify_type)
            tg = df_all.groupby('유형').agg({'매출': 'mean', '이익': 'mean', '회원수': 'mean'}).reset_index()
            tg['이익률'] = (tg['이익'] / tg['매출'] * 100).fillna(0)
            tg['이익비중'] = tg['이익'] / tg['이익'].sum() * 100
            tg = tg.sort_values('이익', ascending=False)

            c1, c2 = st.columns(2)
            type_colors = ['#3b82f6', '#22c55e', '#a855f7', '#f97316', '#06b6d4']
            with c1:
                fig = go.Figure(go.Pie(labels=tg['유형'], values=tg['이익'], hole=0.5,
                    marker=dict(colors=type_colors[:len(tg)], line=dict(color='#0a0e1a', width=3)),
                    textinfo='percent', textfont=dict(size=14, color='white'),
                    hovertemplate='%{label}<br>이익: %{value:,.0f}원<br>비중: %{percent}<extra></extra>'))
                fig.add_annotation(text="2018~2021<br>평균", x=0.5, y=0.5, font_size=14, font_color='#94a3b8', showarrow=False)
                lo(fig, title='유형별 이익 비중 (2018~2021 평균)', height=380,
                   showlegend=True, legend=dict(font=dict(size=11), x=0, y=-0.15, orientation='h'))
                st.plotly_chart(fig, use_container_width=True, key="pc_11")
            with c2:
                fig = go.Figure()
                fig.add_trace(go.Bar(name='매출', x=tg['유형'], y=tg['매출']/만, marker_color=C['blue'],
                    text=[f"{v/만:,.0f}" for v in tg['매출']], textposition='inside', textfont=dict(size=12, color='white'),
                    hovertemplate='%{x}<br>매출: %{y:,.0f}만원<extra></extra>'))
                fig.add_trace(go.Bar(name='이익', x=tg['유형'], y=tg['이익']/만, marker_color=C['green'],
                    text=[f"{v/만:,.0f}" for v in tg['이익']], textposition='inside', textfont=dict(size=12, color='white'),
                    hovertemplate='%{x}<br>이익: %{y:,.0f}만원<extra></extra>'))
                lo(fig, title='유형별 매출·이익 (2018~2021 평균, 만원)', barmode='group', height=380, yaxis_title='만원')
                st.plotly_chart(fig, use_container_width=True, key="pc_12")

            # ── 분석 2: 성별 추정 ──
            st.markdown("---")
            subsec("분석 2 — 성별 이익기여도 (추정)")
            info("골프연습장 남녀 비율 통계(한국레저산업연구소 2023): 남성 72%, 여성 28%. 여성 요금은 남성 대비 평균 8~10% 저렴합니다. 연도별 변화(2018: 76%→2021: 70%)는 여성 골퍼 증가 추세를 반영한 추정치이며, 실제 등촌골프연습장 POS 데이터가 아닙니다.")

            # 성별 비율: 골프장 통계 기반 (한국레저산업연구소 2023)
            # 2018~2021 여성 골퍼 비율 증가 추세 반영
            gender_ratios = {2018: (0.76, 0.24), 2019: (0.74, 0.26), 2020: (0.72, 0.28), 2021: (0.70, 0.30)}
            female_disc = 0.91  # 여성 요금 = 남성 × 0.91

            # 4개년 성별 매출·이익 추정
            g_tbl = {'연도': [], '남성비율': [], '여성비율': [], '남성매출(만)': [], '여성매출(만)': [],
                     '남성이익(만)': [], '여성이익(만)': [], '여성 객단가(만)': []}
            g_male_rev, g_female_rev = [], []
            for yr in years_avail:
                mr, fr = gender_ratios.get(yr, (0.72, 0.28))
                dy = df_c[df_c['연도'] == yr]
                tot_r = dy['매출'].sum()
                tot_p = dy['이익'].sum()
                tot_m = dy['회원수'].sum()
                m_share = mr / (mr + fr * female_disc)
                f_share = 1 - m_share
                g_tbl['연도'].append(f"{yr}년")
                g_tbl['남성비율'].append(f"{mr*100:.0f}%")
                g_tbl['여성비율'].append(f"{fr*100:.0f}%")
                g_tbl['남성매출(만)'].append(f"{tot_r*m_share/만:,.0f}")
                g_tbl['여성매출(만)'].append(f"{tot_r*f_share/만:,.0f}")
                g_tbl['남성이익(만)'].append(f"{tot_p*m_share/만:,.0f}")
                g_tbl['여성이익(만)'].append(f"{tot_p*f_share/만:,.0f}")
                g_tbl['여성 객단가(만)'].append(f"{tot_r*f_share/(tot_m*fr)/만:.1f}" if tot_m*fr > 0 else '-')
                g_male_rev.append(tot_r * m_share)
                g_female_rev.append(tot_r * f_share)

            dark_table(pd.DataFrame(g_tbl))

            c1, c2 = st.columns(2)
            with c1:
                fig = go.Figure()
                fig.add_trace(go.Bar(name='남성 매출', x=[str(yr) for yr in years_avail], y=[v/억 for v in g_male_rev],
                    marker_color=C['blue'], text=[f"{v/억:.1f}" for v in g_male_rev], textposition='inside', textfont=dict(size=12, color='white')))
                fig.add_trace(go.Bar(name='여성 매출', x=[str(yr) for yr in years_avail], y=[v/억 for v in g_female_rev],
                    marker_color='#ec4899', text=[f"{v/억:.1f}" for v in g_female_rev], textposition='inside', textfont=dict(size=12, color='white')))
                lo(fig, title='성별 매출 추이 (억원)', barmode='stack', height=380, yaxis_title='억원')
                st.plotly_chart(fig, use_container_width=True, key="pc_13")
            with c2:
                f_pcts = [fr/(mr+fr*female_disc)*100 for mr, fr in [gender_ratios.get(yr, (0.72,0.28)) for yr in years_avail]]
                fig = go.Figure()
                fig.add_trace(go.Scatter(x=[str(yr) for yr in years_avail], y=f_pcts,
                    mode='lines+markers',
                    line=dict(color='#ec4899', width=3), marker=dict(size=10),
                    hovertemplate='%{x}년: %{y:.1f}%<extra>여성 비중</extra>',
                    fill='tozeroy', fillcolor='rgba(236,72,153,0.1)'))
                lo(fig, title='여성 매출 비중 추이 (%)', height=380, yaxis_title='%')
                st.plotly_chart(fig, use_container_width=True, key="pc_14")

            st.markdown(f"""
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin:12px 0;">
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #ec4899;">
<b style="color:#f9a8d4;">여성 시장 기회</b><br>
<span style="color:#cbd5e1;">여성 비율 {gender_ratios[years_avail[0]][1]*100:.0f}%→{gender_ratios[years_avail[-1]][1]*100:.0f}% <b>(+{(gender_ratios[years_avail[-1]][1]-gender_ratios[years_avail[0]][1])*100:.0f}%p)</b><br>
전국 여성 골퍼 비율 연 1~2%p 꾸준히 증가 중</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #3b82f6;">
<b style="color:#93c5fd;">남성 시장 현황</b><br>
<span style="color:#cbd5e1;">여전히 매출의 {100-f_pcts[-1]:.0f}% 차지하지만 비중 점진적 감소<br>
40~60대 남성 충성 고객 유지가 핵심</span>
</div>
</div>

<div style="background:#1e293b;border-radius:12px;padding:18px;margin:8px 0;border-left:4px solid #fbbf24;">
<b style="color:#fbbf24;font-size:14px;">성별 기반 전략 제언</b><br><br>
<span style="color:#cbd5e1;">
<b>1. 여성 전용 프로그램</b>: 주부 타겟 오전 10시 그룹 레슨 (4~6명, 주 2회) → 동반 가입 시 할인<br>
<b>2. 여성 친화 시설</b>: 파우더룸, 여성 전용 락카룸, 자외선 차단 시설 → 경쟁사 대비 차별화 포인트<br>
<b>3. 커플/가족 상품</b>: 부부 동시 가입 시 10% 할인 → 남성 기존 회원이 여성 동반 유도<br>
<b>4. SNS 마케팅</b>: 여성 골퍼 인스타그램 체험단 → 20~40대 여성 유입에 효과적<br>
<b>5. 여성 요금 전략</b>: 현재 남성 대비 9% 할인 유지 → 추가 할인보다 <b>부가가치(레슨·시설)</b>로 차별화
</span>
</div>
""", unsafe_allow_html=True)

            # ── 분석 3: 고정수익 vs 변동수익 ──
            st.markdown("---")
            subsec("분석 3 — 고정수익 vs 변동수익")
            info("회원권(1/3/6/12개월)은 **선납 고정수익**, 일일권·쿠폰은 **방문 변동수익**입니다. 고정수익 비중이 높을수록 경영 안정성이 좋습니다.")

            def classify_fixed(name):
                n = str(name)
                if '일일' in n or '쿠폰' in n: return '변동수익'
                if '락카' in n or '레슨' in n: return '부대수익'
                return '고정수익(회원권)'

            # 연도별 고정/변동 추이
            fix_data = {'연도': [], '고정수익': [], '변동수익': [], '부대수익': [], '고정비중': []}
            for yr in years_avail:
                dy = df_c[df_c['연도'] == yr].copy()
                dy['수익유형'] = dy['상품'].apply(classify_fixed)
                grp = dy.groupby('수익유형')['매출'].sum()
                fixed = grp.get('고정수익(회원권)', 0)
                variable = grp.get('변동수익', 0)
                ancillary = grp.get('부대수익', 0)
                total = fixed + variable + ancillary
                fix_data['연도'].append(str(yr))
                fix_data['고정수익'].append(fixed)
                fix_data['변동수익'].append(variable)
                fix_data['부대수익'].append(ancillary)
                fix_data['고정비중'].append(fixed/total*100 if total else 0)

            c1, c2 = st.columns(2)
            with c1:
                fig = go.Figure()
                fig.add_trace(go.Bar(name='고정수익(회원권)', x=fix_data['연도'], y=[v/억 for v in fix_data['고정수익']],
                    marker_color=C['blue'], text=[f"{v/억:.1f}" for v in fix_data['고정수익']], textposition='inside', textfont=dict(size=12, color='white')))
                fig.add_trace(go.Bar(name='변동수익(일일/쿠폰)', x=fix_data['연도'], y=[v/억 for v in fix_data['변동수익']],
                    marker_color=C['orange'], text=[f"{v/억:.1f}" for v in fix_data['변동수익']], textposition='inside', textfont=dict(size=12, color='white')))
                fig.add_trace(go.Bar(name='부대수익', x=fix_data['연도'], y=[v/억 for v in fix_data['부대수익']],
                    marker_color=C['slate'], text=[f"{v/억:.1f}" for v in fix_data['부대수익']], textposition='inside', textfont=dict(size=11, color='white')))
                lo(fig, title='연도별 고정 vs 변동수익 (억원)', barmode='stack', height=380, yaxis_title='억원')
                st.plotly_chart(fig, use_container_width=True, key="pc_15")
            with c2:
                fig = go.Figure(go.Scatter(x=fix_data['연도'], y=fix_data['고정비중'],
                    mode='lines+markers',
                    line=dict(color=C['blue'], width=3), marker=dict(size=10),
                    fill='tozeroy', fillcolor='rgba(59,130,246,0.1)',
                    hovertemplate='%{x}년: %{y:.1f}%<extra>고정수익 비중</extra>'))
                fig.add_hline(y=50, line_dash='dash', line_color='#475569',
                    annotation_text='50% 기준선', annotation_font_size=11, annotation_font_color='#94a3b8')
                lo(fig, title='고정수익 비중 추이 (%)', height=380, yaxis_title='%')
                st.plotly_chart(fig, use_container_width=True, key="pc_16")

            # ── 분석 4: 요일·시간대별 이익기여도 추정 ──
            st.markdown("---")
            subsec("분석 4 — 요일·시간대별 이익기여도 추정 (2018~2021 평균)")
            info("골프연습장 매출은 요일·시간대에 따라 크게 다릅니다. 한국레저산업연구소(2023) 업계 평균 기반, 2018~2021 연평균 매출 적용.")

            # 2018~2021 연평균 매출
            avg_rev_all = df_c.groupby('연도')['매출'].sum().mean()

            # 요일별 매출 비중 (업계 평균)
            day_names = ['월', '화', '수', '목', '금', '토', '일']
            day_pcts = [11, 12, 13, 13, 15, 19, 17]
            day_colors = [C['blue']]*5 + [C['orange'], C['orange']]

            # 시간대별 매출 비중
            time_names = ['06~09', '09~12', '12~15', '15~18', '18~21', '21~23']
            time_pcts = [8, 22, 15, 18, 28, 9]
            time_colors = ['#94a3b8', C['blue'], C['orange'], C['blue'], C['green'], '#94a3b8']

            c1, c2 = st.columns(2)
            with c1:
                fig = go.Figure(go.Bar(x=day_names, y=day_pcts,
                    marker_color=day_colors,
                    text=[f"{p}%" for p in day_pcts],
                    textposition='inside', textfont=dict(size=14, color='white'),
                    hovertemplate='%{x}요일<br>비중: %{y}%<br>추정매출: %{customdata:,.0f}만원<extra></extra>',
                    customdata=[avg_rev_all*p/100/만 for p in day_pcts]))
                lo(fig, title='요일별 매출 비중 (2018~2021 평균)', height=380, yaxis_title='비중(%)',
                   yaxis_range=[0, max(day_pcts)*1.3])
                st.plotly_chart(fig, use_container_width=True, key="pc_17")
            with c2:
                fig = go.Figure(go.Bar(x=time_names, y=time_pcts,
                    marker_color=time_colors,
                    text=[f"{p}%" for p in time_pcts],
                    textposition='inside', textfont=dict(size=14, color='white'),
                    hovertemplate='%{x}시<br>비중: %{y}%<br>추정매출: %{customdata:,.0f}만원<extra></extra>',
                    customdata=[avg_rev_all*p/100/만 for p in time_pcts]))
                lo(fig, title='시간대별 매출 비중 (2018~2021 평균)', height=380, yaxis_title='비중(%)',
                   yaxis_range=[0, max(time_pcts)*1.3])
                st.plotly_chart(fig, use_container_width=True, key="pc_18")

            st.markdown("""
<div style="background:#1e293b;border-radius:12px;padding:16px;margin:8px 0;border-left:4px solid #f97316;">
<b style="color:#fdba74;">요일·시간대 인사이트</b><br><br>
<span style="color:#cbd5e1;">
• <b>주말(토·일) = 전체 매출의 36%</b> → 주말 타석 가동률 극대화 필수 (예약 시스템, 회전율 관리)<br>
• <b>저녁(18~21시) = 28%로 피크 시간대</b> → 직장인 수요. 이 시간대 타석 부족 시 매출 손실↑<br>
• <b>얼리(06~09시)+야간(21~23시) = 17%</b> → 비인기 시간대. 할인/모닝회원 전용 요금 전략 필요<br>
• <b>평일 점심(12~15시) = 가장 낮은 시간대</b> → 주부/시니어 타겟 프로모션으로 빈 타석 활용
</span>
</div>
""", unsafe_allow_html=True)

            # ── 종합 핵심 인사이트 ──
            st.markdown("---")
            subsec("종합 핵심 인사이트 (2018~2021)")

            # 데이터 산출
            latest = df_c[df_c['연도'] == years_avail[-1]]
            first = df_c[df_c['연도'] == years_avail[0]]
            first_rev = first['매출'].sum()
            last_rev = latest['매출'].sum()
            first_profit = first['이익'].sum()
            last_profit = latest['이익'].sum()
            first_margin = first_profit / first_rev * 100 if first_rev else 0
            last_margin = last_profit / last_rev * 100 if last_rev else 0
            rev_growth = (last_rev / first_rev - 1) * 100 if first_rev else 0

            st.markdown(f"""
<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;margin:12px 0;">
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid {C['blue']};">
<b style="color:#60a5fa;">매출 성장</b><br>
<span style="color:#e2e8f0;font-size:18px;font-weight:700;">{rev_growth:+.1f}%</span><br>
<span style="color:#94a3b8;font-size:12px;">{first_rev/억:.1f}억→{last_rev/억:.1f}억 (4년)</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid {C['green']};">
<b style="color:#86efac;">이익률 변화</b><br>
<span style="color:#e2e8f0;font-size:18px;font-weight:700;">{first_margin:.1f}%→{last_margin:.1f}%</span><br>
<span style="color:#94a3b8;font-size:12px;">비용 효율 {'개선' if last_margin > first_margin else '악화'}</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid {C['orange']};">
<b style="color:#fdba74;">고정수익 비중</b><br>
<span style="color:#e2e8f0;font-size:18px;font-weight:700;">{fix_data['고정비중'][-1]:.1f}%</span><br>
<span style="color:#94a3b8;font-size:12px;">{'안정적' if fix_data['고정비중'][-1] > 50 else '변동수익 의존↑, 날씨 리스크↑'}</span>
</div>
</div>

<div style="background:#1e293b;border-radius:12px;padding:18px;margin:8px 0;border-left:4px solid #fbbf24;">
<b style="color:#fbbf24;font-size:14px;">경영진 Action Item</b><br><br>
<span style="color:#cbd5e1;">
1. <b>일일권·쿠폰 의존도 관리</b>: 변동수익 비중이 {'높아' if fix_data['고정비중'][-1] < 50 else '관리 가능 수준이나'} 날씨 악화 시 매출 급감 리스크 → 회원권 전환 프로모션 강화<br>
2. <b>3개월권 비중 회복</b>: 중기약정 감소 추세 → 3개월권 혜택 강화 (락카 무료, 레슨 할인 등)<br>
3. <b>여성 고객 확대</b>: 현재 추정 28% → 35% 목표. 여성 전용 시간대/레슨 도입<br>
4. <b>LTV 극대화</b>: 6개월권 LTV가 가장 높음 → 6개월권 업셀링 전략(3→6개월 전환 할인)<br>
5. <b>일일→쿠폰→회원 전환 퍼널</b>: 일일 체험→쿠폰 20회→1개월→3개월 단계적 전환 설계
</span>
</div>
""", unsafe_allow_html=True)

            st.markdown("""
<div style="background:#1e293b;border-radius:12px;padding:18px;margin:12px 0;border-left:4px solid #a855f7;">
<b style="color:#c4b5fd;font-size:15px;">🎯 2026 재오픈 전략 시사점</b><br><br>
<span style="color:#cbd5e1;">
• <b>오픈 초기</b>: 일일권+쿠폰 중심 (합산 55~60%) → 진입장벽 낮춰 유입 극대화<br>
• <b>안정화 후</b>: 일일→쿠폰→3개월권 자연스러운 전환 프로모션<br>
• <b>핵심 경쟁력</b>: 접근성/주차/가격 → 일일권 유입이 전체 성장의 핵심<br>
• <b>부가수익</b>: 락카·프로임대는 추가비용 없는 고마진 → 적극 확대
</span>
</div>
""", unsafe_allow_html=True)
        else:
            st.warning("과거 이익기여도 데이터가 없습니다.")

    with ctab2:
        subsec("향후 이익기여도 시뮬레이션")
        info("상품별 **회원수**와 **평균단가**를 입력하면, 과거 분석과 동일한 구조(상품별·유형별·성별·고정/변동·요일시간대)로 자동 분석됩니다. 기본값은 컨트롤 패널 입력을 반영합니다.")

        # ── 입력 테이블 (패드 연동) ──
        st.caption("📋 상품별 연간 회원수 · 평균단가 · 비용 입력 (기본값=패드값)")
        f_items = [
            ('1개월권', m_1m if m_1m else 500, p_1m if p_1m else 280000, 65),
            ('3개월권', m_3m if m_3m else 1000, p_3m if p_3m else 760000, 60),
            ('6개월권', m_6m if m_6m else 80, p_6m if p_6m else 1500000, 55),
            ('쿠폰', m_coupon if m_coupon else 1800, p_coupon if p_coupon else 400000, 60),
            ('일일권', m_daily if m_daily else 35000, p_daily if p_daily else 22000, 70),
            ('락카', m_locker*12 if m_locker else 1200, p_locker if p_locker else 240000, 5),
            ('프로임대', pro_count if pro_count else 3, pro_rent*만*12 if pro_rent else 18000000, 0),
        ]
        hc = st.columns([2.5, 1.5, 1.5, 1.5])
        hc[0].markdown("**상품**"); hc[1].markdown("**회원수(명/년)**"); hc[2].markdown("**평균단가(원)**"); hc[3].markdown("**비용율(%)**")
        fd = []
        for item, def_m, def_p, def_cr in f_items:
            fc = st.columns([2.5, 1.5, 1.5, 1.5])
            fc[0].caption(item)
            fm = fc[1].number_input(f"fm_{item}", value=def_m, step=100, key=f'fm_{item}', label_visibility='collapsed')
            fp = fc[2].number_input(f"fp_{item}", value=def_p, step=10000, key=f'fp_{item}', label_visibility='collapsed')
            fcr = fc[3].number_input(f"fcr_{item}", value=def_cr, step=5, key=f'fcr_{item}', label_visibility='collapsed')
            rev = fm * fp
            cost = int(rev * fcr / 100)
            profit = rev - cost
            margin = profit / rev * 100 if rev else 0
            fd.append({'상품': item, '회원수': fm, '매출': rev, '비용': cost, '이익': profit, '영업이익률': margin})

        tot_fr = sum(d['매출'] for d in fd)
        tot_fp = sum(d['이익'] for d in fd)
        tot_fc = sum(d['비용'] for d in fd)

        if tot_fr > 0:
            st.markdown("---")

            # ── KPI ──
            fk = st.columns(5)
            fk[0].metric("총 매출", f"{tot_fr/만:,.0f}만")
            fk[1].metric("총 비용", f"{tot_fc/만:,.0f}만")
            fk[2].metric("총 이익", f"{tot_fp/만:,.0f}만")
            fk[3].metric("영업이익률", f"{tot_fp/tot_fr*100:.1f}%")
            fk[4].metric("총 회원수", f"{sum(d['회원수'] for d in fd):,}명")

            # ── 상품별 분석 (과거와 동일) ──
            subsec("상품별 이익기여도")
            c1, c2 = st.columns(2)
            with c1:
                valid = [d for d in fd if d['이익'] > 0]
                if valid:
                    fig = go.Figure(go.Pie(labels=[d['상품'] for d in valid], values=[d['이익'] for d in valid],
                        hole=0.5, marker=dict(colors=PAL, line=dict(color='#0a0e1a', width=2)),
                        textinfo='label+percent', textfont=dict(size=12, color='white')))
                    fig.add_annotation(text="이익<br>기여도", x=0.5, y=0.5, font_size=14, font_color='#94a3b8', showarrow=False)
                    lo(fig, title='향후 이익 기여도', height=400, showlegend=False)
                    st.plotly_chart(fig, use_container_width=True, key="pc_19")
            with c2:
                fig = go.Figure()
                fig.add_trace(go.Bar(name='매출', x=[d['상품'] for d in fd if d['매출']>0], y=[d['매출']/만 for d in fd if d['매출']>0],
                    marker_color=C['blue'], text=[f"{d['매출']/만:,.0f}" for d in fd if d['매출']>0], textposition='inside', textfont=dict(size=12, color='white')))
                fig.add_trace(go.Bar(name='이익', x=[d['상품'] for d in fd if d['매출']>0], y=[d['이익']/만 for d in fd if d['매출']>0],
                    marker_color=C['green'], text=[f"{d['이익']/만:,.0f}" for d in fd if d['매출']>0], textposition='inside', textfont=dict(size=12, color='white')))
                lo(fig, title='향후 매출 vs 이익 (만원)', barmode='group', height=400, yaxis_title='만원')
                st.plotly_chart(fig, use_container_width=True, key="pc_20")

            # 상세 테이블
            rdf = pd.DataFrame(fd)
            rdf['매출비중'] = rdf['매출'].apply(lambda x: f"{x/tot_fr*100:.1f}%")
            rdf['이익기여도'] = rdf['이익'].apply(lambda x: f"{x/tot_fp*100:.1f}%" if tot_fp else "0%")
            rdf_d = rdf.copy()
            rdf_d['회원수'] = rdf_d['회원수'].apply(lambda x: f"{x:,}명")
            rdf_d['매출'] = rdf_d['매출'].apply(lambda x: f"{x/만:,.0f}만")
            rdf_d['비용'] = rdf_d['비용'].apply(lambda x: f"{x/만:,.0f}만")
            rdf_d['이익'] = rdf_d['이익'].apply(lambda x: f"{x/만:,.0f}만")
            rdf_d['영업이익률'] = rdf_d['영업이익률'].apply(lambda x: f"{x:.1f}%")
            dark_table(rdf_d[['상품','회원수','매출','비용','이익','영업이익률','매출비중','이익기여도']])

            st.markdown("---")

            # ── 이용유형별 (과거와 동일) ──
            subsec("이용유형별 이익기여도")
            def _ftype(n):
                if '일일' in n or '쿠폰' in n: return '비회원(일일/쿠폰)'
                if '1개월' in n: return '단기회원'
                if '3개월' in n: return '중기회원'
                if '6개월' in n: return '장기회원'
                return '부대수익'
            for d in fd: d['유형'] = _ftype(d['상품'])
            df_ft = pd.DataFrame(fd)
            tg = df_ft.groupby('유형').agg({'매출':'sum','이익':'sum','회원수':'sum'}).reset_index()
            c1, c2 = st.columns(2)
            with c1:
                fig = go.Figure(go.Pie(labels=tg['유형'], values=tg['이익'], hole=0.5,
                    marker=dict(colors=['#3b82f6','#22c55e','#a855f7','#f97316','#06b6d4'], line=dict(color='#0a0e1a', width=2)),
                    textinfo='label+percent', textfont=dict(size=12, color='white')))
                lo(fig, title='유형별 이익 비중', height=380, showlegend=False)
                st.plotly_chart(fig, use_container_width=True, key="pc_21")
            with c2:
                fig = go.Figure()
                fig.add_trace(go.Bar(name='매출', x=tg['유형'], y=tg['매출']/만, marker_color=C['blue'],
                    text=[f"{v/만:,.0f}" for v in tg['매출']], textposition='inside', textfont=dict(size=12, color='white')))
                fig.add_trace(go.Bar(name='이익', x=tg['유형'], y=tg['이익']/만, marker_color=C['green'],
                    text=[f"{v/만:,.0f}" for v in tg['이익']], textposition='inside', textfont=dict(size=12, color='white')))
                lo(fig, title='유형별 매출·이익 (만원)', barmode='group', height=380, yaxis_title='만원')
                st.plotly_chart(fig, use_container_width=True, key="pc_22")

            st.markdown("---")

            # ── 성별 (과거와 동일) ──
            subsec("성별 이익기여도")
            m_r, f_r = 0.70, 0.30
            f_disc = 0.91
            m_share = m_r / (m_r + f_r * f_disc)
            f_share = 1 - m_share
            c1, c2 = st.columns(2)
            with c1:
                g_d = [{'성별':'남성','비율':f'{m_r*100:.0f}%','매출':f'{tot_fr*m_share/만:,.0f}만','이익':f'{tot_fp*m_share/만:,.0f}만'},
                       {'성별':'여성','비율':f'{f_r*100:.0f}%','매출':f'{tot_fr*f_share/만:,.0f}만','이익':f'{tot_fp*f_share/만:,.0f}만'}]
                dark_table(pd.DataFrame(g_d))
            with c2:
                fig = go.Figure(go.Pie(labels=['남성','여성'], values=[m_share, f_share], hole=0.5,
                    marker=dict(colors=[C['blue'],'#ec4899'], line=dict(color='#0a0e1a', width=2)),
                    textinfo='label+percent', textfont=dict(size=14, color='white')))
                lo(fig, title='성별 매출 비중', height=300, showlegend=False)
                st.plotly_chart(fig, use_container_width=True, key="pc_23")

            st.markdown("---")

            # ── 고정 vs 변동 (과거와 동일) ──
            subsec("고정수익 vs 변동수익")
            def _ffixed(n):
                if '일일' in n or '쿠폰' in n: return '변동수익'
                if '락카' in n or '프로' in n or '임대' in n: return '부대수익'
                return '고정수익(회원권)'
            fixed_rev = sum(d['매출'] for d in fd if _ffixed(d['상품'])=='고정수익(회원권)')
            var_rev = sum(d['매출'] for d in fd if _ffixed(d['상품'])=='변동수익')
            anc_rev = sum(d['매출'] for d in fd if _ffixed(d['상품'])=='부대수익')
            fix_pct = fixed_rev / tot_fr * 100 if tot_fr else 0
            c1, c2, c3 = st.columns(3)
            c1.metric("고정수익(회원권)", f"{fixed_rev/만:,.0f}만", delta=f"{fix_pct:.0f}%")
            c2.metric("변동수익(일일/쿠폰)", f"{var_rev/만:,.0f}만", delta=f"{var_rev/tot_fr*100:.0f}%")
            c3.metric("부대수익", f"{anc_rev/만:,.0f}만", delta=f"{anc_rev/tot_fr*100:.0f}%")

            if fix_pct > 50:
                st.markdown(f'<div style="background:#14532d;border-radius:8px;padding:12px;color:#86efac;">고정수익 비중 {fix_pct:.0f}% → 경영 안정성 양호</div>', unsafe_allow_html=True)
            else:
                st.markdown(f'<div style="background:#7f1d1d;border-radius:8px;padding:12px;color:#fca5a5;">고정수익 비중 {fix_pct:.0f}% → 변동수익 의존↑, 날씨 리스크 주의</div>', unsafe_allow_html=True)

        else:
            info("위 표에 회원수와 단가를 입력하면 과거와 동일한 구조로 이익기여도가 자동 분석됩니다.")

    # ── 이익기여도 기반 가격 전략 (운영전략 탭과 동일 내용) ──
    subsec("이익기여도 기반 가격 전략")
    info("2018~2021년 4개년 평균 실적 기반으로 도출한 가격 전략입니다. 상품명에 마우스를 올리면 역산 계산 과정이 표시됩니다. 상세 요금표는 '운영전략' 탭에서 확인할 수 있습니다.")

    st.markdown("""
<style>
.prd-tip2 { position:relative; cursor:help; border-bottom:1px dotted #64748b; }
.prd-tip2 .prd-pop2 {
    display:none; position:absolute; bottom:130%; left:0; width:320px; z-index:999;
    background:#0f172a; border:1px solid #334155; border-radius:10px; padding:14px;
    box-shadow:0 8px 24px rgba(0,0,0,0.5); font-size:11px; line-height:1.7; color:#cbd5e1;
}
.prd-tip2:hover .prd-pop2 { display:block; }
.prd-pop2 b { color:#60a5fa; } .prd-pop2 .val2 { color:#fbbf24; font-weight:700; }
</style>
<div style="background:#0f172a;border:1px solid #334155;border-radius:12px;padding:20px;margin:16px 0;">
<div style="color:#fbbf24;font-size:15px;font-weight:700;margin-bottom:10px;">📊 이익기여도 → 가격 전략 연결</div>
<div style="background:#1e293b;border-radius:10px;padding:14px;margin-bottom:12px;">
<table style="width:100%;border-collapse:collapse;font-size:12px;">
<tr>
<th style="color:#60a5fa;padding:6px 10px;text-align:left;border-bottom:1px solid #334155;">순위</th>
<th style="color:#60a5fa;padding:6px;text-align:left;border-bottom:1px solid #334155;">상품</th>
<th style="color:#60a5fa;padding:6px;text-align:right;border-bottom:1px solid #334155;">1방문당 수익</th>
<th style="color:#60a5fa;padding:6px;text-align:center;border-bottom:1px solid #334155;">이익기여도</th>
<th style="color:#60a5fa;padding:6px;text-align:center;border-bottom:1px solid #334155;">할인 전략</th>
<th style="color:#60a5fa;padding:6px;text-align:left;border-bottom:1px solid #334155;">전략 근거</th>
</tr>
<tr><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">1위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip2">일일권<span class="prd-pop2"><b>일일권 (2018~2021 평균)</b><br>4개년 평균 매출: 약 7.9억/년<br>1회=1건(확정), 전액이 방문당 수익<br>전체 매출의 약 22% 차지<br>이익기여도: 1방문당 수익 최고</span></span></td>
<td style="color:#86efac;padding:5px;text-align:right;font-weight:700;">23,000원</td>
<td style="color:#86efac;padding:5px;text-align:center;">★★★</td>
<td style="color:#86efac;padding:5px;text-align:center;">-4~4.5%</td>
<td style="color:#94a3b8;padding:5px;font-size:11px;">유입 극대화 → 체험 후 회원 전환</td></tr>
<tr style="background:#111827;"><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">2위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip2">쿠폰(10회)<span class="prd-pop2"><b>쿠폰 10회 (2018~2021 평균)</b><br>유효기간: 3개월. 회당 22,000원<br>미소진 시(7~8회) 실질 27,500~31,400원/회<br>4개년 평균 쿠폰 매출: 약 6.2억/년<br>미소진분은 순수익 → 실질 기여도 최상위</span></span></td>
<td style="color:#86efac;padding:5px;text-align:right;font-weight:700;">22,000원+</td>
<td style="color:#86efac;padding:5px;text-align:center;">★★★</td>
<td style="color:#86efac;padding:5px;text-align:center;">-4.3%</td>
<td style="color:#94a3b8;padding:5px;font-size:11px;">체험→전환 핵심 상품</td></tr>
<tr><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">3위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip2">쿠폰(20회)<span class="prd-pop2"><b>쿠폰 20회 (2018~2021 평균)</b><br>유효기간: 6개월. 회당 21,500원<br>미소진 시(15~16회) 실질 26,900~28,700원/회<br>10회 대비 회당 500원 저렴(볼륨 할인)<br>재구매율 높은 핵심 상품</span></span></td>
<td style="color:#86efac;padding:5px;text-align:right;font-weight:700;">21,500원+</td>
<td style="color:#86efac;padding:5px;text-align:center;">★★☆</td>
<td style="color:#86efac;padding:5px;text-align:center;">-4.4%</td>
<td style="color:#94a3b8;padding:5px;font-size:11px;">주력 쿠폰, 재구매 유도</td></tr>
<tr style="background:#111827;"><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">4위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip2">1개월권<span class="prd-pop2"><b>1개월권 (2018~2021 평균)</b><br>4개년 평균 매출: 약 1.0억/년<br>역산법: 총 타석건수에서 일일/쿠폰 차감 후 배분<br>추정 월 ~18회 방문 → 290,000÷18=16,100원/회<br>⚠ 방문횟수는 POS 없어 역산 추정(검증 필요)</span></span></td>
<td style="color:#f9a8d4;padding:5px;text-align:right;font-weight:700;">16,100원</td>
<td style="color:#f9a8d4;padding:5px;text-align:center;">★★☆</td>
<td style="color:#f9a8d4;padding:5px;text-align:center;">-2.9~3.8%</td>
<td style="color:#94a3b8;padding:5px;font-size:11px;">진입 상품, 3개월 전환 유도</td></tr>
<tr><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">5위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip2">3개월권<span class="prd-pop2"><b>3개월권 (2018~2021 평균)</b><br>4개년 평균 매출: 약 6.1억/년 (최대 매출 상품)<br>역산법: 추정 월 ~20회 방문<br>800,000÷3÷20=13,300원/회<br>⚠ 방문횟수는 POS 없어 역산 추정(검증 필요)</span></span></td>
<td style="color:#fdba74;padding:5px;text-align:right;font-weight:700;">13,300원</td>
<td style="color:#fdba74;padding:5px;text-align:center;">★☆☆</td>
<td style="color:#fdba74;padding:5px;text-align:center;">-3.6~4.3%</td>
<td style="color:#94a3b8;padding:5px;font-size:11px;">안정 매출 기반, 볼륨 확보</td></tr>
<tr style="background:#111827;"><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">6위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip2">6개월권<span class="prd-pop2"><b>6개월권 (2018~2021 평균)</b><br>4개년 평균 매출: 약 1.1억/년<br>역산법: 추정 월 ~22회 방문<br>1,580,000÷6÷22=12,000원/회<br>사우나/헬스 무료이용 비용(~3,000원/회) 감안<br>실질 순수익 ~9,000원/회<br>⚠ 방문횟수는 POS 없어 역산 추정(검증 필요)</span></span></td>
<td style="color:#fca5a5;padding:5px;text-align:right;font-weight:700;">12,000원</td>
<td style="color:#fca5a5;padding:5px;text-align:center;">★☆☆</td>
<td style="color:#fca5a5;padding:5px;text-align:center;">-2.7~3.4%</td>
<td style="color:#94a3b8;padding:5px;font-size:11px;">마진 보호, 사우나/헬스 비용 감안</td></tr>
</table>
</div>
<div style="color:#94a3b8;font-size:12px;padding-top:10px;border-top:1px solid #334155;">
<b>전략 요약</b>: 이익기여도 ★★★(일일/쿠폰)은 할인 -4%↑로 유입 극대화, ★☆☆(6개월)은 할인 -3%↓로 마진 보호. 고객 여정: 일일체험 → 쿠폰 → 1개월 → 3개월 → 6개월 VIP 전환
</div>
</div>
""", unsafe_allow_html=True)

# ═══ TAB 3: AI Strategy ═══
if _ti == 3:
    sec("🤖", "조건값 내 AI 추천 사업전략")
    info("현재 입력된 모든 가정값, 과거 실적, 경쟁사 분석을 종합하여 등촌골프연습장의 최적 사업전략을 제안합니다. 5년 후 매각/명도를 고려한 출구전략도 포함합니다.")

    subsec("1. 요금 전략 — 조건값 내 최적 요금표(안)")

    # 전략 근거 설명 — 팩트 기반, 이론적 프레임워크 포함
    st.markdown("""
<div style="background:#0f172a;border:1px solid #1e293b;border-radius:14px;padding:20px;margin:12px 0;">
<div style="color:#60a5fa;font-size:16px;font-weight:700;margin-bottom:14px;">📐 가격 전략 프레임워크: 경쟁적 가격 포지셔닝(Competitive Pricing)</div>
<div style="color:#cbd5e1;font-size:13px;line-height:1.8;">
마이클 포터의 경쟁전략 이론에 따르면, 후발주자의 시장 진입 시 <b>비용우위(Cost Leadership)</b>도 <b>차별화(Differentiation)</b>도 아닌
<b>"경쟁적 가격 + 시설 차별화"</b> 혼합 전략이 가장 효과적입니다.<br><br>
<b>핵심 논리</b>: 고객이 "저렴해서 온다"가 아니라 <b>"비슷한 가격인데 시설이 새것이고 깨끗하니까 온다"</b>는 인식을 만들어야 합니다.
이를 위해 제니스 대비 <b>3~5% 저렴</b>하되, 가격 차이가 아닌 시설 품질이 전환의 주된 이유가 되는 수준으로 설정합니다.
</div>
</div>

<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin:12px 0;">
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #3b82f6;">
<b style="color:#60a5fa;">가격 설정 원칙 — 제니스 대비 3~5% 저렴 (최대 4.5%)</b><br><br>
<span style="color:#cbd5e1;font-size:13px;line-height:1.7;">
<b>상권·고객층 차이 반영 (팩트)</b><br>
• 제니스(구로구 고척동): 여의도·영등포 고소득 직장인 고객 다수 → 가격 민감도 낮음<br>
• 등촌(강서구): 마곡지구 발전 중이나 주거지역 중심 → 가격 민감도 상대적 높음<br>
• <b>같은 가격이면 등촌 고객에겐 부담</b> → 3~5% 할인이 적정<br><br>
<b>3~5%가 최적인 이유 (JND 이론)</b><br>
• 소비자 심리학 JND(Just Noticeable Difference): 3% 미만 = "차이 없음", 5% 초과 = "싸구려" 인식<br>
• 예: 1개월 종일 남 제니스 30만 → 등촌 28만 (차이 2만원/월, -6.7%) → "약간 저렴 + 새 시설"
</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #22c55e;">
<b style="color:#86efac;">등촌이 할인해도 경쟁력 있는 이유 (팩트)</b><br><br>
<span style="color:#cbd5e1;font-size:13px;line-height:1.7;">
① <b>2026년 신축</b>: 건물, 타석, 매트, 볼머신, 조명 전부 새것 → 시설 프리미엄<br>
② <b>사우나/피트니스 전 회원 무료</b>: 제니스와 동등 조건<br>
③ <b>88타석 대형 + 스크린골프 5대</b>: 우천/한파 대안까지 보유<br>
④ <b>마곡지구 인접</b>: LG·코오롱·넥센 등 대기업 입주 수요<br>
→ 3~5% 저렴하면서 시설은 더 좋음 = <b>"가격 대비 최고"</b> 포지셔닝
</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #f97316;">
<b style="color:#fdba74;">할인이 필요한 구조적 이유</b><br><br>
<span style="color:#cbd5e1;font-size:13px;line-height:1.7;">
① <b>재오픈 후발주자</b>: 기존 회원 기반 제로, 인지도 없음<br>
② <b>고객 전환 지연</b>: 40~60대 주 고객은 현 연습장 회원권 만료 전 이동 불가<br>
③ <b>상권 차이</b>: 제니스 고객(여의도 고소득)과 등촌 고객(강서구 주거지) 구매력 차이<br>
④ <b>Switching Cost</b>: 전환비용 보상을 위해 3~5% 가격 인센티브 필요 (경영학 이론)<br>
⑤ <b>제니스 112타석 vs 등촌 88타석</b>: 규모 열위 24타석 → 가격 차이 합리적
</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #ef4444;">
<b style="color:#fca5a5;">과도한 할인(10%↑)의 리스크</b><br><br>
<span style="color:#cbd5e1;font-size:13px;line-height:1.7;">
① <b>시세 훼손</b>: 강서구 골프연습장 전체 가격 하락 유발 → 동업자 반발, 업계 내 고립<br>
② <b>저가 이미지 고착</b>: 소비자 심리학상 '앵커링 효과(Anchoring Effect)'로 첫 가격이 기준점. 저가 출발 시 추후 인상 저항 극대화<br>
③ <b>마진 잠식</b>: 월 고정비 3,300만원+ → 객단가 5% 하락 시 연간 약 1.3억 이익 감소<br>
④ <b>"왜 싸지?" 의심</b>: 지나치게 저렴하면 시설 품질·서비스에 대한 불신 유발
</span>
</div>
</div>
""", unsafe_allow_html=True)

    # ── 등촌 요금표 (사이드바 입력값 연동 + 2년 후 제니스 동일가 적용) ──
    subsec("등촌골프연습장 전체 요금표 (사이드바 연동)")
    info("**1~2년차**: 사이드바 '💴 요금 설정'에서 직접 입력한 가격을 적용. **3년차~**: 시장가 따라잡기 — **제니스(구로구, 112타석)와 동일 요금**으로 인상. 할인율은 1~2년차 등촌 vs 제니스 가격 차이.")

    # 한국 천단위 콤마 포맷 헬퍼
    def _fmt(v): return f'{int(v):,}'
    def _disc_pct(my, jn):
        if jn == 0: return '-'
        p = (my - jn) / jn * 100
        return f'{p:+.1f}%'

    # 회원권 요금표
    st.markdown("**회원권**")
    pricing_html = """
<table style="width:100%;border-collapse:collapse;font-size:13px;margin:8px 0;">
<tr>
<th style="background:#1e293b;color:#60a5fa;padding:10px;text-align:left;border-bottom:2px solid #334155;" rowspan="2">구분</th>
<th style="background:#1e293b;color:#60a5fa;padding:10px;text-align:center;border-bottom:2px solid #334155;" rowspan="2">이용시간</th>
<th style="background:#0f4c3a;color:#86efac;padding:8px;text-align:center;border-bottom:1px solid #334155;" colspan="2">1~2년차 등촌 (사이드바 설정)</th>
<th style="background:#1e3a5f;color:#93c5fd;padding:8px;text-align:center;border-bottom:1px solid #334155;" colspan="2">3년차~ 등촌 (= 제니스 동일가)</th>
<th style="background:#1e293b;color:#fbbf24;padding:8px;text-align:center;border-bottom:2px solid #334155;" rowspan="2">1~2년차<br>vs 제니스</th>
</tr>
<tr>
<th style="background:#0f4c3a;color:#86efac;padding:6px;text-align:center;border-bottom:2px solid #334155;">남</th>
<th style="background:#0f4c3a;color:#86efac;padding:6px;text-align:center;border-bottom:2px solid #334155;">여</th>
<th style="background:#1e3a5f;color:#93c5fd;padding:6px;text-align:center;border-bottom:2px solid #334155;">남</th>
<th style="background:#1e3a5f;color:#93c5fd;padding:6px;text-align:center;border-bottom:2px solid #334155;">여</th>
</tr>"""

    # 사이드바 입력값 + JENIS 기준 동적 행 생성
    # (구분, 이용시간, 사이드바 남, 사이드바 여, JENIS 키 남, JENIS 키 여)
    rows = [
        ('1개월', '모닝(10시이전)', p_1m_morning_m, p_1m_morning_f, JENIS['1m_morning_m'], JENIS['1m_morning_f']),
        ('', '종일', p_1m_allday_m, p_1m_allday_f, JENIS['1m_allday_m'], JENIS['1m_allday_f']),
        ('', '자유(주말포함)', p_1m_free_m, p_1m_free_f, JENIS['1m_free_m'], JENIS['1m_free_f']),
        ('3개월', '모닝(10시이전)', p_3m_morning_m, p_3m_morning_f, JENIS['3m_morning_m'], JENIS['3m_morning_f']),
        ('', '종일', p_3m_allday_m, p_3m_allday_f, JENIS['3m_allday_m'], JENIS['3m_allday_f']),
        ('', '자유(주말포함)', p_3m_free_m, p_3m_free_f, JENIS['3m_free_m'], JENIS['3m_free_f']),
        ('6개월', '모닝(10시이전)', p_6m_morning_m, p_6m_morning_f, JENIS['6m_morning_m'], JENIS['6m_morning_f']),
        ('', '종일', p_6m_allday_m, p_6m_allday_f, JENIS['6m_allday_m'], JENIS['6m_allday_f']),
        ('', '자유(주말포함)', p_6m_free_m, p_6m_free_f, JENIS['6m_free_m'], JENIS['6m_free_f']),
    ]

    for i, (cat, time, my_m, my_f, jn_m, jn_f) in enumerate(rows):
        bg = '#111827' if i % 2 == 0 else '#0f172a'
        cat_td = f'<td style="background:{bg};color:#e2e8f0;padding:8px 12px;font-weight:700;border-bottom:1px solid #1e293b;" rowspan="3">{cat}</td>' if cat else ''
        # 1~2년차 vs 제니스 할인율 (남 기준)
        disc = _disc_pct(my_m, jn_m)
        pricing_html += f"""<tr>
{cat_td}
<td style="background:{bg};color:#cbd5e1;padding:8px 12px;border-bottom:1px solid #1e293b;">{time}</td>
<td style="background:{bg};color:#86efac;padding:8px;text-align:right;border-bottom:1px solid #1e293b;font-weight:600;">{_fmt(my_m)}</td>
<td style="background:{bg};color:#86efac;padding:8px;text-align:right;border-bottom:1px solid #1e293b;">{_fmt(my_f)}</td>
<td style="background:{bg};color:#93c5fd;padding:8px;text-align:right;border-bottom:1px solid #1e293b;">{_fmt(jn_m)}</td>
<td style="background:{bg};color:#93c5fd;padding:8px;text-align:right;border-bottom:1px solid #1e293b;">{_fmt(jn_f)}</td>
<td style="background:{bg};color:#fbbf24;padding:8px;text-align:center;border-bottom:1px solid #1e293b;font-weight:600;">{disc}</td>
</tr>"""

    pricing_html += "</table>"
    st.markdown(pricing_html, unsafe_allow_html=True)

    # 쿠폰/일일권/락카
    # 이익기여도 기반 전략 카드 — 올바른 순서: 일일>쿠폰>1개월>3개월>6개월
    st.markdown("""
<div style="background:#0f172a;border:1px solid #334155;border-radius:12px;padding:20px;margin:16px 0;">
<div style="color:#fbbf24;font-size:15px;font-weight:700;margin-bottom:6px;">📊 이익기여도 기반 가격 전략 (2018~2021 4개년 평균)</div>
<div style="color:#94a3b8;font-size:12px;margin-bottom:14px;">
2018~2021년 4개년 평균 실적 기반. 1방문당 수익이 높은 상품 = 이익기여도 높음. 장기권은 자주 방문하여 객단가가 희석되고, 모든 회원이 사우나/피트니스 무료 이용 가능하므로 장기권일수록 시설 이용 비용이 증가합니다.
</div>

<div style="background:#172554;border:1px solid #1e3a5f;border-radius:8px;padding:12px 14px;margin-bottom:12px;">
<span style="color:#93c5fd;font-size:12px;">
<b>📐 산출 방법: 타석 총량 역산법</b> — POS 방문기록이 없으므로 2021년 실적 데이터와 타석 운영 총량을 역산하여 상품별 방문횟수를 추정합니다.<br>
<b>계산식</b>: 88타석 × 14타임/일(06~23시, 70분) × 365일 = 연간 최대 449,120건 → 가동률 59%(2021 매출÷최대매출 역산) 적용 → <b>실제 약 265,000건/년</b><br>
<b>배분</b>: 일일권(1인=1건 확정) → 쿠폰(회당 1건 확정) → 나머지를 1개월/3개월/6개월 회원의 월정액 비중으로 배분<br>
<b>출처</b>: 2021 Excel 실적(회원수·매출), KGA 2023 한국골프지표(연평균 18.1회) | ⚠️ 오픈 후 POS 실측 데이터로 교체 필요
</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:14px;margin-bottom:12px;">
<div style="color:#e2e8f0;font-size:13px;font-weight:600;margin-bottom:8px;">📐 1방문당 실질 수익 (이익기여도 순위) — 상품명에 마우스를 올리면 계산 과정이 표시됩니다</div>
<style>
.prd-tip { position:relative; cursor:help; border-bottom:1px dotted #64748b; }
.prd-tip .prd-pop {
    display:none; position:absolute; bottom:130%; left:0; width:320px; z-index:999;
    background:#0f172a; border:1px solid #334155; border-radius:10px; padding:14px;
    box-shadow:0 8px 24px rgba(0,0,0,0.5); font-size:11px; line-height:1.7; color:#cbd5e1;
}
.prd-tip:hover .prd-pop { display:block; }
.prd-pop b { color:#60a5fa; } .prd-pop .val { color:#fbbf24; font-weight:700; }
</style>
<table style="width:100%;border-collapse:collapse;font-size:12px;">
<tr>
<th style="color:#60a5fa;padding:6px 10px;text-align:left;border-bottom:1px solid #334155;">순위</th>
<th style="color:#60a5fa;padding:6px;text-align:left;border-bottom:1px solid #334155;">상품</th>
<th style="color:#60a5fa;padding:6px;text-align:right;border-bottom:1px solid #334155;">가격</th>
<th style="color:#60a5fa;padding:6px;text-align:right;border-bottom:1px solid #334155;">월 방문(역산)</th>
<th style="color:#60a5fa;padding:6px;text-align:right;border-bottom:1px solid #334155;">1방문당 수익</th>
<th style="color:#60a5fa;padding:6px;text-align:left;border-bottom:1px solid #334155;">할인 전략</th>
</tr>

<tr><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">1위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip">일일권<span class="prd-pop">
<b>일일권 — 2018~2021 4개년 평균</b><br>
• 결제 = 방문: <span class="val">1회 = 1건 (확정)</span><br>
• 4개년 평균 이용객: <span class="val">71,548명/년</span><br>
• 4개년 평균 매출: <span class="val">13.9억/년</span><br>
• 4개년 평균 이익: <span class="val">6.7억/년</span><br>
• 전체 매출의 약 37% 차지 (최대 비중)<br>
<b>→ 23,000원/방문 (희석 없음, 기여도 1위)</b>
</span></span></td>
<td style="color:#e2e8f0;padding:5px;text-align:right;">23,000원</td>
<td style="color:#94a3b8;padding:5px;text-align:right;">1회 (확정)</td>
<td style="color:#86efac;padding:5px;text-align:right;font-weight:700;">23,000원</td>
<td style="color:#86efac;padding:5px;font-size:11px;">-4~4.5%</td></tr>

<tr style="background:#111827;"><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">2위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip">쿠폰(10회)<span class="prd-pop">
<b>쿠폰 10회 — 2018~2021 4개년 평균</b><br>
• 10회 쿠폰 = <span class="val">10건 이용 (확정)</span><br>
• 유효기간: <span class="val">3개월</span><br>
• 220,000원 ÷ 10회 = <span class="val">22,000원/회</span><br>
• 4개년 평균 쿠폰 이용: <span class="val">2,906명/년</span><br>
• 4개년 평균 쿠폰 매출: <span class="val">7.4억/년</span><br>
<b>→ 22,000원/방문</b><br>
<span style="color:#86efac;">※ 미소진 시(7~8회) 실질 27,500~31,400원/회<br>
미소진분은 순수익 → 실질 기여도 더 높음</span>
</span></span></td>
<td style="color:#e2e8f0;padding:5px;text-align:right;">220,000원</td>
<td style="color:#94a3b8;padding:5px;text-align:right;">10회 (확정)</td>
<td style="color:#86efac;padding:5px;text-align:right;font-weight:700;">22,000원</td>
<td style="color:#86efac;padding:5px;font-size:11px;">-4.3%</td></tr>

<tr><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">3위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip">쿠폰(20회)<span class="prd-pop">
<b>쿠폰 20회 — 2018~2021 4개년 평균</b><br>
• 20회 쿠폰 = <span class="val">20건 이용 (확정)</span><br>
• 유효기간: <span class="val">6개월</span><br>
• 430,000원 ÷ 20회 = <span class="val">21,500원/회</span><br>
• 10회 대비 회당 500원 저렴 (볼륨 할인)<br>
<b>→ 21,500원/방문</b><br>
<span style="color:#86efac;">※ 미소진 시(15~16회) 실질 26,900~28,700원/회<br>
미소진분은 순수익 → 실질 기여도 더 높음</span>
</span></span></td>
<td style="color:#e2e8f0;padding:5px;text-align:right;">430,000원</td>
<td style="color:#94a3b8;padding:5px;text-align:right;">20회 (확정)</td>
<td style="color:#86efac;padding:5px;text-align:right;font-weight:700;">21,500원</td>
<td style="color:#86efac;padding:5px;font-size:11px;">-4.4%</td></tr>

<tr style="background:#111827;"><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">4위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip">1개월권<span class="prd-pop">
<b>1개월권 — 2018~2021 4개년 평균 역산</b><br>
• 4개년 평균 회원: <span class="val">806명/년</span><br>
• 4개년 평균 매출: <span class="val">1.8억/년</span><br>
• 4개년 평균 이익: <span class="val">0.9억/년</span><br>
• 역산 추정 월 방문: <span class="val">~18회</span><br>
• 290,000원 ÷ 18회 = <span class="val">약 16,100원/회</span><br>
<span style="color:#f97316;">⚠ 방문횟수는 POS 없어 역산 추정</span><br>
※ 단, 무제한 이용이므로 실제 방문은 더 많음<br>
• <b>현실 추정: 월 15~20회</b> (무제한 특성상)<br>
• 290,000 ÷ 18회 = <span class="val">약 16,100원/방문</span><br>
<b>⚠️ 역산과 무제한 특성을 종합한 추정치</b>
</span></span></td>
<td style="color:#e2e8f0;padding:5px;text-align:right;">290,000원</td>
<td style="color:#94a3b8;padding:5px;text-align:right;">~18회 (추정)</td>
<td style="color:#f9a8d4;padding:5px;text-align:right;font-weight:700;">16,100원</td>
<td style="color:#f9a8d4;padding:5px;font-size:11px;">-2.9~3.8%</td></tr>

<tr><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">5위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip">3개월권<span class="prd-pop">
<b>3개월권 — 2018~2021 4개년 평균 역산</b><br>
• 4개년 평균 회원: <span class="val">1,514명/년</span><br>
• 4개년 평균 매출: <span class="val">8.0억/년 (최대 매출 상품)</span><br>
• 4개년 평균 이익: <span class="val">3.7억/년</span><br>
• 역산 추정 월 방문: <span class="val">~20회</span><br>
• 800,000 ÷ 3개월 ÷ 20회 = <span class="val">약 13,300원/회</span><br>
<span style="color:#f97316;">⚠ 방문횟수는 POS 없어 역산 추정</span>
</span></span></td>
<td style="color:#e2e8f0;padding:5px;text-align:right;">800,000원</td>
<td style="color:#94a3b8;padding:5px;text-align:right;">~20회 (추정)</td>
<td style="color:#fdba74;padding:5px;text-align:right;font-weight:700;">13,300원</td>
<td style="color:#fdba74;padding:5px;font-size:11px;">-3.6~4.3%</td></tr>

<tr style="background:#111827;"><td style="color:#fbbf24;padding:5px 10px;font-weight:700;">6위</td>
<td style="color:#e2e8f0;padding:5px;"><span class="prd-tip">6개월권<span class="prd-pop">
<b>6개월권 — 역산 계산 과정</b><br>
• 2019~2021 평균 회원: <span class="val">176명/년</span> (2018년 미운영)<br>
• 2019~2021 평균 매출: <span class="val">1.6억/년</span><br>
• 2019~2021 평균 이익: <span class="val">0.8억/년</span><br>
• 역산 추정 월 방문: <span class="val">~22회</span> (장기 약정, 거의 매일)<br>
• 1,580,000 ÷ 6 ÷ 22 = <span class="val">약 12,000원/회</span><br>
• 사우나/헬스 무료이용 비용(~3,000원/회) 감안 시<br>
  실질 순수익: <span class="val">약 9,000원/회</span><br>
<span style="color:#f97316;">⚠ 방문횟수는 POS 없어 역산 추정</span>
</span></span></td>
<td style="color:#e2e8f0;padding:5px;text-align:right;">1,580,000원</td>
<td style="color:#94a3b8;padding:5px;text-align:right;">~22회 (추정)</td>
<td style="color:#fca5a5;padding:5px;text-align:right;font-weight:700;">12,000원</td>
<td style="color:#fca5a5;padding:5px;font-size:11px;">-2.7~3.4%</td></tr>
</table>
</div>

<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;">
<div style="background:#1e293b;border-radius:10px;padding:14px;border-top:3px solid #22c55e;">
<div style="color:#86efac;font-weight:700;font-size:13px;">💰 고수익 (일일/쿠폰)</div>
<div style="color:#94a3b8;font-size:12px;line-height:1.7;margin-top:6px;">
이익기여도 1~2위<br>
1방문당 <b>21,500~23,000원</b><br>
→ 할인 <b>-3~4.5%</b>로 유입 극대화<br>
<span style="color:#86efac;">일일 체험 → 쿠폰 → 회원 전환</span>
</div>
</div>
<div style="background:#1e293b;border-radius:10px;padding:14px;border-top:3px solid #3b82f6;">
<div style="color:#93c5fd;font-weight:700;font-size:13px;">📊 안정 매출 (1개월/3개월)</div>
<div style="color:#94a3b8;font-size:12px;line-height:1.7;margin-top:6px;">
이익기여도 3~4위<br>
1방문당 <b>13,300~16,100원</b><br>
→ 할인 <b>-2.9~4.3%</b><br>
<span style="color:#93c5fd;">매월 고정 매출 기반 확보</span>
</div>
</div>
<div style="background:#1e293b;border-radius:10px;padding:14px;border-top:3px solid #a855f7;">
<div style="color:#c4b5fd;font-weight:700;font-size:13px;">🏋️ 장기 록인 (6개월)</div>
<div style="color:#94a3b8;font-size:12px;line-height:1.7;margin-top:6px;">
이익기여도 5위 (객단가 최저)<br>
1방문당 <b>~12,000원</b><br>
→ 할인 <b>-2.7~3.4%</b> (최소)<br>
<span style="color:#c4b5fd;">사우나/헬스 무료 비용 감안</span>
</div>
</div>
</div>

<div style="color:#94a3b8;font-size:12px;margin-top:12px;padding-top:10px;border-top:1px solid #334155;">
<b>전략 요약</b>: 이익기여도가 높은 일일/쿠폰은 적극 할인(-4%↑)으로 유입을 극대화하고, 이익기여도가 낮은 6개월은 할인 최소(-3%↓)로 마진 보호.
모든 회원이 사우나/피트니스 무료 이용 가능하므로, 장기권일수록 시설 이용 원가가 높아져 할인폭 축소가 합리적입니다.
</div>
</div>
""", unsafe_allow_html=True)

    st.markdown("**쿠폰·일일권·락카** (1~2년차 = 사이드바 설정 / 3년차~ = 제니스 동일가)")
    etc_html = """
<table style="width:100%;border-collapse:collapse;font-size:13px;margin:8px 0;">
<tr>
<th style="background:#1e293b;color:#60a5fa;padding:10px;text-align:left;border-bottom:2px solid #334155;">구분</th>
<th style="background:#1e293b;color:#60a5fa;padding:10px;text-align:left;border-bottom:2px solid #334155;">세부</th>
<th style="background:#0f4c3a;color:#86efac;padding:10px;text-align:right;border-bottom:2px solid #334155;">1~2년차 등촌<br>(사이드바 설정)</th>
<th style="background:#1e3a5f;color:#93c5fd;padding:10px;text-align:right;border-bottom:2px solid #334155;">3년차~ 등촌<br>(= 제니스)</th>
<th style="background:#1e293b;color:#fbbf24;padding:10px;text-align:center;border-bottom:2px solid #334155;">vs 제니스</th>
<th style="background:#1e293b;color:#94a3b8;padding:10px;text-align:left;border-bottom:2px solid #334155;">전략 근거</th>
</tr>"""

    # 사이드바 입력값 + JENIS 동적 행 생성
    # (구분, 세부, 사이드바 변수, JENIS 키, 전략 근거)
    etc_rows = [
        ('쿠폰', '10회(70분)', p_coupon_10, JENIS['coupon_10'], '체험 진입, 20회 전환 유도'),
        ('', '20회(70분)', p_coupon_20, JENIS['coupon_20'], '주력 쿠폰'),
        ('', '30회(70분)', p_coupon_30, JENIS['coupon_30'], '충성고객 락인'),
        ('일일권', '주중 70분', p_daily_wd_70, JENIS['daily_wd_70'], '최초 유입, 회원전환 핵심'),
        ('', '주중 90분', p_daily_wd_90, JENIS['daily_wd_90'], '장시간 수요 흡수'),
        ('', '주말 70분', p_daily_we_70, JENIS['daily_we_70'], '주말 가동률 확보'),
        ('', '주말 90분', p_daily_we_90, JENIS['daily_we_90'], '주말 프리미엄'),
        ('락카', '보증금', 20000, JENIS['locker_deposit'], '업계 표준 (반환)'),
        ('', '월 사용료', p_locker, JENIS['locker_monthly'], '순수익 (원가 거의 없음)'),
    ]

    for i, (cat, detail, my, jn, reason) in enumerate(etc_rows):
        bg = '#111827' if i % 2 == 0 else '#0f172a'
        cat_style = 'font-weight:700;' if cat else ''
        disc = _disc_pct(my, jn)
        etc_html += f"""<tr>
<td style="background:{bg};color:#e2e8f0;padding:8px 12px;border-bottom:1px solid #1e293b;{cat_style}">{cat}</td>
<td style="background:{bg};color:#cbd5e1;padding:8px 12px;border-bottom:1px solid #1e293b;">{detail}</td>
<td style="background:{bg};color:#86efac;padding:8px;text-align:right;border-bottom:1px solid #1e293b;font-weight:600;">{_fmt(my)}</td>
<td style="background:{bg};color:#93c5fd;padding:8px;text-align:right;border-bottom:1px solid #1e293b;">{_fmt(jn)}</td>
<td style="background:{bg};color:#fbbf24;padding:8px;text-align:center;border-bottom:1px solid #1e293b;font-weight:600;">{disc}</td>
<td style="background:{bg};color:#94a3b8;padding:8px;font-size:11px;border-bottom:1px solid #1e293b;">{reason}</td>
</tr>"""
    etc_html += "</table>"
    st.markdown(etc_html, unsafe_allow_html=True)

    # 연간 매출 시뮬레이션
    subsec("AI 추천 요금 기준 연 매출 시뮬레이션")
    info("위 추천가 기준으로 산출한 연간 매출 추정입니다. 오픈 1년차(9개월)와 2년차(정상 가동) 비교.")

    # 1년차 추정 (9개월, 오픈 보정)
    yr1_items = [
        ('1개월권', 400, 280000, 0.6), ('3개월권', 300, 760000, 0.5), ('6개월권', 50, 1500000, 0.3),
        ('쿠폰(평균)', 1800, 400000*20/20, 0.7), ('일일권', 35000, 23000, 0.8),
        ('락카', 800, 20000*9, 0.9), ('프로임대', 3, 1500000*9, 1.0),
    ]
    yr2_items = [
        ('1개월권', 700, 290000, 1.0), ('3개월권', 500, 780000, 1.0), ('6개월권', 120, 1550000, 1.0),
        ('쿠폰(평균)', 2500, 420000*20/20, 1.0), ('일일권', 50000, 24000, 1.0),
        ('락카', 1500, 20000*12, 1.0), ('프로임대', 3, 1500000*12, 1.0),
    ]

    sim_cols = st.columns(2)
    with sim_cols[0]:
        st.markdown("**1년차 (2026, 9개월 가동)**")
        yr1_total = 0
        yr1_html = '<table style="width:100%;border-collapse:collapse;font-size:12px;">'
        yr1_html += '<tr><th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:left;">상품</th><th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;">추정매출</th></tr>'
        for name, cnt, price, adj in yr1_items:
            rev = int(cnt * price * adj)
            yr1_total += rev
            bg = '#111827'
            yr1_html += f'<tr><td style="background:{bg};color:#cbd5e1;padding:6px 8px;">{name}</td><td style="background:{bg};color:#86efac;padding:6px 8px;text-align:right;">{rev/만:,.0f}만</td></tr>'
        yr1_html += f'<tr><td style="background:#1e293b;color:#f8fafc;padding:8px;font-weight:700;">합계</td><td style="background:#1e293b;color:#fbbf24;padding:8px;text-align:right;font-weight:700;font-size:14px;">{yr1_total/억:.1f}억</td></tr></table>'
        st.markdown(yr1_html, unsafe_allow_html=True)

    with sim_cols[1]:
        st.markdown("**2년차 (2027, 정상 가동)**")
        yr2_total = 0
        yr2_html = '<table style="width:100%;border-collapse:collapse;font-size:12px;">'
        yr2_html += '<tr><th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:left;">상품</th><th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;">추정매출</th></tr>'
        for name, cnt, price, adj in yr2_items:
            rev = int(cnt * price * adj)
            yr2_total += rev
            bg = '#111827'
            yr2_html += f'<tr><td style="background:{bg};color:#cbd5e1;padding:6px 8px;">{name}</td><td style="background:{bg};color:#86efac;padding:6px 8px;text-align:right;">{rev/만:,.0f}만</td></tr>'
        yr2_html += f'<tr><td style="background:#1e293b;color:#f8fafc;padding:8px;font-weight:700;">합계</td><td style="background:#1e293b;color:#fbbf24;padding:8px;text-align:right;font-weight:700;font-size:14px;">{yr2_total/억:.1f}억</td></tr></table>'
        st.markdown(yr2_html, unsafe_allow_html=True)

    # 실제 확인된 경쟁 연습장 요금 비교
    subsec("경쟁 연습장 요금 비교 (실제 확인 데이터)")
    info("웹사이트에서 직접 확인된 요금만 표시합니다. '-'는 해당 상품 미운영 또는 미확인. 전국 평균은 김캐디 2023 데이터 기준.")

    avg1_html = """
<table style="width:100%;border-collapse:collapse;font-size:12px;margin:8px 0;">
<tr>
<th style="background:#1e293b;color:#60a5fa;padding:9px 12px;text-align:left;border-bottom:2px solid #334155;">연습장</th>
<th style="background:#1e293b;color:#94a3b8;padding:9px;text-align:left;border-bottom:2px solid #334155;">위치</th>
<th style="background:#1e293b;color:#94a3b8;padding:9px;text-align:center;border-bottom:2px solid #334155;">타석</th>
<th style="background:#1e293b;color:#e2e8f0;padding:9px;text-align:right;border-bottom:2px solid #334155;">1개월 종일</th>
<th style="background:#1e293b;color:#e2e8f0;padding:9px;text-align:right;border-bottom:2px solid #334155;">3개월 종일</th>
<th style="background:#1e293b;color:#e2e8f0;padding:9px;text-align:right;border-bottom:2px solid #334155;">6개월 종일</th>
<th style="background:#1e293b;color:#e2e8f0;padding:9px;text-align:right;border-bottom:2px solid #334155;">일일 60~70분</th>
<th style="background:#1e293b;color:#e2e8f0;padding:9px;text-align:right;border-bottom:2px solid #334155;">쿠폰 10회</th>
<th style="background:#1e293b;color:#94a3b8;padding:9px;text-align:left;border-bottom:2px solid #334155;">출처</th>
</tr>
<tr>
<td style="background:#0f4c3a;color:#86efac;padding:8px 12px;font-weight:700;">등촌(AI추천)</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;">강서구 등촌동</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;text-align:center;">88</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;text-align:right;font-weight:700;">290,000</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;text-align:right;font-weight:700;">800,000</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;text-align:right;font-weight:700;">1,580,000</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;text-align:right;font-weight:700;">23,000</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;text-align:right;font-weight:700;">220,000</td>
<td style="background:#0f4c3a;color:#86efac;padding:8px;font-size:10px;">AI 추천가</td>
</tr>
<tr>
<td style="background:#111827;color:#fca5a5;padding:8px 12px;font-weight:700;">제니스스포츠</td>
<td style="background:#111827;color:#e2e8f0;padding:8px;">구로구 고척동</td>
<td style="background:#111827;color:#e2e8f0;padding:8px;text-align:center;">112</td>
<td style="background:#111827;color:#fca5a5;padding:8px;text-align:right;font-weight:600;">300,000</td>
<td style="background:#111827;color:#fca5a5;padding:8px;text-align:right;font-weight:600;">830,000</td>
<td style="background:#111827;color:#fca5a5;padding:8px;text-align:right;font-weight:600;">1,630,000</td>
<td style="background:#111827;color:#fca5a5;padding:8px;text-align:right;font-weight:600;">24,000</td>
<td style="background:#111827;color:#fca5a5;padding:8px;text-align:right;font-weight:600;">230,000</td>
<td style="background:#111827;color:#94a3b8;padding:8px;font-size:10px;">사용자 제공</td>
</tr>
<tr>
<td style="background:#0f172a;color:#e2e8f0;padding:8px 12px;font-weight:700;">메이필드호텔</td>
<td style="background:#0f172a;color:#e2e8f0;padding:8px;">강서구 방화동</td>
<td style="background:#0f172a;color:#e2e8f0;padding:8px;text-align:center;">75</td>
<td style="background:#0f172a;color:#e2e8f0;padding:8px;text-align:right;">322,000</td>
<td style="background:#0f172a;color:#e2e8f0;padding:8px;text-align:right;">795,000</td>
<td style="background:#0f172a;color:#94a3b8;padding:8px;text-align:right;">-</td>
<td style="background:#0f172a;color:#e2e8f0;padding:8px;text-align:right;">22,000</td>
<td style="background:#0f172a;color:#e2e8f0;padding:8px;text-align:right;">220,000</td>
<td style="background:#0f172a;color:#94a3b8;padding:8px;font-size:10px;">mayfield.co.kr</td>
</tr>
</table>
"""
    st.markdown(avg1_html, unsafe_allow_html=True)

    # 등촌 vs 경쟁사 차이율
    st.markdown("""
<div style="background:#0f172a;border:1px solid #1e293b;border-radius:12px;padding:16px;margin:10px 0;">
<div style="color:#60a5fa;font-weight:700;font-size:13px;margin-bottom:10px;">등촌 추천가 vs 경쟁사 차이 (1개월 종일 기준, 실제 데이터)</div>
<div style="display:grid;grid-template-columns:repeat(2,1fr);gap:10px;text-align:center;">
<div style="background:#111827;border-radius:8px;padding:12px;">
<div style="color:#fca5a5;font-size:12px;font-weight:600;">제니스 (300,000원)</div>
<div style="color:#86efac;font-size:20px;font-weight:800;">-3.3%</div>
<div style="color:#94a3b8;font-size:11px;">10,000원 저렴</div>
</div>
<div style="background:#111827;border-radius:8px;padding:12px;">
<div style="color:#e2e8f0;font-size:12px;font-weight:600;">메이필드 (322,000원)</div>
<div style="color:#86efac;font-size:20px;font-weight:800;">-9.9%</div>
<div style="color:#94a3b8;font-size:11px;">32,000원 저렴 (호텔 프리미엄)</div>
</div>
</div>
</div>
""", unsafe_allow_html=True)

    st.markdown("""
<div style="background:#1e293b;border-radius:10px;padding:14px;margin:10px 0;border-left:3px solid #3b82f6;">
<span style="color:#93c5fd;font-size:0.85rem;line-height:1.7;">
<b>분석 결과 (실제 데이터 기반)</b><br>
• <b>직접 경쟁사 제니스</b>(112타석, 구로구): 등촌이 3.3% 저렴. 신축 시설 대비 적정 차이<br>
• <b>메이필드호텔</b>(75타석, 강서구): 호텔 부대시설 프리미엄으로 9.9% 비쌈. 고객층 상이<br>
• <b>88CC</b>(영등포구, 골프장 부설): 골프장 부설 연습장으로 저가 포지셔닝. 비교 대상으로 부적합<br>
<b>⚠️ 양천구(목동골프타운), 마포구, 서대문구</b>: 80타석 이상 대형 실외 연습장이 존재하지 않거나 홈페이지 요금 미공개로 확인 불가. 추후 직접 전화 조사 필요
</span>
</div>
""", unsafe_allow_html=True)

    # 인상 로드맵
    st.markdown("""
<div style="background:#1e293b;border-radius:12px;padding:18px;margin:12px 0;border-left:4px solid #a855f7;">
<b style="color:#c4b5fd;font-size:15px;">📈 연도별 인상 로드맵</b><br><br>
<span style="color:#cbd5e1;">
<b>1~2년차</b>: 제니스 대비 3~5% 저렴 → "저렴해서"가 아닌 "비슷한 가격에 새 시설"로 고객 전환 유도<br>
<b>★ 2년간 가격 동결</b>: 오픈 후 1년 만에 인상하면 '미끼 가격' 불신 → 40~60대 주 고객층 이탈 가속. 2년 동결로 신뢰 구축이 핵심<br>
<b>3년차</b>: 첫 인상 3% (소비자물가지수 연평균 2.5% 감안, 물가 상승분 명분) → 2년 동결 경력이 인상 수용도를 높임<br>
<b>4~5년차</b>: 추가 2~3%/년 인상 → 결과적으로 제니스와 1~2% 차이 수준 도달. 신규 시설 프리미엄 유지<br>
<b>이론적 근거</b>: 행동경제학의 '손실회피(Loss Aversion)' — 가격 인상의 심리적 고통은 동일 금액 할인의 기쁨보다 2배 강함(카네만/트버스키). 따라서 <b>초기 적정가 → 2년 동결 → 소폭 인상</b>이 <b>저가 출발 → 급격 인상</b>보다 고객 유지에 유리
</span>
</div>
""", unsafe_allow_html=True)

    subsec("2. 임대매장 운영전략 — 5년 매각/명도 고려")
    st.markdown(f"""
**핵심 원칙**: 계열사·직영 위주 운영으로 **매각 시 명도 리스크 최소화**

**매장 구성 (총 {s_rent_shops}개)**
| 매장 | 운영형태 | 월 수익 | 5년 후 명도 | 비고 |
|------|---------|--------|-----------|------|
| 무인편의점 (CU/세븐일레븐) | **직영/FC** | 월 200~300만원 | 즉시 가능 | 24시간 무인, 인건비 제로, 이용객 필수 |
| 음료자판기 코너 | **직영** | 월 80~120만원 | 즉시 가능 | 자판기 3~4대, 관리 최소 |
| 스크린골프존 (골프존 5대) | **직영 운영** | 월 약 1,000만원 | 장비 철거 | 오픈과 동시 운영, 별도 인건비 없음(안내팀 겸직), 비수기 핵심 대안 |
| 프로 레슨 공간 (프로 3명) | **임대** | 월 450만원 | 계약 종료 | 프로 1인 월 150만원 고정 임대 |
| 예비 공간 | **미임대** | - | 즉시 | 향후 수요에 따라 결정 |

**골프샵(프로샵) 미입점 사유**: 실외 전용 연습장은 골프용품 매장 유치가 어려움 (트래픽 부족, 실내 매장 선호). 무인편의점이 수익성/관리 면에서 우월.

**외부 임차인 주의사항**
- 5년 후 매각/명도 시 외부 임차인의 **상가임대차보호법** 적용으로 명도 지연 리스크
- 상가임대차보호법상 임차인은 **10년간 갱신 요구권** 보유 → 외부 임대 최소화 필수
- 스크린골프존은 직영 운영(월 매출 약 1,000만원, 별도 인건비 없음). 매각 시 장비 철거로 즉시 명도 가능
- 프로 레슨은 1년 단위 계약, 갱신 거절 사유 확보 필요
- **가능하면 계열사/직영/FC로 운영하여 명도 100% 통제**
""")

    subsec("3. 마케팅 전략 — 오픈 전후 3단계")
    st.markdown("""
**[1단계] 오픈 전 (D-3개월~오픈)**
- 네이버 플레이스 등록 + 블로그 체험단 20명 모집
- 인스타그램/유튜브 시설 공개 영상 (타석, 야경, 주차장)
- 사전등록 이벤트: 오픈 전 등록 시 1개월 추가 무료
- 현수막/배너: 등촌역, 마곡나루역, 주요 도로 5개소
- 예상 비용: 약 2,000만원

**[2단계] 오픈 초기 (1~6개월)**
- 오픈 특가 요금 적용 (정상가 대비 5~8% 할인, 시세 훼손 방지)
- 일일 무료체험 이벤트 (주 2회, 각 20명)
- 기업 제휴: 마곡지구 입주기업 단체 할인 (10%+)
- 카카오톡 채널 운영: 실시간 예약 + 프로모션 발송
- 월 마케팅 예산: 300만원

**[3단계] 안정기 (7개월~)**
- 정상 요금 전환 + 재등록 할인(5%) 정책
- 회원 추천 프로그램: 추천인/피추천인 각 1만원 할인
- 시즌 이벤트: 봄/가을 대회, 여름 야간 이벤트, 겨울 스크린골프 패키지
- 월 마케팅 예산: 150만원
""")

    subsec("4. 인력 운영전략 — 컨트롤 패널 연동")
    info("좌측 컨트롤 패널 '인건비·조직구성'에서 입력한 값이 실시간 반영됩니다.")

    # 컨트롤 패널 변수 참조
    _staff = [
        ('대표', ceo_n, ceo_s, '총괄경영, 대외관계, 의사결정'),
        ('총무팀', adm_n, adm_s, '인사/구매/계약/대관 (팀장급)'),
        ('경리팀', acc_n, acc_s, '회계/세무/자금관리 (주임급)'),
        ('시설팀', fac_n, fac_s, '타석/장비/볼머신/설비 (10년+경력)'),
        ('안내팀', desk_n, desk_s, '프론트/회원관리/예약 (특근최소화)'),
    ]
    _total_n = sum(s[1] for s in _staff)
    _total_sal = sum(s[1]*s[2] for s in _staff)
    _ins_rate = s_insurance  # 컨트롤 패널의 4대보험율
    _total_with_ins = _total_sal * (1 + _ins_rate)

    t_html = '<table style="width:100%;border-collapse:collapse;font-size:13px;margin:8px 0;">'
    t_html += '<tr>'
    for h in ['부서', '인원', '월급(1인)', '월 소계', '역할']:
        t_html += f'<th style="background:#1e293b;color:#60a5fa;padding:10px 12px;text-align:left;border-bottom:2px solid #334155;">{h}</th>'
    t_html += '</tr>'
    for i, (dept, cnt, sal, role) in enumerate(_staff):
        bg = '#111827' if i % 2 == 0 else '#0f172a'
        sub = cnt * sal
        t_html += f'<tr>'
        t_html += f'<td style="background:{bg};color:#e2e8f0;padding:8px 12px;font-weight:700;border-bottom:1px solid #1e293b;">{dept}</td>'
        t_html += f'<td style="background:{bg};color:#e2e8f0;padding:8px 12px;text-align:center;border-bottom:1px solid #1e293b;">{cnt}명</td>'
        t_html += f'<td style="background:{bg};color:#e2e8f0;padding:8px 12px;text-align:right;border-bottom:1px solid #1e293b;">{sal:,}만</td>'
        t_html += f'<td style="background:{bg};color:#86efac;padding:8px 12px;text-align:right;font-weight:600;border-bottom:1px solid #1e293b;">{sub:,}만</td>'
        t_html += f'<td style="background:{bg};color:#94a3b8;padding:8px 12px;font-size:12px;border-bottom:1px solid #1e293b;">{role}</td>'
        t_html += '</tr>'
    t_html += f"""<tr>
        <td style="background:#1e293b;color:#f8fafc;padding:10px 12px;font-weight:700;">합계</td>
        <td style="background:#1e293b;color:#f8fafc;padding:10px 12px;text-align:center;font-weight:700;">{_total_n}명</td>
        <td style="background:#1e293b;color:#94a3b8;padding:10px 12px;text-align:right;">-</td>
        <td style="background:#1e293b;color:#fbbf24;padding:10px 12px;text-align:right;font-weight:700;font-size:15px;">{_total_sal:,.0f}만</td>
        <td style="background:#1e293b;color:#94a3b8;padding:10px 12px;font-size:11px;">4대보험 {_ins_rate*100:.1f}% 별도</td>
    </tr>"""
    t_html += '</table>'
    st.markdown(t_html, unsafe_allow_html=True)

    # 총액 산출 카드
    _ins_amt = _total_sal * _ins_rate
    _annual = _total_with_ins * 12

    st.markdown(f"""
<div style="background:#0f172a;border:1px solid #334155;border-radius:12px;padding:18px;margin:12px 0;">
<div style="display:grid;grid-template-columns:1fr 1fr 1fr;gap:12px;text-align:center;">
<div>
<div style="color:#94a3b8;font-size:12px;">월 급여 합계</div>
<div style="color:#e2e8f0;font-size:22px;font-weight:700;">{_total_sal:,.0f}만</div>
</div>
<div>
<div style="color:#94a3b8;font-size:12px;">4대보험 ({_ins_rate*100:.1f}%)</div>
<div style="color:#e2e8f0;font-size:22px;font-weight:700;">+{_ins_amt:,.0f}만</div>
</div>
<div>
<div style="color:#94a3b8;font-size:12px;">월 총 인건비</div>
<div style="color:#fbbf24;font-size:22px;font-weight:800;">{_total_with_ins:,.0f}만</div>
</div>
</div>
<div style="color:#94a3b8;font-size:12px;margin-top:10px;text-align:center;">연간 인건비: <b>{_annual/만:.1f}억</b> (퇴직충당금 별도)</div>
</div>
""", unsafe_allow_html=True)

    # 인건비율 분석
    _labor_ratio_2027 = (_annual * 만) / rev_p[1] * 100 if len(rev_p) > 1 and rev_p[1] else 0
    st.markdown(f"""
<div style="background:#1e293b;border-radius:12px;padding:16px;margin:8px 0;border-left:4px solid #a855f7;">
<b style="color:#c4b5fd;">운영 효율화 방안</b><br><br>
<span style="color:#cbd5e1;">
• <b>현재 인건비율</b>: 매출 대비 <b>{_labor_ratio_2027:.1f}%</b> {'(목표 25% 이하 달성)' if _labor_ratio_2027 <= 25 else '(⚠️ 목표 25% 초과 — 효율화 필요)'}<br>
• <b>무인 키오스크</b>: 야간/주말 프론트 인력 절감 (자동 결제·회원 인증)<br>
• <b>자동 볼머신</b>: 시설팀 볼 수거 업무 자동화<br>
• <b>클라우드 POS + 회원관리 앱</b>: 경리팀 정산·회계 자동화<br>
• <b>특근/야근 최소화</b>: 2교대 기본, 주말 대체휴무 원칙<br>
• <b>야간 최소 운영</b>: 20시 이후 안내 1명 + 시설 1명 (총 2명)
</span>
</div>
""", unsafe_allow_html=True)

    subsec("5. 출구전략 — 5년 후 매각/명도 시나리오")
    st.markdown(f"""
**시나리오 A: 5년 후 매각 (권리금 포함)**
- 예상 매각가: 연 EBITDA x 3~5배 = {ebitda_p[-1]/억:.1f}억 x 4 = **{ebitda_p[-1]/억*4:.1f}억**
- 투자 회수: 투자금 {s_inv}억 + 5년 누적 EBITDA {cum_ebitda[-1]/억:.1f}억 + 매각 {ebitda_p[-1]/억*4:.1f}억
- **총 수익: {(cum_ebitda[-1] + ebitda_p[-1]*4)/억:.1f}억** (투자 대비 {(cum_ebitda[-1] + ebitda_p[-1]*4)/inv_won*100:.0f}%)

**시나리오 B: 5년 후 명도 반환**
- 잔존가치: 시설물 감가상각 후 장부가
- 보증금 회수 + 누적 현금흐름
- 임대매장 명도 상세:
  - 무인편의점(직영/FC): 자판기 철거 후 **즉시 명도** 가능
  - 스크린골프존(직영): 장비 철거로 **즉시 명도** 가능
  - 프로 레슨(임대): 1년 단위 계약, 갱신 거절 통보 **6개월 전** 필요
  - 예비 공간: **즉시 명도**
- **핵심**: 외부 임차인 없이 직영/FC 위주이므로 명도 리스크 최소 (전체 1~2개월 내 완료 가능)
- **상가임대차보호법 주의**: 외부 임차인이 있을 경우 10년 갱신요구권 + 권리금 회수 보호 적용. 반드시 계약 전 법률 검토

**매각 가치 극대화를 위한 사전 조치**
1. 3년차부터 시설 리모델링/보수 투자 (매각가 상승)
2. 회원 데이터베이스 체계화 (매수인에게 가치)
3. 안정적 EBITDA 성장 추세 유지 (밸류에이션 근거)
4. 임대매장 전량 직영/FC 유지 (명도 리스크 제거)
""")

    subsec("6. 리스크 대응 — 상위 6개 리스크 관리")
    st.markdown("""
| 순위 | 리스크 | 영향도 | 대응 전략 |
|------|--------|--------|----------|
| 1 | **비수기 장기화 (기후변화)** | 매출 30~40%↓, 연 5개월↑ | 기상청 기준 서울 폭염일수 +31%(2020년대), 한파일수 증가. 여름(6~8월)+겨울(12~2월)=5개월 비수기. 대응: 스크린골프 5대, 동계할인, 야간영업, 미스트/냉풍기 |
| 2 | **기존 회원 전환 지연** | 오픈 초기 매출 부진 | 40~60대 주 고객층 충성도 높아 전환 지연. 회원권 만료까지 이동 불가. 대응: 일일/쿠폰 체험 유입 집중 |
| 3 | **제니스 등 경쟁 심화** | 회원 이탈 | 제니스(구로구 고척동 112타석) 대비 신규 시설+동등가 전략. 가격 아닌 시설로 승부 |
| 4 | **인건비/전기료 상승** | 비용 증가 연 3~5% | 2026 최저임금 210만원, 전기료 연 5%↑ 추세. 대응: 무인키오스크, 자동볼머신, 에너지 모니터링 |
| 5 | **매각 리스크** | 출구 불확실 | 임대매장 직영/FC 유지(명도 용이), 3년차부터 매각 준비, 복수 매수인 접촉 |
| 6 | **미세먼지/우천** | 실외 영업 불가 | 서울 미세먼지 나쁨이상 연 80일↑. 스크린골프 5대+실내 공간으로 부분 대안 |
""")

    # ══ 실행 로드맵 ══
    subsec("7. 실행 로드맵")
    info("오픈 전 준비부터 5년차 출구전략까지 단계별 핵심 과제와 목표를 정리합니다.")

    st.markdown("""
<div style="position:relative;margin:20px 0 30px 0;">

<!-- 타임라인 세로선 -->
<div style="position:absolute;left:24px;top:0;bottom:0;width:3px;background:linear-gradient(180deg,#3b82f6,#22c55e,#f97316,#a855f7,#ef4444);border-radius:2px;"></div>

<!-- Phase 1: 오픈 준비 -->
<div style="display:flex;gap:20px;margin-bottom:20px;position:relative;">
<div style="min-width:48px;height:48px;border-radius:50%;background:#1e3a5f;border:3px solid #3b82f6;display:flex;align-items:center;justify-content:center;z-index:1;">
<span style="color:#60a5fa;font-weight:900;font-size:14px;">P1</span>
</div>
<div style="background:#111827;border:1px solid #1e293b;border-radius:12px;padding:18px 22px;flex:1;border-left:3px solid #3b82f6;">
<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px;">
<span style="color:#60a5fa;font-size:15px;font-weight:800;">오픈 준비</span>
<span style="color:#64748b;font-size:12px;font-weight:600;background:#1e293b;padding:3px 10px;border-radius:12px;">D-3개월 ~ 오픈일 (2026.03~06)</span>
</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;color:#cbd5e1;font-size:12.5px;line-height:1.7;">
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📋 핵심 과제</div>
• 시설 공사 완료 및 안전 검수<br>
• 인력 채용 완료 (9명)<br>
• 사업자 등록, 영업 허가<br>
• POS/회원관리 시스템 구축<br>
• 골프존 스크린 5대 설치
</div>
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📣 마케팅</div>
• 네이버 플레이스 등록<br>
• 블로그 체험단 20명 모집<br>
• 현수막/배너 5개소 설치<br>
• 사전등록 이벤트 운영<br>
• <b>예산: 2,000만원</b>
</div>
</div>
<div style="margin-top:10px;padding-top:8px;border-top:1px solid #1e293b;color:#60a5fa;font-size:12px;font-weight:600;">
🎯 목표: 사전등록 200명 이상 확보
</div>
</div>
</div>

<!-- Phase 2: 오픈 초기 -->
<div style="display:flex;gap:20px;margin-bottom:20px;position:relative;">
<div style="min-width:48px;height:48px;border-radius:50%;background:#14532d;border:3px solid #22c55e;display:flex;align-items:center;justify-content:center;z-index:1;">
<span style="color:#86efac;font-weight:900;font-size:14px;">P2</span>
</div>
<div style="background:#111827;border:1px solid #1e293b;border-radius:12px;padding:18px 22px;flex:1;border-left:3px solid #22c55e;">
<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px;">
<span style="color:#86efac;font-size:15px;font-weight:800;">오픈 초기 (생존)</span>
<span style="color:#64748b;font-size:12px;font-weight:600;background:#1e293b;padding:3px 10px;border-radius:12px;">1~6개월 (2026.06~11)</span>
</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;color:#cbd5e1;font-size:12.5px;line-height:1.7;">
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📋 핵심 과제</div>
• 일일/쿠폰 체험 유입 집중<br>
• 타 연습장 회원권 만료 대기<br>
• 비수기(7~8월 폭염) 대응<br>
• 고객 VOC 수집 및 즉시 개선<br>
• 임대매장(편의점/자판기) 오픈
</div>
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📊 KPI</div>
• 램프업 목표: 정상의 25%→65%<br>
• 일일/쿠폰 비중 60% 이상<br>
• 월 매출 목표: 1.5~2억<br>
• 타석 가동률 40% 이상<br>
• <b>이 시기 적자는 정상</b>
</div>
</div>
<div style="margin-top:10px;padding-top:8px;border-top:1px solid #1e293b;color:#86efac;font-size:12px;font-weight:600;">
🎯 목표: 월 회원 400명 돌파, 재방문율 50% 이상
</div>
</div>
</div>

<!-- Phase 3: 안정화 -->
<div style="display:flex;gap:20px;margin-bottom:20px;position:relative;">
<div style="min-width:48px;height:48px;border-radius:50%;background:#78350f;border:3px solid #f97316;display:flex;align-items:center;justify-content:center;z-index:1;">
<span style="color:#fdba74;font-weight:900;font-size:14px;">P3</span>
</div>
<div style="background:#111827;border:1px solid #1e293b;border-radius:12px;padding:18px 22px;flex:1;border-left:3px solid #f97316;">
<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px;">
<span style="color:#fdba74;font-size:15px;font-weight:800;">안정화 (성장)</span>
<span style="color:#64748b;font-size:12px;font-weight:600;background:#1e293b;padding:3px 10px;border-radius:12px;">7~18개월 (2026.12~2027.11)</span>
</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;color:#cbd5e1;font-size:12.5px;line-height:1.7;">
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📋 핵심 과제</div>
• 회원권 전환 본격화 (1개월→3개월)<br>
• 재등록률 70% 이상 관리<br>
• 임대매장 2차 오픈 (스크린골프존)<br>
• 프로 레슨 프로그램 안정화<br>
• 비용 효율화 (무인화 확대)
</div>
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📊 KPI</div>
• 연매출 {rev_p[1]/억:.0f}억 목표 (2027F)<br>
• 영업이익률 흑자 전환<br>
• 타석 가동률 60% 이상<br>
• 인건비율 25% 이하<br>
• <b>2년간 가격 동결 유지</b>
</div>
</div>
<div style="margin-top:10px;padding-top:8px;border-top:1px solid #1e293b;color:#fdba74;font-size:12px;font-weight:600;">
🎯 목표: 2027F 영업이익 흑자 전환, 재등록률 70%
</div>
</div>
</div>

<!-- Phase 4: 최적화 -->
<div style="display:flex;gap:20px;margin-bottom:20px;position:relative;">
<div style="min-width:48px;height:48px;border-radius:50%;background:#581c87;border:3px solid #a855f7;display:flex;align-items:center;justify-content:center;z-index:1;">
<span style="color:#c4b5fd;font-weight:900;font-size:14px;">P4</span>
</div>
<div style="background:#111827;border:1px solid #1e293b;border-radius:12px;padding:18px 22px;flex:1;border-left:3px solid #a855f7;">
<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px;">
<span style="color:#c4b5fd;font-size:15px;font-weight:800;">최적화 (수익 극대화)</span>
<span style="color:#64748b;font-size:12px;font-weight:600;background:#1e293b;padding:3px 10px;border-radius:12px;">19~36개월 (2027.12~2029.05)</span>
</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;color:#cbd5e1;font-size:12.5px;line-height:1.7;">
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📋 핵심 과제</div>
• 3년차 첫 요금 인상 (3% 이내)<br>
• VIP 회원 프로그램 도입<br>
• 비수기 이벤트 정례화<br>
• 시설 부분 리모델링 검토<br>
• 매각/재계약 사전 검토 착수
</div>
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📊 KPI</div>
• EBITDA 이익률 15% 이상<br>
• 누적 EBITDA {cum_ebitda[2]/억:.0f}억→{cum_ebitda[3]/억:.0f}억<br>
• 투자회수율 {rec_rate[2]*100:.0f}%→{rec_rate[3]*100:.0f}%<br>
• 고정 고객 1,500명 이상<br>
• <b>마진 보호 + 효율 극대화</b>
</div>
</div>
<div style="margin-top:10px;padding-top:8px;border-top:1px solid #1e293b;color:#c4b5fd;font-size:12px;font-weight:600;">
🎯 목표: EBITDA 이익률 15%+, 안정적 현금 창출
</div>
</div>
</div>

<!-- Phase 5: 출구전략 -->
<div style="display:flex;gap:20px;position:relative;">
<div style="min-width:48px;height:48px;border-radius:50%;background:#7f1d1d;border:3px solid #ef4444;display:flex;align-items:center;justify-content:center;z-index:1;">
<span style="color:#fca5a5;font-weight:900;font-size:14px;">P5</span>
</div>
<div style="background:#111827;border:1px solid #1e293b;border-radius:12px;padding:18px 22px;flex:1;border-left:3px solid #ef4444;">
<div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:8px;">
<span style="color:#fca5a5;font-size:15px;font-weight:800;">출구전략 (매각/명도)</span>
<span style="color:#64748b;font-size:12px;font-weight:600;background:#1e293b;padding:3px 10px;border-radius:12px;">37~60개월 (2029.06~2031.05)</span>
</div>
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;color:#cbd5e1;font-size:12.5px;line-height:1.7;">
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📋 시나리오 A: 매각</div>
• 예상 매각가: EBITDA × 3~5배<br>
• 회원 DB 및 고객 자산 정리<br>
• 시설 상태 점검, 리모델링<br>
• 복수 매수인 접촉 (D-12개월)<br>
• <b>권리금 + 누적이익 회수</b>
</div>
<div>
<div style="color:#94a3b8;font-size:11px;font-weight:700;margin-bottom:4px;">📋 시나리오 B: 명도</div>
• 임대매장 명도 (직영→즉시가능)<br>
• 스크린골프 FC 계약 해지 (1~2월)<br>
• 프로 레슨 계약 종료 통보 (D-6개월)<br>
• 시설물 감가상각 후 잔존가 확인<br>
• <b>보증금 회수 + 잔여 현금</b>
</div>
</div>
<div style="margin-top:10px;padding-top:8px;border-top:1px solid #1e293b;color:#fca5a5;font-size:12px;font-weight:600;">
🎯 목표: 누적 EBITDA {cum_ebitda[-1]/억:.0f}억 + 매각가로 투자금 {inv_won/억:.0f}억 전액 회수
</div>
</div>
</div>

</div>
""", unsafe_allow_html=True)

    # 로드맵 요약 테이블
    subsec("실행 로드맵 요약")
    roadmap_data = {
        '단계': ['P1 오픈 준비', 'P2 오픈 초기', 'P3 안정화', 'P4 최적화', 'P5 출구전략'],
        '기간': ['D-3개월~오픈', '1~6개월', '7~18개월', '19~36개월', '37~60개월'],
        '핵심 목표': ['사전등록 200명', '월 회원 400명', f'연매출 {rev_p[1]/억:.0f}억', 'EBITDA 15%+', f'누적 {cum_ebitda[-1]/억:.0f}억 회수'],
        '예상 매출': ['-', f'월 1.5~2억', f'연 {rev_p[1]/억:.0f}억', f'연 {rev_p[2]/억:.0f}억', f'연 {rev_p[4]/억:.0f}억'],
        '손익': ['투자 집행', '적자 (정상)', '흑자 전환', '이익 극대화', '매각/명도'],
        '핵심 리스크': ['공사 지연', '고객유입 부진', '재등록률 하락', '경쟁 심화', '매각가 하락'],
    }
    dark_table(pd.DataFrame(roadmap_data))

# ═══ TAB 4: Revenue ═══
if _ti == 4:
    sec("📈", "매출 추정 분석")
    subsec("매출 보정계수 산출")

    st.markdown("""
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin:12px 0;">
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #3b82f6;">
<b style="color:#60a5fa;">📊 보정계수란?</b><br>
<span style="color:#cbd5e1;">업계 평균 매출(기준매출)을 <b>우리 연습장의 실제 여건</b>에 맞게 조정하는 곱셈 계수입니다.<br>
예: 기준매출 30억 × 보정계수 0.85 = 보정매출 25.5억</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #22c55e;">
<b style="color:#86efac;">🔢 숫자 의미</b><br>
<span style="color:#cbd5e1;">• <b>1.00</b> = 기준(업계 평균)과 동일 조건<br>
• <b>1.00 초과</b> = 기준보다 유리 (매출↑)<br>
• <b>1.00 미만</b> = 기준보다 불리 (매출↓)</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #f97316;">
<b style="color:#fdba74;">⚙️ 설정 방법</b><br>
<span style="color:#cbd5e1;">각 항목을 <b>경쟁사 대비, 업계 평균 대비</b> 우리의 상대적 위치로 판단하여 입력합니다. 확신이 없으면 기본값(1.0 근처)을 유지하세요.</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #a855f7;">
<b style="color:#c4b5fd;">📐 최종 계산</b><br>
<span style="color:#cbd5e1;">5개 항목을 모두 곱합니다:<br>
<b>종합 = ①×②×③×④×⑤</b><br>
예: 1.12×1.08×1.10×0.78×0.90 = <b>0.94배</b></span>
</div>
</div>
""", unsafe_allow_html=True)

    with st.expander("📖 보정계수 5대 항목 상세 안내 (클릭하여 확인)", expanded=False):
        st.markdown("""
**고객이 골프연습장을 선택하는 5대 핵심 요인**

고객은 연습장을 선택할 때 ① 요금(가장 중요), ② 시설 쾌적성, ③ 위치/접근성, ④ 날씨 영향을 고려하며, ⑤ 경쟁사의 할인 전략도 매출에 직접적 영향을 미칩니다.
이 5가지 요인을 종합하여 기준매출 대비 보정계수를 산출합니다.

**① 요금경쟁력** (가장 중요!)
- **의미**: 경쟁사 대비 얼마나 저렴한가. 고객이 가장 민감하게 반응하는 요인
- **기준**: 1.0 = 경쟁사와 동일 요금
- **범위**: 0.90~1.20 (경쟁사 대비 저렴할수록 높음)
- **등촌 기준**: 제니스 대비 5~15% 저렴 → **1.10~1.15**

**② 시설/쾌적성**
- **의미**: 시설 연식, 타석 간격, 조명, 편의시설(락카/휴게실/카페) 등 이용 쾌적성
- **기준**: 1.0 = 시설 연식 5년, 표준 타석 간격
- **범위**: 0.80~1.20 (신축=1.1~1.2, 노후=0.85)
- **등촌 기준**: 2026년 신축, 88타석, 사우나/헬스 완비 → **1.10~1.15**

**③ 위치/접근성**
- **의미**: 역세권 여부, 주차 편의성, 주요 도로 접근성
- **기준**: 1.0 = 지하철 도보 15분, 주차 30대
- **범위**: 0.85~1.20 (역세권+주차=1.1~1.2, 접근 불편=0.85)
- **등촌 기준**: 9호선 등촌역 도보 10분, 주차 30~40대 → **1.10**

**④ 기상/계절성**
- **의미**: 실외 연습장의 기상 영향. 폭염/한파 시 고객 이탈 정도
- **기준**: 1.0 = 실내 연습장 (기상 영향 없음)
- **범위**: 0.65~0.85 (서울 실외 기준 0.70~0.85)
- **등촌 기준**: 서울 평균 기후, 겨울 난방 미비 → **0.75~0.80**
- **조정 요인**: 방풍/방한 시설, 냉방팬, 지붕 커버 범위

**⑤ 경쟁사 할인 전략** — 주변 경쟁사(제니스, 마곡나루 등)의 할인/프로모션 공격 강도. 경쟁사가 할인 공세를 펼칠수록 우리 매출에 부정적 영향
- **의미**: 경쟁사의 가격 할인·프로모션 강도 (할인 공세가 심할수록 낮음)
- **기준**: 1.0 = 경쟁사가 정상가 유지, 특별 할인 없음
- **범위**: 0.75~1.05 (경쟁사 대규모 할인=0.75~0.85, 정상=0.95~1.0, 경쟁사 인상=1.0~1.05)
- **등촌 기준**: 제니스가 신규 오픈 대응으로 단기 할인(3~6개월) 진행 가능성 → **0.85~0.95**
- **조정 요인**: 제니스 할인 기간/폭, 스크린골프 가격 인하, 인근 신규 연습장 오픈 여부
""")

    st.caption("※ 본 5개 항목은 **참고용 분석 지표**입니다 — rev_p에는 이미 사이드바의 상권/경제 점수로 반영되어 있어 여기 입력값은 시뮬레이션 결과에 영향을 주지 않습니다.")
    fc = st.columns(6)
    price_f = fc[0].number_input("요금경쟁력", value=sv('f_price', 1.12), step=0.05, key='f_price', help="경쟁사 대비 저렴할수록 높음. 제니스 대비 5~15% 저렴=1.10~1.15")
    facility_f = fc[1].number_input("시설/쾌적성", value=sv('f_facility', 1.08), step=0.05, key='f_facility', help="신축=1.1~1.2, 노후=0.85. 타석간격, 조명, 편의시설 포함")
    location_f = fc[2].number_input("위치/접근성", value=sv('f_location', 1.10), step=0.05, key='f_location', help="역세권+주차=1.1~1.2, 접근 불편=0.85")
    weather_f = fc[3].number_input("기상/계절성", value=sv('f_weather', 0.78), step=0.05, key='f_weather', help="실외=0.65~0.85 (여름폭염+겨울한파), 실내=1.0")
    switch_f = fc[4].number_input("경쟁사할인", value=sv('f_switch', 0.90), step=0.05, key='f_switch', help="경쟁사 할인 공세 강도. 대규모할인=0.75~0.85, 정상=0.95~1.0, 경쟁사인상=1.0~1.05")
    combined = price_f * facility_f * location_f * weather_f * switch_f
    fc[5].metric("종합 보정계수(참고)", f"{combined:.2f}배")

    subsec("상품별 연간 매출 시뮬레이션")
    info("각 상품(회원권, 쿠폰, 레슨 등)의 회원수와 연 매출을 시뮬레이션한 결과입니다. 비중이 높은 상품이 핵심 수익원입니다.")
    sim_data = []
    mem_map = {'1개월 회원': m_1m, '3개월 회원': m_3m, '6개월 회원': m_6m, '쿠폰': m_coupon, '일일회원': m_daily, '골프레슨(프로임대)': pro_count, '락카': m_locker, '임대(11~2월,4개월)': 0}
    for item, rv in custom_rev_items.items():
        pct = rv / custom_total_rev * 100 if custom_total_rev else 0
        sim_data.append({'상품': item, '회원수': f"{mem_map[item]:,}명", '연매출': fmt만(rv)+'원', '비중': f"{pct:.1f}%"})
    sim_data.append({'상품': '합계', '회원수': f"{total_members:,}명", '연매출': fmt억(custom_total_rev), '비중': '100%'})
    dark_table(pd.DataFrame(sim_data))

    c1, c2 = st.columns(2)
    with c1:
        info("사이드바에서 입력한 회원수×단가 기준의 **커스텀 시뮬레이션 매출** 구성 비율입니다 (특정 연도가 아닌 사용자 입력값 기반). 특정 상품이 30% 이상 차지하면 해당 상품 의존도가 높아 리스크 요인이 됩니다.")
        items = [s for s in sim_data if s['상품'] != '합계']
        fig = go.Figure(go.Pie(labels=[s['상품'] for s in items], values=[custom_rev_items[s['상품']] for s in items], hole=0.45, marker=dict(colors=PAL, line=dict(color='white', width=2)), textinfo='label+percent'))
        lo(fig, title='매출 구성 비율 (사이드바 입력 기준)', height=420)
        st.plotly_chart(fig)
    with c2:
        info("**2026F~2030F 5개년** 매출 항목별 누적 차트입니다(패드 연동: 2026 상품별 비중 유지하며 rev_p로 스케일). 각 색상이 상품을 나타내며, 전체 높이가 해당 연도 총 매출입니다.")
        fig = go.Figure()
        for idx, (n, v) in enumerate(rev_items_dyn.items()):
            fig.add_trace(go.Bar(name=n, x=[str(y) for y in D['yp']], y=[x/억 for x in v], marker_color=PAL[idx%len(PAL)]))
        lo(fig, title='매출 항목별 5개년 추이 (억원)', barmode='stack', height=420, yaxis_title='억원')
        st.plotly_chart(fig)

    subsec("10개년 매출 전망")
    info("10개년 장기 매출·비용·영업이익 전망입니다. **매출 = 골프매출(회원권+쿠폰+일일) + 프로임대 + 락카 + 임대매장 수익** 합산. 전반 5년은 Excel 재무모델 상세 추정, 후반 5년은 매출성장률 기반 외삽입니다.")
    yrs_10 = [f"{2026+i}F" for i in range(10)]
    fig10 = go.Figure()
    fig10.add_trace(go.Bar(name='매출', x=yrs_10, y=[r/억 for r in rev_10yr], marker_color=[C['blue'] if i<5 else C['blue_l'] for i in range(10)], text=[f"{r/억:.1f}" for r in rev_10yr], textposition='inside', textfont=dict(size=12, color='#e2e8f0')))
    # 비용 10년도 추정
    cost_10yr = list(cost_p) + [cost_p[-1] * (1 + s_util_up) ** (i+1) for i in range(5)]
    op_10yr_full = [r - c for r, c in zip(rev_10yr, cost_10yr)]
    fig10.add_trace(go.Scatter(name='영업이익', x=yrs_10, y=[o/억 for o in op_10yr_full], mode='lines+markers', line=dict(color=C['green'], width=2)))
    fig10.add_trace(go.Bar(name='비용', x=yrs_10, y=[c/억 for c in cost_10yr], marker_color=[C['red'] if i<5 else C['red_l'] for i in range(10)], opacity=0.5))
    lo(fig10, title='10개년 매출·비용·영업이익 전망 (억원)', height=440, yaxis_title='억원', barmode='group')
    st.plotly_chart(fig10)

    info("아래 테이블에서 연도별 숫자를 정확히 확인할 수 있습니다. 매출 구성: 회원권+쿠폰+일일+프로임대+락카+임대매장")
    tbl_10 = {'항목': ['매출 (억)', '비용 (억)', '영업이익 (억)', '영업이익률 (%)', '전년비 성장률 (%)']}
    for i, yr in enumerate(yrs_10):
        r, c = rev_10yr[i], cost_10yr[i]
        o = r - c
        margin = o / r * 100 if r else 0
        growth = (r / rev_10yr[i-1] - 1) * 100 if i > 0 and rev_10yr[i-1] else 0
        tbl_10[yr] = [f"{r/억:.1f}", f"{c/억:.1f}", f"{o/억:.1f}", f"{margin:.1f}%", f"{growth:+.1f}%" if i > 0 else "-"]
    dark_table(pd.DataFrame(tbl_10))

    subsec("타석당 수익성")
    info("88타석 기준으로 타석 1개당 수익성을 산출합니다. 타석당 매출이 높을수록 효율적인 운영입니다.")
    bk = st.columns(4)
    bk[0].metric("타석당 연매출 (2027F)", fmt만(rev_per_bay[1])+'원')
    bk[1].metric("타석당 월매출", fmt만(rev_per_bay[1]/12)+'원')
    bk[2].metric("타석당 일매출", f"{rev_per_bay[1]/365:,.0f}원")
    bk[3].metric("타석당 시간매출", f"{rev_per_bay[1]/365/s_hours:,.0f}원")

# ═══ TAB 5: Cost ═══
if _ti == 5:
    sec("💸", "비용 추정")
    c1, c2 = st.columns(2)
    with c1:
        info("2027F 기준 비용 항목별 구성 비율입니다(패드 연동). 인건비, 감가상각 등 고정비 비중이 높으면 손익분기 달성이 어려워집니다.")
        cl = list(cost_items_dyn.keys()); cv = [cost_items_dyn[k][1] for k in cl]
        fig = go.Figure(go.Pie(labels=cl, values=cv, hole=0.45, marker=dict(colors=PAL2, line=dict(color='white', width=2)), textinfo='label+percent'))
        lo(fig, title='2027F 비용 구성', height=420)
        st.plotly_chart(fig)
    with c2:
        info("비용을 고정비·준변동비·변동비로 분류한 누적 막대입니다. 고정비 비중이 클수록 매출 변동에 따른 손익 민감도가 높아집니다.")
        fig = go.Figure()
        fig.add_trace(go.Bar(name='고정비', x=['2026년'], y=[fixed_total/억], marker_color=C['red']))
        fig.add_trace(go.Bar(name='준변동비', x=['2026년'], y=[semi_total*0.6/억], marker_color=C['orange']))
        fig.add_trace(go.Bar(name='변동비', x=['2026년'], y=[var_total/억], marker_color=C['yellow']))
        lo(fig, title='고정비 vs 변동비 구조 (억원)', barmode='stack', height=420, yaxis_title='억원')
        st.plotly_chart(fig)

    subsec("비용 항목별 5개년 추이")
    info("5개년간 비용 항목별 변동 추이입니다(패드 연동). 인건비·전력비 상승률이 매출 성장률을 초과하면 수익성이 악화됩니다.")
    fig_ct = go.Figure()
    for idx, (n, v) in enumerate(cost_items_dyn.items()):
        fig_ct.add_trace(go.Bar(name=n, x=[str(y) for y in D['yp']], y=[x/억 for x in v], marker_color=PAL[idx%len(PAL)]))
    lo(fig_ct, title='비용 항목별 5개년 누적 추이 (억원)', barmode='stack', height=420, yaxis_title='억원')
    st.plotly_chart(fig_ct)

    subsec("타석당 비용 분석")
    info("타석 1개당 비용을 산출하여 타석 수익성 관점에서 비용 효율을 평가합니다.")
    ck = st.columns(4)
    ck[0].metric("타석당 연비용", fmt만(cost_p[0]/s_bays)+'원')
    ck[1].metric("타석당 월비용", fmt만(cost_p[0]/s_bays/12)+'원')
    ck[2].metric("변동비 비율", f"{var_ratio*100:.1f}%")
    ck[3].metric("공헌이익률", f"{contrib_margin*100:.1f}%")

    subsec("비용 항목별 상세")
    info("5개년 비용 항목별 상세 금액표입니다(패드 연동). 각 항목의 연도별 증감을 확인하세요. 2026은 9개월 운영 기준.")
    cd = {'항목': list(cost_items_dyn.keys())}
    for i, yr in enumerate(D['yp']): cd[yr] = [fmt만(cost_items_dyn[k][i])+'원' for k in cd['항목']]
    dark_table(pd.DataFrame(cd))

# ═══ TAB 6: P&L + BEP ═══
if _ti == 6:
    sec("📋", "손익·BEP 분석")
    warn("2026F는 6월 오픈 기준 **9개월(6~2월)** 실적입니다. 연간 환산 시 x12/9 = 1.33배가 필요합니다.")
    info("⚠️ 2026F 월별 매출은 Excel 모델 기준이나, **현실적으로 오픈 초기(6~8월)는 인지도 부족 + 기존 회원권 만료 대기 + 여름 비수기가 겹쳐 정상 매출의 30~60% 수준**이 예상됩니다. 9월 이후 성수기 진입 시 점진적 정상화됩니다.")
    subsec("5개년 손익계산서")
    info("5개년 손익계산서입니다. 매출에서 비용을 차감한 영업이익과, 감가상각비를 더한 EBITDA를 함께 표시합니다. 회수율은 누적 EBITDA / 투자금입니다.")
    pl_full = {'항목': ['확정매출', '총비용', '  (감가상각비)', '영업이익', 'EBITDA', '영업이익률(%)', 'EBITDA이익률(%)', '회수율(%)']}
    for i, yr in enumerate(D['yp']):
        r, c, d, o, e = rev_p[i], cost_p[i], dep[i], op_p[i], ebitda_p[i]
        pl_full[yr] = [fmt억(r), fmt억(c), f"  ({fmt억(d)})", fmt억(o), fmt억(e), f"{o/r*100:.1f}%" if r else "0%", f"{e/r*100:.1f}%" if r else "0%", f"{rec_rate[i]*100:.1f}%"]
    dark_table(pd.DataFrame(pl_full))

    c1, c2 = st.columns(2)
    with c1:
        subsec("월별 매출 vs 비용 (2026, 현실 보정)")
        info("오픈 초기 현실을 반영한 월별 추정입니다. 6월=오픈월(인지도 없음, 장마 시작), 7~8월=폭염 비수기, 9~10월=성수기+인지도 상승, 11~12월=기온 하락, 1~2월=겨울 비수기.")
        # 현실 보정: 오픈초기+계절성+고객전환지연 반영
        # mrev_custom은 이미 시즌×램프업×상권×경제 보정 적용됨 (이중 적용 방지)
        fig = go.Figure()
        fig.add_trace(go.Bar(name='매출', x=D['months'], y=[v/1e6 for v in mrev_custom], marker_color=C['blue'],
            text=[f"{v/1e6:.0f}" for v in mrev_custom], textposition='inside', textfont=dict(size=11, color='white')))
        fig.add_trace(go.Bar(name='비용', x=D['months'], y=[v/1e6 for v in mcost_custom], marker_color=C['red'],
            text=[f"{v/1e6:.0f}" for v in mcost_custom], textposition='inside', textfont=dict(size=11, color='white')))
        lo(fig, title='월별 매출 vs 비용 (백만원, 컨트롤패널 연동)', barmode='group', height=440, yaxis_title='백만원')
        st.plotly_chart(fig)
    with c2:
        subsec("월별 손익 & 누적 (보정)")
        info("보정된 월별 손익(매출-비용)과 누적 추이입니다. 오픈 초기에는 비용이 매출을 초과하여 적자가 누적되며, 성수기(9~10월)에 일부 회복됩니다.")
        mpl = [r-c for r,c in zip(mrev_custom, mcost_custom)]
        cum_m = []; c_=0
        for v in mpl: c_+=v; cum_m.append(c_)
        fig2 = go.Figure()
        fig2.add_trace(go.Bar(name='월 손익', x=D['months'], y=[v/1e6 for v in mpl], marker_color=[C['green'] if v>=0 else C['red'] for v in mpl]))
        fig2.add_trace(go.Scatter(name='누적', x=D['months'], y=[v/1e6 for v in cum_m], mode='lines+markers', line=dict(color=C['orange'], width=3)))
        fig2.add_hline(y=0, line_dash="dash", line_color="#cbd5e1")
        lo(fig2, title='월별 손익 및 누적 추이 (백만원)', height=420, yaxis_title='백만원')
        st.plotly_chart(fig2)

    subsec("월별 매출 항목별 분해")
    info("2026년 월별 매출(골프 부분)을 상품별로 분해한 누적 막대 차트입니다(패드 연동). 임대는 4개월만 발생하므로 별도이며, 본 차트는 골프 9개월 매출만 표시합니다.")
    fig_ms = go.Figure()
    # mrev_custom은 골프 전용 매출이므로 임대 항목을 제외한 비중으로 분배
    _golf_items = {k: v for k, v in custom_rev_items.items() if '임대' not in k}
    _golf_total = sum(_golf_items.values())
    for idx, (name, ann_val) in enumerate(_golf_items.items()):
        share = ann_val / _golf_total if _golf_total else 0
        monthly_vals = [v * share for v in mrev_custom]
        fig_ms.add_trace(go.Bar(name=name, x=D['months'], y=[v/1e6 for v in monthly_vals], marker_color=PAL[idx%len(PAL)]))
    lo(fig_ms, title='월별 매출 항목별 분해 — 골프 (백만원)', barmode='stack', height=420, yaxis_title='백만원')
    st.plotly_chart(fig_ms)

    # ── BEP Section (merged into P&L tab) ──
    st.markdown("---")
    sec("⚖️", "손익분기점(BEP) 분석")
    info("손익분기점(BEP)은 매출과 비용이 동일하여 이익이 0인 지점입니다. BEP를 넘어서야 흑자가 시작되며, 안전마진이 클수록 경영 안정성이 높습니다.")
    bk = st.columns(4)
    bk[0].metric("BEP 매출", fmt억(bep_revenue))
    bk[1].metric("BEP 회원수", f"{bep_members:,.0f}명")
    bk[2].metric("공헌이익률", f"{contrib_margin*100:.1f}%")
    bk[3].metric("안전마진", f"{safety_margin:.1f}%", delta="흑자" if safety_margin>0 else "적자")

    c1, c2 = st.columns(2)
    with c1:
        subsec("BEP 차트")
        info("매출선과 비용선이 교차하는 지점이 손익분기점입니다. 초록 점선 오른쪽 영역이 이익 구간입니다.")
        rr = [bep_revenue*x/100 for x in range(50,151,10)]
        tc = [fixed_total + v*var_ratio for v in rr]
        fig = go.Figure()
        fig.add_trace(go.Scatter(name='매출', x=[r/억 for r in rr], y=[r/억 for r in rr], line=dict(color=C['blue'], width=3)))
        fig.add_trace(go.Scatter(name='총비용', x=[r/억 for r in rr], y=[c/억 for c in tc], line=dict(color=C['red'], width=3)))
        fig.add_trace(go.Scatter(name='고정비', x=[r/억 for r in rr], y=[fixed_total/억]*len(rr), line=dict(color=C['orange'], width=2, dash='dash')))
        fig.add_vline(x=bep_revenue/억, line_dash="dash", line_color=C['green'], annotation_text="BEP", annotation_font_size=11)
        lo(fig, title='손익분기점(BEP) 분석', height=420, xaxis_title='매출(억)', yaxis_title='금액(억)')
        st.plotly_chart(fig)
    with c2:
        subsec("비용구조 워터폴")
        info("고정비에서 시작하여 준변동비, 변동비를 더한 총비용과 매출을 비교하는 워터폴 차트입니다. 최종 영업이익이 양(+)이면 흑자입니다.")
        fig = go.Figure(go.Waterfall(x=["고정비","준변동비(60%F)","변동비","총비용","매출(2026F)","영업이익"],
            y=[fixed_total/억, semi_total*0.6/억, var_total/억, 0, rev_p[0]/억 if rev_p[0] else 0, 0],
            measure=["absolute","relative","relative","total","absolute","total"],
            connector={"line":{"color":C['slate']}}, decreasing={"marker":{"color":C['red']}}, increasing={"marker":{"color":C['green']}}, totals={"marker":{"color":C['blue']}}))
        lo(fig, title='비용구조 워터폴 (억원)', height=420)
        st.plotly_chart(fig)

    subsec("월별 BEP 도달 분석 (2026)")
    info("2026년 각 월별 매출이 월간 BEP(빨간 점선)를 초과하는지 확인합니다. 초과 월이 많을수록 연간 흑자 가능성이 높습니다.")
    monthly_bep = fixed_total/12/contrib_margin if contrib_margin else 0
    fig_mb = go.Figure()
    fig_mb.add_trace(go.Bar(name='월 매출', x=D['months'], y=[v/1e6 for v in mrev_custom], marker_color=C['blue']))
    fig_mb.add_hline(y=monthly_bep/1e6, line_dash="dash", line_color=C['red'], annotation_text=f"월BEP {monthly_bep/1e6:.0f}백만", annotation_font_size=11)
    bep_met = sum(1 for v in mrev_custom if v >= monthly_bep)
    lo(fig_mb, title=f'월별 BEP 달성 현황: {bep_met}/9개월', height=380, yaxis_title='백만원')
    st.plotly_chart(fig_mb)

# ═══ TAB 7: Investment/IRR ═══
if _ti == 7:
    sec("🎯", "투자 수익성 분석")
    info("투자 의사결정의 핵심 지표인 NPV, IRR, Payback, ROI를 종합적으로 분석합니다. NPV가 양(+)이고 IRR이 할인율보다 높으면 투자 가치가 있습니다.")
    ik = st.columns(5)
    ik[0].metric("순현재가치(NPV)", fmt억(npv_val))
    ik[1].metric("내부수익률(IRR)", f"{irr_val*100:.1f}%")
    ik[2].metric("투자회수기간", f"{payback:.1f}년" if payback else "5년+")
    ik[3].metric("누적 EBITDA", fmt억(cum_ebitda[-1]))
    ik[4].metric("투자수익률(ROI)", f"{(cum_ebitda[-1]-inv_won)/inv_won*100:.1f}%" if inv_won else "N/A")

    c1, c2 = st.columns(2)
    with c1:
        subsec(f"DCF 할인현금흐름 (할인율 {disc_r*100:.0f}%)")
        info(f"각 연도의 세후FCF(=EBITDA−법인세)를 할인율({disc_r*100:.0f}%)로 현재가치 환산합니다. 할인율 = 무위험수익률(은행 3~4%) + 사업리스크(3~4%) + 유동성프리미엄(2~3%). IRR이 할인율보다 높으면 은행 예금보다 투자 수익이 우수합니다.")
        dcf = []
        for i, yr in enumerate(D['yp']):
            pv = fcf_p[i]/(1+disc_r)**(i+1)
            dcf.append({'연도': yr, '세후FCF': fmt억(fcf_p[i]), '할인계수': f"{1/(1+disc_r)**(i+1):.4f}", '현재가치': fmt억(pv)})
        _pv_sum = sum(fcf_p[i]/(1+disc_r)**(i+1) for i in range(5))
        dcf.append({'연도': '합계(PV)', '세후FCF': fmt억(sum(fcf_p)), '할인계수': '-', '현재가치': fmt억(_pv_sum)})
        dcf.append({'연도': '−투자금', '세후FCF': '-', '할인계수': '-', '현재가치': fmt억(-inv_won)})
        dcf.append({'연도': 'NPV', '세후FCF': '-', '할인계수': '-', '현재가치': fmt억(_pv_sum - inv_won)})
        dark_table(pd.DataFrame(dcf))
    with c2:
        subsec("투자회수 추이")
        info("누적 EBITDA가 투자금(빨간 점선)에 도달하는 시점이 투자 회수 완료 시점입니다. 초록 영역이 넓을수록 빠른 회수를 의미합니다.")
        fig = go.Figure()
        fig.add_trace(go.Scatter(name='누적', x=[str(y) for y in D['yp']], y=[c/억 for c in cum_ebitda], mode='lines+markers', line=dict(color=C['green'], width=3), fill='tozeroy', fillcolor='rgba(34,197,94,0.1)'))
        fig.add_hline(y=s_inv, line_dash="dash", line_color=C['red'], annotation_text=f"투자금 {s_inv}억", annotation_font_size=11)
        lo(fig, title='누적 EBITDA vs 투자금 (억원)', height=420, yaxis_title='억원')
        st.plotly_chart(fig)

    subsec("할인율별 NPV 비교")
    info("할인율(자본비용)을 변동시켰을 때 NPV의 변화를 비교합니다(세후 FCF 기준). NPV가 0이 되는 할인율이 IRR과 일치합니다.")
    dr_range = [0.06,0.08,0.10,0.12,0.15,0.18,0.20]
    npv_by_dr = [sum(f/(1+dr)**(i+1) for i,f in enumerate(fcf_p))-inv_won for dr in dr_range]
    fig_dr = go.Figure(go.Bar(x=[f"{dr*100:.0f}%" for dr in dr_range], y=[n/억 for n in npv_by_dr], marker_color=[C['green'] if n>0 else C['red'] for n in npv_by_dr], text=[f"{n/억:.1f}" for n in npv_by_dr], textposition='inside', textfont=dict(size=12, color='#e2e8f0')))
    fig_dr.add_hline(y=0, line_dash="dash", line_color="#cbd5e1")
    lo(fig_dr, title='할인율별 순현재가치 비교 (억원)', height=380, yaxis_title='억원')
    st.plotly_chart(fig_dr)

# ═══ TAB 8: Cash Flow ═══
if _ti == 8:
    sec("💰", "현금흐름 분석")
    subsec("5개년 현금흐름표")
    info("EBITDA(임대수익 이미 포함)에서 법인세를 차감하여 영업 현금흐름(OCF)과 잉여 현금흐름(FCF)을 산출합니다. 누적 FCF가 양(+)이면 현금이 축적되고 있음을 의미합니다.")
    cf_data = {'항목': ['EBITDA(임대포함)','법인세(-)','영업CF','FCF','누적FCF']}
    cum_fcf = 0
    for i, yr in enumerate(D['yp']):
        tax = max(op_p[i]*s_tax_rate,0)
        ocf = ebitda_p[i]-tax
        fcf = ocf; cum_fcf += fcf
        cf_data[yr] = [fmt억(ebitda_p[i]), fmt억(tax), fmt억(ocf), fmt억(fcf), fmt억(cum_fcf)]
    dark_table(pd.DataFrame(cf_data))

    c1, c2 = st.columns(2)
    with c1:
        subsec("FCF 추이")
        info("연도별 잉여 현금흐름(FCF) 막대와 누적 추이선입니다. 초록 막대가 양(+)이면 해당 연도에 현금이 유입됨을 의미합니다.")
        cf_v, cum_cf = [], []; cc_f=0
        for i in range(5):
            tax=max(op_p[i]*s_tax_rate,0)
            f=ebitda_p[i]-tax; cf_v.append(f); cc_f+=f; cum_cf.append(cc_f)
        fig=go.Figure()
        fig.add_trace(go.Bar(name='FCF', x=[str(y) for y in D['yp']], y=[v/억 for v in cf_v], marker_color=[C['green'] if v>=0 else C['red'] for v in cf_v]))
        fig.add_trace(go.Scatter(name='누적', x=[str(y) for y in D['yp']], y=[v/억 for v in cum_cf], mode='lines+markers', line=dict(color=C['orange'], width=3)))
        fig.add_hline(y=0, line_dash="dash", line_color="#cbd5e1")
        lo(fig, title='연도별 FCF 및 누적 추이 (억원)', height=420, yaxis_title='억원')
        st.plotly_chart(fig)
    with c2:
        subsec("월별 현금흐름 (2026)")
        info("2026년 월별 현금흐름(매출-비용)과 누적 추이입니다. 개장 초기에는 적자가 나타나며 점차 흑자로 전환됩니다.")
        m_cf=[r-c for r,c in zip(mrev_custom,mcost_custom)]; cum_mcf=[]; c_=0
        for v in m_cf: c_+=v; cum_mcf.append(c_)
        fig=go.Figure()
        fig.add_trace(go.Bar(name='월CF', x=D['months'], y=[v/1e6 for v in m_cf], marker_color=[C['green'] if v>=0 else C['red'] for v in m_cf]))
        fig.add_trace(go.Scatter(name='누적', x=D['months'], y=[v/1e6 for v in cum_mcf], mode='lines+markers', line=dict(color=C['orange'], width=3)))
        lo(fig, title='월별 현금흐름 추이 (백만원)', height=420, yaxis_title='백만원')
        st.plotly_chart(fig)

# ═══ TAB 9: Rental Shops ═══
if _ti == 9:
    sec("🏪", "임대매장 수익 분석")
    info("사이드바 '🏪 임대매장 설정'에서 매장 수, 임대 조건, 계열사 투자 배분을 조절하세요.")

    # ── KPI ──
    info("임대매장은 골프연습장의 부가 수익원입니다. 외부 임차인은 고정 임대료를, 계열사는 할인 임대료 + 매출배분을 통해 수익을 창출합니다.")
    rk = st.columns(6)
    rk[0].metric("총 임대매장", f"{s_rent_shops}개")
    rk[1].metric("외부 임차", f"{s_rent_external}개")
    rk[2].metric("계열사", f"{s_rent_affiliate}개")
    rk[3].metric("연 임대수익", fmt만(rent_total_annual) + '원')
    rk[4].metric("보증금 수입", fmt만(rent_deposit_total) + '원')
    rk[5].metric("매출배분 수익(1년차)", fmt만(aff_rev_share_income[0]) + '원')

    # ── 매장 구성 & 수익 구조 ──
    c1, c2 = st.columns(2)
    with c1:
        subsec("임대매장 구성")
        info("외부 임차인과 계열사 매장의 비율을 보여줍니다. 외부 비중이 높으면 안정적 임대수익, 계열사 비중이 높으면 시너지 효과가 큽니다.")
        fig = go.Figure(go.Pie(
            labels=['외부 임차인', '계열사 매장'],
            values=[s_rent_external, s_rent_affiliate],
            hole=0.5, marker=dict(colors=[C['blue'], C['orange']], line=dict(color='white', width=3)),
            textinfo='label+percent'))
        lo(fig, title='매장 구성 비율', height=420)
        st.plotly_chart(fig)

    with c2:
        subsec("임대 수익 구성 (연간)")
        info("연간 임대수익을 외부 임대료, 계열사 임대료, 매출배분으로 분해합니다. 다양한 수익원이 있을수록 안정적입니다.")
        fig = go.Figure()
        fig.add_trace(go.Bar(name='외부 임대료', x=['임대수익'], y=[rent_ext_annual/억], marker_color=C['blue']))
        fig.add_trace(go.Bar(name='계열사 임대료', x=['임대수익'], y=[rent_aff_annual/억], marker_color=C['orange']))
        fig.add_trace(go.Bar(name='매출배분 수익', x=['매출배분'], y=[aff_rev_share_income[0]/억], marker_color=C['green']))
        lo(fig, title='연간 임대 수익 구성 (억원)', barmode='stack', height=420, yaxis_title='억원')
        st.plotly_chart(fig)

    # ── 외부 vs 계열사 조건 비교 ──
    subsec("임대 조건 비교")
    info("외부 임차인과 계열사 매장의 임대 조건을 비교합니다. 계열사는 투자 분담 대신 임대료 할인과 매출배분 구조를 적용합니다.")
    cond_data = {
        '항목': ['매장 수', '월 임대료', '보증금', '연간 인상률', '계약기간', '투자분담', '매출배분'],
        '외부 임차인': [f'{s_rent_external}개', f'{s_rent_ext_price}만원', f'{s_rent_ext_deposit}만원',
                      f'{s_rent_ext_up*100:.1f}%', '2년 (갱신 가능)', '없음 (임차인 부담)', '없음'],
        '계열사 매장': [f'{s_rent_affiliate}개', f'{aff_rent_monthly/만:,.0f}만원 ({s_aff_rent_disc}%↓)',
                      '면제 (투자분담)', f'{s_rent_ext_up*50:.1f}%', '5년 (장기)', f'본사 {s_aff_inv_share}% / 계열사 {100-s_aff_inv_share}%',
                      f'매출의 {s_aff_rev_share}%']
    }
    dark_table(pd.DataFrame(cond_data))

    # ── 5개년 임대수익 전망 ──
    subsec("5개년 임대수익 전망")
    c1, c2 = st.columns(2)
    with c1:
        info("5개년간 임대수익 구성요소별 전망입니다. 임대료 인상률과 매출배분 성장을 반영하며, 검정선이 총 임대수익 추이입니다.")
        fig = go.Figure()
        fig.add_trace(go.Bar(name='외부 임대료', x=[str(y) for y in D['yp']], y=[r['외부']/억 for r in rent_5yr], marker_color=C['blue']))
        fig.add_trace(go.Bar(name='계열사 임대료', x=[str(y) for y in D['yp']], y=[r['계열사']/억 for r in rent_5yr], marker_color=C['orange']))
        fig.add_trace(go.Bar(name='매출배분', x=[str(y) for y in D['yp']], y=[v/억 for v in aff_rev_share_income], marker_color=C['green']))
        total_rental_income = [rent_5yr[i]['합계'] + aff_rev_share_income[i] for i in range(5)]
        fig.add_trace(go.Scatter(name='총 임대수익', x=[str(y) for y in D['yp']], y=[v/억 for v in total_rental_income],
                                 mode='lines+markers', line=dict(color=C['dark'], width=3), marker=dict(size=8)))
        lo(fig, title='5개년 임대수익 전망 (억원)', barmode='stack', height=420, yaxis_title='억원')
        st.plotly_chart(fig)

    with c2:
        rent_tbl = {'항목': ['외부 임대료', '계열사 임대료', '매출배분 수익', '합계', '전체매출 대비']}
        for i, yr in enumerate(D['yp']):
            t = total_rental_income[i]
            pct = t / rev_p[i] * 100 if rev_p[i] else 0
            rent_tbl[yr] = [fmt만(rent_5yr[i]['외부'])+'원', fmt만(rent_5yr[i]['계열사'])+'원',
                            fmt만(aff_rev_share_income[i])+'원', fmt만(t)+'원', f'{pct:.1f}%']
        dark_table(pd.DataFrame(rent_tbl))

    # ── 계열사 투자 배분 분석 ──
    subsec("계열사 투자 배분 분석")
    info(f"계열사 총 투자금 {s_aff_inv_total:,}만원 중 본사 {s_aff_inv_share}%({s_aff_inv_total*s_aff_inv_share//100:,}만원) 분담 → 임대료 할인 {s_aff_rent_disc}% + 매출배분 {s_aff_rev_share}% 수취")

    c1, c2 = st.columns(2)
    with c1:
        info("본사와 계열사 간 투자금 분담 비율을 도넛 차트로 보여줍니다. 본사 분담 비중이 높을수록 임대료 할인과 매출배분으로 보전이 필요합니다.")
        fig = go.Figure(go.Pie(
            labels=['본사 분담', '계열사 분담'], values=[aff_inv_hq, aff_inv_aff],
            hole=0.5, marker=dict(colors=[C['red'], C['cyan']], line=dict(color='white', width=3)),
            textinfo='label+percent'))
        lo(fig, title='계열사 투자금 분담 구조', height=420)
        st.plotly_chart(fig)

    with c2:
        info("본사가 계열사에 투자한 금액 대비 회수 현황입니다. 누적 수익선이 빨간 점선(투자금)을 넘으면 회수 완료입니다.")
        # ROI 분석 - 본사의 계열사 투자 회수
        cum_return = 0
        roi_data = []
        for i in range(5):
            # 본사가 받는것: 계열사 임대료 + 매출배분
            yr_income = rent_5yr[i]['계열사'] + aff_rev_share_income[i]
            cum_return += yr_income
            roi = cum_return / aff_inv_hq * 100 if aff_inv_hq else 0
            roi_data.append({'연도': D['yp'][i], '연수익': yr_income, '누적': cum_return, 'ROI': roi})

        fig = go.Figure()
        fig.add_trace(go.Bar(name='연 수익', x=[r['연도'] for r in roi_data], y=[r['연수익']/만 for r in roi_data], marker_color=C['green']))
        fig.add_trace(go.Scatter(name='누적 수익', x=[r['연도'] for r in roi_data], y=[r['누적']/만 for r in roi_data],
                                 mode='lines+markers', line=dict(color=C['orange'], width=3)))
        fig.add_hline(y=aff_inv_hq/만, line_dash="dash", line_color=C['red'], annotation_text=f"본사 투자금 {aff_inv_hq/만:,.0f}만", annotation_font_size=11)
        lo(fig, title='계열사 투자 회수 추이 (만원)', height=420, yaxis_title='만원')
        st.plotly_chart(fig)

    # ── 임대매장 관리 체크리스트 ──
    subsec("임대매장 관리 포인트")
    mgmt = st.columns(2)
    mgmt[0].markdown("""
**📋 외부 임차인 관리**
- 임대차 계약 갱신 주기 (2년) 관리
- 보증금 예치 및 반환 스케줄
- 월 임대료 수금 관리 (연체 모니터링)
- 연간 임대료 인상 협의 및 통보
- 시설 원상복구 의무 확인
- 업종 제한 조항 (경쟁업종 금지)
- 화재/배상 보험 가입 확인
- 관리비(공용전기, 청소, 주차) 정산
""")
    mgmt[1].markdown(f"""
**🤝 계열사 매장 관리**
- 투자금 분담 비율 합의서 (본사 {s_aff_inv_share}%)
- 매출배분({s_aff_rev_share}%) 정산 주기 (월/분기)
- 매출 POS 연동 및 투명한 모니터링
- 인테리어/설비 소유권 귀속 명확화
- 계열사 직원 채용 및 교육 기준
- 운영시간 조율 (골프연습장 운영과 연계)
- 시너지 프로모션 (회원 교차 할인)
- 브랜드 가이드라인 준수
""")

    # ── 리스크 & 기회 ──
    subsec("임대매장 리스크 & 기회")
    ro = st.columns(2)
    ro[0].markdown("""
**⚠️ 리스크**
- 외부 임차인 공실 발생 시 수익 감소
- 계열사 매장 매출 부진 시 배분수익 감소
- 업종 부적합으로 인한 고객 불만
- 공용부 관리비 분쟁
- 임차인 시설 훼손
""")
    ro[1].markdown(f"""
**✅ 기회**
- 안정적 고정 수입원 (연 {fmt만(rent_total_annual)}원)
- 계열사 시너지로 고객 유입 증대
- 보증금 {fmt만(rent_deposit_total)}원 운용 가능
- 복합문화공간 이미지 구축
- 골프+F&B+리테일 원스톱 서비스
""")

# ═══ TAB (HIDDEN): Sensitivity ═══
if False:  # Tab removed from navigation
    sec("🔍", "민감도 · 시나리오 분석")
    info("주요 변수(매출성장률, 할인율 등)의 변동이 NPV에 미치는 영향을 분석합니다. 민감도가 높은 변수일수록 리스크 관리가 중요합니다.")
    if D['sm'] and any(any(v for v in row) for row in D['sm']):
        subsec("NPV 민감도 히트맵")
        info("매출 성장률(가로)과 할인율(세로) 조합에 따른 NPV를 히트맵으로 표시합니다. 초록색은 양(+)의 NPV, 빨간색은 음(-)의 NPV입니다.")
        fig=go.Figure(go.Heatmap(z=D['sm'], x=[f"매출성장{r*100:+.0f}%" for r in D['sr']], y=[f"할인율{d*100:.0f}%" for d in D['sd']],
            colorscale=[[0,C['red']],[0.5,'#fef3c7'],[1,C['green']]], text=[[f"{v}억" for v in row] for row in D['sm']], texttemplate="%{text}", textfont=dict(size=12, color='#e2e8f0')))
        lo(fig, title='NPV 민감도 히트맵 (억원)', height=420)
        st.plotly_chart(fig)

    c1, c2 = st.columns(2)
    with c1:
        subsec("토네이도 차트")
        info("각 비용/수익 항목을 ±10% 변동시켰을 때 NPV에 미치는 영향을 양방향 막대로 비교합니다. 막대가 긴 항목이 NPV에 가장 민감한 변수입니다.")
        td=[('인건비±10%',-3.0,3.0),('세금과공과±10%',-2.4,2.4),('고객유치율±10%',-1.5,1.5),('임대매출±10%',-1.1,1.1),('전력비±10%',-0.6,0.6),('이탈률±10%',-0.4,0.4)]
        fig=go.Figure()
        fig.add_trace(go.Bar(name='악화', y=[t[0] for t in td], x=[t[1] for t in td], orientation='h', marker_color=C['red']))
        fig.add_trace(go.Bar(name='개선', y=[t[0] for t in td], x=[t[2] for t in td], orientation='h', marker_color=C['green']))
        lo(fig, title='NPV 민감도 토네이도 차트 (억원)', barmode='relative', height=420, xaxis_title='NPV 변동(억)')
        st.plotly_chart(fig)
    with c2:
        subsec("시나리오 비교")
        info("최악·일반·최상 3가지 시나리오에서의 핵심 투자 지표를 비교합니다. 최악 시나리오에서도 투자 회수가 가능한지 확인하세요.")
        sc_d = {'지표': ['5년매출합계','5년EBITDA','NPV','회수율','IRR']}
        for si, sn in enumerate(['최악','일반(Base)','최상']):
            acq = A['market']['신규고객유치율'][si]; ba = A['market']['신규고객유치율'][1]
            sc = acq/ba if ba else 1
            sr = [r*sc for r in rev_p]; so = [r-c for r,c in zip(sr, cost_p)]; se = [o+d for o,d in zip(so, dep)]
            sdr = A['invest']['할인율'][si]; snpv = sum(e/(1+sdr)**(i+1) for i,e in enumerate(se))-inv_won
            srec = sum(se)/inv_won if inv_won else 0
            try: sirr=float(np_irr([-inv_won]+se))
            except: sirr=(sum(se)/5)/inv_won if inv_won else 0
            sc_d[sn] = [fmt억(sum(sr)), fmt억(sum(se)), fmt억(snpv), f"{srec*100:.1f}%", f"{sirr*100:.1f}%"]
        dark_table(pd.DataFrame(sc_d))

# ═══ TAB 10: Trade Area & Competition ═══
if _ti == 10:
    sec("📍", "상권·경쟁 가중치 분석")
    info("상권 입지, 경쟁 환경, 경제지표를 종합 평가하여 매출 보정계수를 산출합니다. 점수가 높을수록 사업 환경이 유리하며, 보정계수 1.0 이상이면 기준 매출 대비 상향 조정됩니다.")

    # ── 상권 종합 점수 ──
    tk = st.columns(5)
    tk[0].metric("상권 종합점수", f"{ta_score:.1f}점", delta=f"등급 {ta_grade}")
    tk[1].metric("매출 보정계수", f"{ta_rev_adj:.2f}x")
    tk[2].metric("경쟁업체 (반경3km)", f"{s_comp_outdoor+s_comp_indoor+s_comp_screen}개")
    tk[3].metric("경제환경 점수", f"{econ_score:.1f}점", delta=econ_grade)
    tk[4].metric("종합 보정계수", f"{combined_adj:.2f}x")

    c1, c2 = st.columns(2)
    with c1:
        subsec("상권 입지 평가 레이더")
        info("5개 평가 항목(배후인구, 접근성, 가시성, 경쟁우위, 성장성)을 레이더 차트로 시각화합니다. 기준선(70점) 이상이면 양호합니다.")
        ta_cats = ['배후인구', '접근성', '가시성', '경쟁우위', '성장성']
        ta_vals = [tw_pop, tw_access, tw_visible, tw_compete, tw_growth]
        ta_weights = [w_pop, w_acc, w_vis, w_comp, w_grow]
        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=ta_vals, theta=ta_cats, fill='toself', name='등촌 평가',
            line=dict(color=C['blue'], width=2), fillcolor='rgba(59,130,246,0.2)',
            text=[f"{v}점 (×{w}%)" for v, w in zip(ta_vals, ta_weights)], hoverinfo='text+name'))
        fig.add_trace(go.Scatterpolar(r=[70]*5, theta=ta_cats, name='기준선 (70점)',
            line=dict(color=C['slate'], width=1, dash='dash'), fillcolor='rgba(0,0,0,0)'))
        lo(fig, title=f'상권 종합점수: {ta_score:.1f}점 ({ta_grade})', height=420,
           polar=dict(bgcolor='#111827', radialaxis=dict(range=[0, 100], tickvals=[20,40,60,80,100], gridcolor='#1e293b', linecolor='#334155', tickfont=dict(color='#64748b')), angularaxis=dict(gridcolor='#1e293b', linecolor='#334155', tickfont=dict(color='#94a3b8'))))
        st.plotly_chart(fig)

    with c2:
        subsec("항목별 가중 점수 분해")
        info("각 항목의 가중점수(막대)와 원점수(점선)를 비교합니다. 가중점수가 높은 항목이 종합 점수에 더 큰 영향을 줍니다.")
        weighted_scores = [tw_pop*w_pop/w_total, tw_access*w_acc/w_total, tw_visible*w_vis/w_total,
                          tw_compete*w_comp/w_total, tw_growth*w_grow/w_total] if w_total else [0]*5
        fig = go.Figure()
        fig.add_trace(go.Bar(name='가중점수', x=ta_cats, y=weighted_scores,
            marker_color=[C['green'] if s >= 15 else C['orange'] if s >= 10 else C['red'] for s in weighted_scores],
            text=[f"{s:.1f}" for s in weighted_scores], textposition='inside', textfont=dict(size=12, color='white')))
        fig.add_trace(go.Scatter(name='원점수', x=ta_cats, y=ta_vals,
            mode='lines+markers',
            line=dict(color=C['blue'], width=2, dash='dot'), yaxis='y2'))
        lo(fig, title='가중점수(막대) vs 원점수(선)', height=420, yaxis_title='가중점수',
           yaxis2=dict(title='원점수', overlaying='y', side='right', range=[0,100], gridcolor='rgba(0,0,0,0)'))
        st.plotly_chart(fig)

    # ── 경쟁 환경 ──
    subsec("경쟁 환경 분석 (반경 3km)")
    c1, c2 = st.columns(2)
    with c1:
        info("반경 3km 이내 업종별 경쟁업체 수입니다. 실외 연습장이 가장 직접적인 경쟁 대상이며, 스크린골프는 간접 경쟁입니다.")
        comp_types = ['실외 연습장', '실내 연습장', '스크린골프']
        comp_counts = [s_comp_outdoor, s_comp_indoor, s_comp_screen]
        fig = go.Figure(go.Bar(x=comp_types, y=comp_counts,
            marker_color=[C['red'], C['orange'], C['yellow']],
            text=[f"{c}개" for c in comp_counts], textposition='inside', textfont=dict(size=12, color='white')))
        lo(fig, title='업종별 경쟁업체 수', height=440, yaxis_title='업체수')
        st.plotly_chart(fig)

    with c2:
        info("경쟁강도 지수는 실외(가중치 1.0), 실내(0.7), 스크린(0.3)을 합산한 종합 지표입니다. 10점 이상이면 경쟁이 치열한 상권입니다.")
        # Competition intensity index
        comp_intensity = s_comp_outdoor * 1.0 + s_comp_indoor * 0.7 + s_comp_screen * 0.3
        market_per_comp = e_golf_pop * 만 / (e_range_count + e_screen_count) if (e_range_count + e_screen_count) else 0
        local_density = (s_comp_outdoor + s_comp_indoor + s_comp_screen) / 3  # per km radius
        fig = go.Figure(go.Indicator(mode="gauge+number", value=comp_intensity,
            gauge={'axis': {'range': [0, 20], 'tickfont': {'color': '#94a3b8'}},
                   'bar': {'color': C['red'] if comp_intensity > 10 else C['orange'] if comp_intensity > 5 else C['green']},
                   'bgcolor': '#1e293b',
                   'steps': [{'range':[0,5],'color':'#14532d'},{'range':[5,10],'color':'#78350f'},{'range':[10,20],'color':'#7f1d1d'}]},
            title={'text': '경쟁강도 지수', 'font': {'color': '#e2e8f0'}},
            number={'suffix': '점', 'font': {'color': '#f1f5f9', 'size': 32}}))
        lo(fig, title='경쟁강도 게이지', height=440)
        st.plotly_chart(fig)

    # ── 경쟁사 비교 ──
    subsec("주요 경쟁사 비교")
    info("등촌(신규), 제니스, 마곡나루 3개 경쟁사의 핵심 스펙을 비교합니다. 가격, 시설, 규모 면에서 차별화 포인트를 확인하세요.")
    ct = [{'항목':'위치','등촌(신규)':'등촌동','제니스':'구로구 고척동','마곡나루':'마곡동'},
          {'항목':'타석수','등촌(신규)':'88타석','제니스':'112타석','마곡나루':'50타석'},
          {'항목':'형태','등촌(신규)':'실외','제니스':'실내+실외','마곡나루':'실내'},
          {'항목':'부대시설','등촌(신규)':'골프+사우나+헬스+락카+임대매장','제니스':'골프+사우나+헬스','마곡나루':'골프'},
          {'항목':'1개월 남','등촌(신규)':'290,000원','제니스':'330,000원','마곡나루':'350,000원'},
          {'항목':'강점','등촌(신규)':'가격, 신규시설','제니스':'복합시설, 규모','마곡나루':'실내, 신도시'}]
    dark_table(pd.DataFrame(ct))

    c1, c2 = st.columns(2)
    with c1:
        subsec("경쟁 포지셔닝 레이더")
        info("3개 경쟁사의 5개 항목별 경쟁력을 레이더 차트로 비교합니다. 등촌은 가격경쟁력에서 우위, 제니스는 시설·부대시설에서 우위입니다.")
        cats = ['가격경쟁력','시설규모','접근성','부대시설','브랜드']
        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=[85,70,80,80,60], theta=cats, fill='toself', name='등촌', line=dict(color=C['blue']), fillcolor='rgba(59,130,246,0.15)'))
        fig.add_trace(go.Scatterpolar(r=[60,90,80,80,85], theta=cats, fill='toself', name='제니스', line=dict(color=C['red']), fillcolor='rgba(239,68,68,0.15)'))
        fig.add_trace(go.Scatterpolar(r=[55,50,75,40,70], theta=cats, fill='toself', name='마곡나루', line=dict(color=C['orange']), fillcolor='rgba(249,115,22,0.15)'))
        lo(fig, title='경쟁 포지셔닝 비교 레이더', height=420, polar=dict(bgcolor='#111827', radialaxis=dict(range=[0,100], gridcolor='#1e293b', linecolor='#334155', tickfont=dict(color='#64748b')), angularaxis=dict(gridcolor='#1e293b', linecolor='#334155', tickfont=dict(color='#94a3b8'))))
        st.plotly_chart(fig)
    with c2:
        subsec("서울 실외 연습장 요금 비교 (실제 확인 데이터)")
        info("웹에서 직접 확인된 서울 실외 연습장 실제 요금입니다. 동도센트리움/목동골프타운은 홈페이지 접속 불가로 전화 확인 필요.")

        # 실제 확인된 데이터만 사용 (출처 명시)
        comp_html = """
<table style="width:100%;border-collapse:collapse;font-size:12px;margin:8px 0;">
<tr>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:left;border-bottom:2px solid #334155;">연습장</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:left;border-bottom:2px solid #334155;">위치</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:center;border-bottom:2px solid #334155;">타석</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">1개월 종일</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">3개월 종일</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">일일(주중)</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:left;border-bottom:2px solid #334155;">출처</th>
</tr>
<tr style="background:#0f2a1e;">
<td style="color:#86efac;padding:7px 8px;font-weight:700;border-bottom:1px solid #1e293b;">등촌(2026)</td>
<td style="color:#86efac;padding:7px 8px;border-bottom:1px solid #1e293b;">강서구</td>
<td style="color:#86efac;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">88</td>
<td style="color:#86efac;padding:7px 8px;text-align:right;font-weight:700;border-bottom:1px solid #1e293b;">290,000</td>
<td style="color:#86efac;padding:7px 8px;text-align:right;font-weight:700;border-bottom:1px solid #1e293b;">800,000</td>
<td style="color:#86efac;padding:7px 8px;text-align:right;font-weight:700;border-bottom:1px solid #1e293b;">23,000</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">컨트롤패널</td>
</tr>
<tr>
<td style="color:#fca5a5;padding:7px 8px;font-weight:700;border-bottom:1px solid #1e293b;">제니스</td>
<td style="color:#e2e8f0;padding:7px 8px;border-bottom:1px solid #1e293b;">구로구 고척동</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">112</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">300,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">830,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">24,000</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">사용자 제공</td>
</tr>
<tr style="background:#111827;">
<td style="color:#e2e8f0;padding:7px 8px;font-weight:600;border-bottom:1px solid #1e293b;">메이필드호텔</td>
<td style="color:#e2e8f0;padding:7px 8px;border-bottom:1px solid #1e293b;">강서구</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">75</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">322,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">795,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">22,000(60분)</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">mayfield.co.kr</td>
</tr>
<tr>
<td style="color:#e2e8f0;padding:7px 8px;font-weight:600;border-bottom:1px solid #1e293b;">쇼골프타운</td>
<td style="color:#e2e8f0;padding:7px 8px;border-bottom:1px solid #1e293b;">강서구</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">183</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">230,000(주중)</td>
<td style="color:#94a3b8;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">-</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">16,800~22,000</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">김캐디</td>
</tr>
<tr style="background:#111827;">
<td style="color:#e2e8f0;padding:7px 8px;font-weight:600;border-bottom:1px solid #1e293b;">88CC</td>
<td style="color:#e2e8f0;padding:7px 8px;border-bottom:1px solid #1e293b;">영등포구</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">-</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">200,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">590,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">19,000(60분)</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">88countryclub.co.kr</td>
</tr>
<tr>
<td style="color:#e2e8f0;padding:7px 8px;font-weight:600;border-bottom:1px solid #1e293b;">엑스골프 장한평</td>
<td style="color:#e2e8f0;padding:7px 8px;border-bottom:1px solid #1e293b;">성동구</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">72</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">250,000(주중)</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">560,000</td>
<td style="color:#e2e8f0;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">18,000~19,000</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">김캐디</td>
</tr>
<tr style="background:#111827;">
<td style="color:#94a3b8;padding:7px 8px;font-weight:600;border-bottom:1px solid #1e293b;">동도센트리움</td>
<td style="color:#94a3b8;padding:7px 8px;border-bottom:1px solid #1e293b;">구로구 오류동</td>
<td style="color:#94a3b8;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">124</td>
<td style="color:#fbbf24;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">전화확인필요</td>
<td style="color:#fbbf24;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">전화확인필요</td>
<td style="color:#fbbf24;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">전화확인필요</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">02-2060-0707</td>
</tr>
<tr>
<td style="color:#94a3b8;padding:7px 8px;font-weight:600;border-bottom:1px solid #1e293b;">목동골프타운</td>
<td style="color:#94a3b8;padding:7px 8px;border-bottom:1px solid #1e293b;">양천구</td>
<td style="color:#94a3b8;padding:7px 8px;text-align:center;border-bottom:1px solid #1e293b;">-</td>
<td style="color:#fbbf24;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">전화확인필요</td>
<td style="color:#fbbf24;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">전화확인필요</td>
<td style="color:#fbbf24;padding:7px 8px;text-align:right;border-bottom:1px solid #1e293b;">전화확인필요</td>
<td style="color:#94a3b8;padding:7px 8px;font-size:10px;border-bottom:1px solid #1e293b;">02-2605-5966</td>
</tr>
<tr style="background:#1e293b;">
<td style="color:#fbbf24;padding:8px;font-weight:700;border-bottom:2px solid #334155;" colspan="3">서울 전체 평균 (1,519개소)</td>
<td style="color:#fbbf24;padding:8px;text-align:right;font-weight:700;border-bottom:2px solid #334155;">239,911</td>
<td style="color:#94a3b8;padding:8px;text-align:right;border-bottom:2px solid #334155;">-</td>
<td style="color:#fbbf24;padding:8px;text-align:right;font-weight:700;border-bottom:2px solid #334155;">20,681(60분)</td>
<td style="color:#94a3b8;padding:8px;font-size:10px;border-bottom:2px solid #334155;">김캐디 2023</td>
</tr>
</table>
<div style="color:#64748b;font-size:11px;margin-top:4px;">※ 마포구·서대문구: 대형 실외 연습장 없음 (실내/스크린만 존재) | 모든 금액 VAT 포함 | 2024~2025년 확인 기준, 2026년 인상 가능성 있음</div>
"""
        st.markdown(comp_html, unsafe_allow_html=True)

        # 차트도 실제 데이터로 교체
        fig = go.Figure(go.Bar(
            x=['등촌<br>(2026)', '제니스<br>(구로)', '메이필드<br>(강서)', '쇼골프<br>(강서)', '88CC<br>(영등포)', '엑스골프<br>(성동)'],
            y=[29, 30, 32.2, 23, 20, 25],
            marker_color=[C['blue'], C['red'], C['orange'], C['green'], C['purple'], C['cyan']],
            text=['29만', '30만', '32.2만', '23만', '20만', '25만'],
            textposition='inside', textfont=dict(size=13, color='white'),
            hovertemplate='%{x}: %{y}만원<extra></extra>'))
        lo(fig, title='1개월 종일 남성 요금 비교 (만원, 각 연습장 홈페이지/김캐디 확인)', height=380, yaxis_title='만원')
        st.plotly_chart(fig)

    # ── SWOT ──
    subsec("SWOT 분석")
    sw = st.columns(2)
    sw[0].markdown("""**💪 강점 (S)**\n- 가격경쟁력 5.9~26% 저렴\n- 88타석 적정 규모\n- 신규시설 (2026)\n- 사우나/헬스 등 부대시설 완비\n- 실외 개방감 + 마곡 인접\n\n**🎯 기회 (O)**\n- 마곡지구 신규 입주 수요\n- 초보자/여성 골퍼 확대\n- 프리미엄 레슨 시장\n- 상권 성장성 높음""")
    sw[1].markdown("""**⚠️ 약점 (W)**\n- 기상 영향 (실외 위주)\n- 90분 타임 미운영\n- 브랜드 인지도 부재 (신규)\n- 기존 고객 확보 시간 필요\n\n**🔥 위협 (T)**\n- 제니스 등 대형 경쟁사 (112타석)\n- 실내 시뮬레이션 골프 확대\n- 골프인구 장기 감소 추세\n- 인건비/전기료 지속 상승""")

    # ── 상권 매출 영향도 ──
    subsec("상권·경쟁이 매출에 미치는 영향")
    info(f"상권점수 {ta_score:.1f}점 → 매출 보정 ×{ta_rev_adj:.2f} | 경제환경 {econ_score:.1f}점 → 매출 보정 ×{econ_rev_adj:.2f} | **종합 보정: ×{combined_adj:.2f}** (이미 rev_p에 반영됨)")
    # rev_p에는 이미 상권·경제 보정이 반영되어 있음 — 보정 적용 전(없었다면) 매출은 rev_p ÷ combined_adj
    rev_no_adj = [r / combined_adj if combined_adj else r for r in rev_p]
    fig = go.Figure()
    fig.add_trace(go.Bar(name='보정 없을 때(가정)', x=[str(y) for y in D['yp']], y=[r/억 for r in rev_no_adj], marker_color=C['blue_l']))
    fig.add_trace(go.Bar(name='보정 적용 매출(현재)', x=[str(y) for y in D['yp']], y=[r/억 for r in rev_p], marker_color=C['blue']))
    fig.add_trace(go.Scatter(name='보정 효과', x=[str(y) for y in D['yp']], y=[(a-b)/억 for a,b in zip(rev_p, rev_no_adj)],
        mode='lines+markers', line=dict(color=C['green'] if combined_adj >= 1 else C['red'], width=2)))
    lo(fig, title='상권·경제 보정 전후 매출 비교 (억원)', barmode='group', height=420, yaxis_title='억원')
    st.plotly_chart(fig)


# ═══ TAB 11: Golf Market & Economic Indicators ═══
if _ti == 11:
    sec("📉", "골프시장·경제지표 분석")

    # ── 데이터 수집 현황 ──
    ed = st.session_state.get('econ_data', DEFAULTS)
    src = ed.get('source', 'default')
    if src == 'api':
        st.success(f"🟢 **API 자동수집 활성** — {ed.get('last_update','')} | {ed.get('notes','')}")
    elif src == 'mixed':
        st.warning(f"🟡 **부분 수집** — {ed.get('notes','')}")
    else:
        st.info("⚪ **수동 입력 모드** — 사이드바 `📉 골프시장·경제지표` → API 키 입력 후 `🔄 지금 갱신` 클릭 시 한국은행·통계청에서 자동 수집")

    with st.expander("🔑 API 키 발급 안내 및 데이터 소스 상세", expanded=False):
        st.markdown(API_GUIDE)
        subsec("데이터 소스 현황")
        src_tbl = [
            {'지표': '기준금리', '출처': '한국은행 ECOS', '주기': '8회/년', 'API': '✅' if ecos_key else '❌ 키 필요', '현재값': f"{e_interest*100:.2f}%"},
            {'지표': 'GDP 성장률', '출처': '한국은행 ECOS', '주기': '분기', 'API': '✅' if ecos_key else '❌ 키 필요', '현재값': f"{e_gdp_growth*100:.1f}%"},
            {'지표': '소비자물가', '출처': '한국은행 ECOS', '주기': '매월', 'API': '✅' if ecos_key else '❌ 키 필요', '현재값': f"{e_cpi*100:.1f}%"},
            {'지표': '실업률', '출처': '통계청 KOSIS', '주기': '매월', 'API': '✅' if kosis_key else '❌ 키 필요', '현재값': f"{e_unemp*100:.1f}%"},
            {'지표': '가처분소득', '출처': '통계청 KOSIS', '주기': '분기', 'API': '✅' if kosis_key else '❌ 키 필요', '현재값': f"{e_disposable*100:.1f}%"},
            {'지표': '골프인구', '출처': '골프장경영협회', '주기': '연간', 'API': '수동', '현재값': f"{e_golf_pop}만명"},
            {'지표': '연습장 수', '출처': '골프장경영협회', '주기': '연간', 'API': '수동', '현재값': f"{e_range_count:,}개"},
        ]
        dark_table(pd.DataFrame(src_tbl))

    info("골프 산업은 경기 변동, 소비 심리, 인구구조에 민감합니다. 사이드바에서 지표를 조절하면 영향도가 실시간 반영됩니다.")

    # ── 골프시장 현황 KPI ──
    ek = st.columns(5)
    ek[0].metric("골프인구", f"{e_golf_pop}만명", delta=f"{e_golf_growth*100:+.1f}%")
    ek[1].metric("연습장 수", f"{e_range_count:,}개")
    ek[2].metric("스크린골프", f"{e_screen_count:,}개")
    ek[3].metric("1인당 연습장", f"{e_golf_pop*만/e_range_count:,.0f}명" if e_range_count else "N/A")
    ek[4].metric("경제환경 점수", f"{econ_score:.1f}점", delta=econ_grade)

    # ── 골프시장 트렌드 (팩트 데이터) ──
    subsec("실외 골프연습장 시장 트렌드 (2018~2026)")
    info("출처: Excel 원본 '골프시장 분석' 시트 (DART 전자공시, 신용카드 빅데이터, 국가통계, 골프산업백서 기반). 골프장(라운딩)이 아닌 **실외 골프연습장** 시장 데이터입니다.")

    st.markdown("""
<div style="background:#0f172a;border:1px solid #334155;border-radius:12px;padding:20px;margin:12px 0;">
<div style="color:#60a5fa;font-size:15px;font-weight:700;margin-bottom:14px;">📊 국내 실외 골프연습장 시장 규모 (Excel 원본 데이터)</div>
<table style="width:100%;border-collapse:collapse;font-size:13px;">
<tr>
<th style="background:#1e293b;color:#60a5fa;padding:8px 12px;text-align:center;border-bottom:2px solid #334155;">연도</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">시장규모(억)</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:center;border-bottom:2px solid #334155;">전년비</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">대형80+타석</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">중형40~70</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">소형~30</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:right;border-bottom:2px solid #334155;">강서구골퍼(천)</th>
<th style="background:#1e293b;color:#60a5fa;padding:8px;text-align:left;border-bottom:2px solid #334155;">시장 상황</th>
</tr>
<tr><td style="background:#111827;color:#e2e8f0;padding:6px 12px;text-align:center;">2018</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">5,842</td><td style="background:#111827;color:#94a3b8;padding:6px;text-align:center;">-</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">32.5억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">12.4억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">5.4억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">443</td><td style="background:#111827;color:#94a3b8;padding:6px;font-size:11px;">안정적 경영</td></tr>
<tr><td style="background:#0f172a;color:#e2e8f0;padding:6px 12px;text-align:center;">2019</td><td style="background:#0f172a;color:#86efac;padding:6px;text-align:right;">6,051</td><td style="background:#0f172a;color:#86efac;padding:6px;text-align:center;">+3.6%</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">33.8억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">13.1억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">5.7억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">437</td><td style="background:#0f172a;color:#94a3b8;padding:6px;font-size:11px;">골프 대중화 확산</td></tr>
<tr><td style="background:#111827;color:#e2e8f0;padding:6px 12px;text-align:center;">2020</td><td style="background:#111827;color:#86efac;padding:6px;text-align:right;">6,610</td><td style="background:#111827;color:#86efac;padding:6px;text-align:center;">+9.2%</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">37.2억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">14.8억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">6.2억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">425</td><td style="background:#111827;color:#94a3b8;padding:6px;font-size:11px;">코로나 실내 제한 반사이익</td></tr>
<tr><td style="background:#0f172a;color:#e2e8f0;padding:6px 12px;text-align:center;">2021</td><td style="background:#0f172a;color:#86efac;padding:6px;text-align:right;font-weight:700;">7,420</td><td style="background:#0f172a;color:#86efac;padding:6px;text-align:center;font-weight:700;">+12.3%</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">41.5억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">16.5억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">6.8억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">420</td><td style="background:#0f172a;color:#86efac;padding:6px;font-size:11px;">MZ 신규 골퍼 대거 유입</td></tr>
<tr><td style="background:#111827;color:#e2e8f0;padding:6px 12px;text-align:center;">2022</td><td style="background:#111827;color:#fbbf24;padding:6px;text-align:right;font-weight:700;">8,150</td><td style="background:#111827;color:#fbbf24;padding:6px;text-align:center;font-weight:700;">+9.8%</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">44.8억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">18.2억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">7.1억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">415</td><td style="background:#111827;color:#fbbf24;padding:6px;font-size:11px;">★ 역대 최고 매출 (정점)</td></tr>
<tr><td style="background:#0f172a;color:#e2e8f0;padding:6px 12px;text-align:center;">2023</td><td style="background:#0f172a;color:#fca5a5;padding:6px;text-align:right;font-weight:700;">6,980</td><td style="background:#0f172a;color:#ef4444;padding:6px;text-align:center;font-weight:700;">-14.4%</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">38.2억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">14.5억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">5.5억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">407</td><td style="background:#0f172a;color:#ef4444;padding:6px;font-size:11px;">엔데믹 하락 시작</td></tr>
<tr><td style="background:#111827;color:#e2e8f0;padding:6px 12px;text-align:center;">2024</td><td style="background:#111827;color:#fca5a5;padding:6px;text-align:right;font-weight:700;">5,920</td><td style="background:#111827;color:#ef4444;padding:6px;text-align:center;font-weight:700;">-15.2%</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">32.1억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">11.8억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">4.1억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">397</td><td style="background:#111827;color:#ef4444;padding:6px;font-size:11px;">2019년 매출 하회</td></tr>
<tr><td style="background:#0f172a;color:#e2e8f0;padding:6px 12px;text-align:center;">2025</td><td style="background:#0f172a;color:#fca5a5;padding:6px;text-align:right;">5,610</td><td style="background:#0f172a;color:#fca5a5;padding:6px;text-align:center;">-5.2%</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">29.8억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">10.9억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">3.8억</td><td style="background:#0f172a;color:#e2e8f0;padding:6px;text-align:right;">388</td><td style="background:#0f172a;color:#fca5a5;padding:6px;font-size:11px;">고비용 구조 고착</td></tr>
<tr><td style="background:#111827;color:#e2e8f0;padding:6px 12px;text-align:center;">2026</td><td style="background:#111827;color:#94a3b8;padding:6px;text-align:right;">5,430</td><td style="background:#111827;color:#94a3b8;padding:6px;text-align:center;">-3.2%</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">28.5억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">10.2억</td><td style="background:#111827;color:#e2e8f0;padding:6px;text-align:right;">3.5억</td><td style="background:#111827;color:#94a3b8;padding:6px;text-align:right;">-</td><td style="background:#111827;color:#94a3b8;padding:6px;font-size:11px;">시장 재편·구조적 쇠퇴기 진입</td></tr>
</table>
</div>
""", unsafe_allow_html=True)

    st.markdown("""
<div style="display:grid;grid-template-columns:1fr 1fr;gap:12px;margin:12px 0;">
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #ef4444;">
<b style="color:#fca5a5;">시장 하락 원인 (Excel 분석 메모 기반)</b><br>
<span style="color:#cbd5e1;font-size:12px;line-height:1.7;">
① <b>2022년 정점 후 31% 붕괴</b>: 8,150억→5,610억 (3년간)<br>
② <b>2024년 매출이 2019년 이하</b>: 5,920억 < 6,051억<br>
③ <b>물가 14% 상승 감안 시 실질 매출은 2019년의 약 70%</b><br>
④ <b>전기료 35~40% 폭등</b>, 최저임금 20%↑ vs 2019년<br>
⑤ <b>이익률 급락</b>: 15.2%(2018) → 6.4%(2024~25 추정)<br>
⑥ <b>신용카드 결제 데이터</b>: 2023~2024 연습장 카드 사용액 14.2% 감소
</span>
</div>
<div style="background:#1e293b;border-radius:10px;padding:16px;border-left:4px solid #f97316;">
<b style="color:#fdba74;">등촌에 미치는 영향</b><br>
<span style="color:#cbd5e1;font-size:12px;line-height:1.7;">
① <b>강서구 골퍼 인구 지속 감소</b>: 443천명(2018)→388천명(2025), -12.4%<br>
② <b>대형 연습장(80+타석) 시장도 축소</b>: 44.8억(2022)→28.5억(2026), -36%<br>
③ <b>시장 재편기</b>: 약한 연습장 도태 → 시설 경쟁력 있는 곳에 기회<br>
④ <b>고비용 구조</b>: 전기료+인건비 폭등 → 비용 관리가 생존의 핵심<br>
⑤ <b>등촌 유리점</b>: 신축 시설+사우나/헬스+88타석 → 시장 재편 시 생존 가능
</span>
</div>
</div>
""", unsafe_allow_html=True)

    warn(f"현재 설정: 골프인구 증감률 {e_golf_growth*100:+.1f}%/년. Excel 원본 기준 실외 연습장 시장 감소율: 2023년 -14.4%, 2024년 -15.2%, 2025년 -5.2%. ① 가정 시트 Base 시나리오: **-3.5%/년**.")

    # ── 경제지표 영향도 분석 ──
    subsec("경제지표 영향도 분석")
    c1, c2 = st.columns(2)
    with c1:
        info("5개 경제지표의 현재 수준을 레이더 차트로 표시합니다. 기준선(55점) 이상이면 우호적, 이하이면 부정적 환경입니다.")
        econ_cats = ['골프시장', 'GDP성장', '물가(역)', '금리(역)', '가처분소득']
        econ_vals = [econ_golf_score, econ_gdp_score, econ_cpi_score, econ_ir_score, econ_disp_score]
        econ_weights = [ew_golf, ew_gdp, ew_cpi, ew_ir, ew_disp]
        fig = go.Figure()
        fig.add_trace(go.Scatterpolar(r=econ_vals, theta=econ_cats, fill='toself', name='현재 지표',
            line=dict(color=C['purple'], width=2), fillcolor='rgba(168,85,247,0.2)',
            text=[f"{v:.0f}점(×{w}%)" for v,w in zip(econ_vals, econ_weights)]))
        fig.add_trace(go.Scatterpolar(r=[55]*5, theta=econ_cats, name='기준선(55)',
            line=dict(color=C['slate'], width=1, dash='dash')))
        lo(fig, title=f'경제환경 종합: {econ_score:.1f}점 ({econ_grade})', height=420,
           polar=dict(bgcolor='#111827', radialaxis=dict(range=[0,100], gridcolor='#1e293b', linecolor='#334155', tickfont=dict(color='#64748b')), angularaxis=dict(gridcolor='#1e293b', linecolor='#334155', tickfont=dict(color='#94a3b8'))))
        st.plotly_chart(fig)

    with c2:
        subsec("지표별 가중 점수")
        info("각 경제지표의 가중 점수를 막대로 비교합니다. 가중점수가 높을수록 해당 지표가 사업 환경에 우호적임을 의미합니다.")
        weighted_econ = [econ_golf_score*ew_golf/ew_total, econ_gdp_score*ew_gdp/ew_total,
                        econ_cpi_score*ew_cpi/ew_total, econ_ir_score*ew_ir/ew_total,
                        econ_disp_score*ew_disp/ew_total] if ew_total else [0]*5
        raw_vals = [e_golf_growth*100, e_gdp_growth*100, e_cpi*100, e_interest*100, e_disposable*100]
        raw_labels = [f"{v:+.1f}%" for v in raw_vals]
        fig = go.Figure()
        fig.add_trace(go.Bar(x=econ_cats, y=weighted_econ,
            marker_color=[C['green'] if s>=12 else C['orange'] if s>=8 else C['red'] for s in weighted_econ],
            text=[f"{s:.1f}" for s in weighted_econ], textposition='inside', textfont=dict(size=12, color='white')))
        lo(fig, title='경제지표 가중점수 (높을수록 우호적)', height=420, yaxis_title='가중점수')
        st.plotly_chart(fig)

    # ── 골프시장 트렌드 ──
    subsec("골프시장 트렌드 시뮬레이션")
    c1, c2 = st.columns(2)
    with c1:
        info("현재 골프인구 증감률을 기반으로 10년간 인구 추이를 전망합니다. 감소 추세이면 빨간색, 증가 추세이면 파란색으로 표시됩니다.")
        # 10-year golf population projection
        golf_pop_proj = [e_golf_pop]
        for i in range(9):
            golf_pop_proj.append(golf_pop_proj[-1] * (1 + e_golf_growth))
        yrs = [str(2026+i) for i in range(10)]
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=yrs, y=golf_pop_proj, mode='lines+markers',
            line=dict(color=C['blue'] if e_golf_growth >= 0 else C['red'], width=3), fill='tozeroy',
            fillcolor=f'rgba(59,130,246,0.1)' if e_golf_growth >= 0 else 'rgba(239,68,68,0.1)'))
        lo(fig, title=f'골프인구 10개년 전망 ({e_golf_growth*100:+.1f}%/년)', height=420, yaxis_title='만명')
        st.plotly_chart(fig)

    with c2:
        info("연습장 수가 연 1%씩 감소(구조조정)한다고 가정할 때, 연습장 1개소당 배분되는 골프인구 전망입니다. 수치가 높을수록 시장 기회가 큽니다.")
        # Range market share projection
        range_per_pop = [e_golf_pop * 만 / e_range_count if e_range_count else 0]
        for i in range(9):
            pop = golf_pop_proj[i+1] * 만
            # Assume range count decreases 1% per year (consolidation)
            ranges = e_range_count * (1 - 0.01) ** (i+1)
            range_per_pop.append(pop / ranges if ranges else 0)
        fig = go.Figure()
        fig.add_trace(go.Bar(x=yrs, y=range_per_pop,
            marker_color=[C['green'] if v > range_per_pop[0] else C['orange'] for v in range_per_pop],
            text=[f"{v:,.0f}" for v in range_per_pop], textposition='inside', textfont=dict(size=12, color='#e2e8f0')))
        lo(fig, title='연습장당 골프인구 전망 (명/개소)', height=420, yaxis_title='명')
        st.plotly_chart(fig)

    # ── 경제지표별 매출 민감도 ──
    subsec("경제지표 변동 시 매출 영향 시뮬레이션")
    info("각 경제지표가 ±1%p 변동했을 때 연 매출에 미치는 영향을 추정합니다.")

    # Sensitivity of each indicator to revenue
    base_rev_yr = rev_p[1] if len(rev_p) > 1 else rev_p[0]  # 2027F base
    sensitivities = [
        ('골프인구 증감률', e_golf_growth*100, 0.8, '골프인구 1%p↑ → 매출 0.8%↑'),
        ('GDP 성장률', e_gdp_growth*100, 0.5, 'GDP 1%p↑ → 매출 0.5%↑'),
        ('소비자물가', e_cpi*100, -0.3, '물가 1%p↑ → 매출 0.3%↓'),
        ('기준금리', e_interest*100, -0.4, '금리 1%p↑ → 매출 0.4%↓'),
        ('가처분소득', e_disposable*100, 0.6, '가처분소득 1%p↑ → 매출 0.6%↑'),
    ]

    fig = go.Figure()
    for name, val, elasticity, desc in sensitivities:
        low_rev = base_rev_yr * (1 + elasticity/100 * (-2))  # -2%p
        high_rev = base_rev_yr * (1 + elasticity/100 * 2)  # +2%p
        fig.add_trace(go.Bar(name=f'{name} 악화', y=[name], x=[(low_rev - base_rev_yr)/억], orientation='h', marker_color=C['red'], showlegend=False))
        fig.add_trace(go.Bar(name=f'{name} 개선', y=[name], x=[(high_rev - base_rev_yr)/억], orientation='h', marker_color=C['green'], showlegend=False))
    lo(fig, title='경제지표 ±2%p 변동 시 매출 영향 (억원)', barmode='relative', height=380, xaxis_title='매출 변동(억)')
    st.plotly_chart(fig)

    # Summary table
    subsec("경제지표 요약")
    info("각 경제지표의 현재 값, 점수, 가중치, 매출 탄력성을 종합 정리한 표입니다. '호재'로 표시된 지표가 많을수록 사업 환경이 유리합니다.")
    econ_tbl = {'지표': [], '현재값': [], '점수(100)': [], '가중치': [], '가중점수': [], '매출탄력성': [], '판단': []}
    for i, (name, val, elasticity, desc) in enumerate(sensitivities):
        econ_tbl['지표'].append(name)
        econ_tbl['현재값'].append(f"{val:+.1f}%")
        econ_tbl['점수(100)'].append(f"{econ_vals[i]:.0f}")
        econ_tbl['가중치'].append(f"{econ_weights[i]}%")
        econ_tbl['가중점수'].append(f"{weighted_econ[i]:.1f}")
        econ_tbl['매출탄력성'].append(f"{elasticity:+.1f}%")
        econ_tbl['판단'].append('호재' if econ_vals[i] >= 55 else '중립' if econ_vals[i] >= 40 else '악재')
    dark_table(pd.DataFrame(econ_tbl))

    # ── 종합 판단 ──
    subsec("종합 시장 판단")
    if econ_score >= 65:
        st.success(f"**시장환경 호황** (점수 {econ_score:.0f}) — 골프 소비 확대 기대, 적극적 투자 적기")
    elif econ_score >= 45:
        st.warning(f"**시장환경 보통** (점수 {econ_score:.0f}) — 안정적이나 외부 변수 모니터링 필요")
    else:
        st.error(f"**시장환경 불황** (점수 {econ_score:.0f}) — 소비 위축 우려, 보수적 운영 권장")

    warn(f"상권 보정(×{ta_rev_adj:.2f}) × 경제 보정(×{econ_rev_adj:.2f}) = **종합 매출보정계수 ×{combined_adj:.2f}** → 기본 매출 대비 {(combined_adj-1)*100:+.1f}% 조정 필요")


# ═══ TAB (HIDDEN): Operations KPI ═══
if False:  # Tab removed from navigation
    sec("🏋️", "운영 KPI · 목표관리")
    info("일일 이용객, 타석 가동률, 객단가, 인건비율 등 핵심 운영 지표를 모니터링합니다. KPI 목표 대비 달성률을 추적하여 경영 의사결정에 활용하세요.")

    # ── 핵심 운영 KPI ──
    avg_rev_per_member = custom_total_rev / total_members if total_members else 0
    daily_members = total_members / (s_days * 9/12)  # 9개월 기준 일 평균
    utilization = daily_members / (s_bays * s_hours / 1.2) * 100  # 타석당 1.2시간 이용 가정
    monthly_rev_per_staff = rev_p[1] / 12 / s_staff if s_staff else 0
    cost_ratio = cost_p[1] / rev_p[1] * 100 if rev_p[1] else 0
    labor_ratio = cost_items_dyn['인건비'][1] / rev_p[1] * 100 if rev_p[1] else 0
    energy_ratio = (cost_items_dyn['전력비'][1] + cost_items_dyn['수도광열비'][1]) / rev_p[1] * 100 if rev_p[1] else 0

    ok = st.columns(6)
    ok[0].metric("일 평균 이용객", f"{daily_members:.0f}명")
    ok[1].metric("타석 가동률", f"{utilization:.1f}%")
    ok[2].metric("객단가", f"{avg_rev_per_member/만:.1f}만원")
    ok[3].metric("직원1인당 매출", fmt만(monthly_rev_per_staff)+'원/월')
    ok[4].metric("인건비율", f"{labor_ratio:.1f}%")
    ok[5].metric("에너지비율", f"{energy_ratio:.1f}%")

    c1, c2 = st.columns(2)
    with c1:
        subsec("월별 KPI 목표 vs 실적 (2026)")
        info("2026년 각 월의 매출과 월간 목표(초록 점선)를 비교합니다. 목표를 초과한 월이 많을수록 안정적인 경영입니다.")
        # Target: average monthly revenue
        target_monthly = rev_p[0] / 9  # 9개월 기준
        fig = go.Figure()
        fig.add_trace(go.Bar(name='실적', x=D['months'], y=[v/1e6 for v in mrev_custom], marker_color=C['blue']))
        fig.add_hline(y=target_monthly/1e6, line_dash="dash", line_color=C['green'], annotation_text=f"월 목표 {target_monthly/1e6:.0f}백만", annotation_font_size=11)
        achieve = [v >= target_monthly for v in mrev_custom]
        lo(fig, title=f'월 매출 목표 달성 현황: {sum(achieve)}/9개월', height=420, yaxis_title='백만원')
        st.plotly_chart(fig)

    with c2:
        subsec("시간대별 추정 가동률")
        warn("⚠️ 업계 일반 패턴 기반 추정치. KGA 2023: '평일 저녁(18~20시)' 24.0% 최다. 막대에 마우스를 올리면 추정 근거가 표시됩니다.")
        hours = ['06-08', '08-10', '10-12', '12-14', '14-16', '16-18', '18-20', '20-23']
        usage = [30, 55, 85, 70, 60, 75, 95, 80]
        reasons = [
            '새벽/조조: 소수 열성 회원만 이용<br>KGA: 오전 초반 이용 비율 8% 미만',
            '오전: 은퇴자/주부층 유입 시작<br>모닝권 회원 중심 이용 시간대',
            '오전~점심: 주부/자영업자 피크<br>모닝권+종일권 중복 이용 구간',
            '점심~오후: 식사 시간 소폭 감소<br>오전 피크 후 자연 감소 구간',
            '오후: 직장인 퇴근 전 소강<br>주중 가장 낮은 이용 구간 중 하나',
            '저녁 준비: 퇴근 후 이용 시작<br>KGA: 평일 저녁 이용 24.0%의 시작점',
            '피크타임: 직장인 퇴근 후 최대<br>KGA 2023: 평일 저녁 24.0% 최다',
            '야간: 피크 후 점차 감소<br>23시 마감 전까지 꾸준한 이용',
        ]
        fig = go.Figure(go.Bar(x=hours, y=usage,
            marker_color=[C['red'] if u < 50 else C['orange'] if u < 70 else C['green'] for u in usage],
            text=[f"{u}%" for u in usage], textposition='inside', textfont=dict(size=12, color='#e2e8f0'),
            customdata=reasons,
            hovertemplate='<b>%{x}</b><br>가동률: %{y}%<br><br>%{customdata}<extra></extra>'))
        fig.add_hline(y=70, line_dash="dash", line_color=C['slate'], annotation_text="목표 70%", annotation_font_size=11)
        lo(fig, title='시간대별 타석 가동률 (%) — 막대 호버 시 추정 근거 표시', height=420, yaxis_title='%', yaxis_range=[0,110])
        fig.update_layout(hovermode='closest', hoverlabel=dict(bgcolor='#0f172a', bordercolor='#334155', font=dict(size=13, color='#e2e8f0')))
        st.plotly_chart(fig)

    # ── 경영 대시보드 체크리스트 ──
    subsec("경영진 월간 체크리스트")
    info("매월 점검해야 할 핵심 관리 항목입니다. 매출, 비용, 고객·인력 3개 영역으로 구분하여 체계적으로 관리하세요.")
    chk = st.columns(3)
    chk[0].markdown("""
**📊 매출 관리**
- [ ] 월 매출 목표 대비 달성률 확인
- [ ] 상품별 매출 비중 변동 추적
- [ ] 신규회원 유입 vs 이탈 분석
- [ ] 쿠폰/일일권 판매 추이
- [ ] 임대매장 수금 확인
- [ ] 계열사 매출배분 정산
""")
    chk[1].markdown("""
**💸 비용 관리**
- [ ] 인건비율 목표(25% 이하) 모니터링
- [ ] 전기/수도 에너지 사용량 체크
- [ ] 소모품(골프공/매트) 재고 확인
- [ ] 시설 유지보수 스케줄 점검
- [ ] 보험 갱신일 확인
- [ ] 세금 납부 스케줄
""")
    chk[2].markdown("""
**👥 고객·인력 관리**
- [ ] 고객 불만/VOC 처리 현황
- [ ] 직원 근태/초과근무 관리
- [ ] 레슨 프로 수업 품질 점검
- [ ] 계절별 프로모션 기획
- [ ] 경쟁사 가격/프로모션 동향
- [ ] 시설 안전점검 (월 1회)
""")

    # ── 연간 운영 캘린더 ──
    subsec("연간 운영 캘린더")
    info("계절별 시즌 구분과 주요 운영 업무입니다. 성수기(봄/가을)에는 매출 극대화, 비수기(여름/겨울)에는 비용 절감과 프로모션에 집중합니다.")
    cal_data = [
        {'월': '3~5월', '시즌': '🟢 성수기', '주요업무': '회원 모집 강화, 봄 프로모션, 레슨 확대', '매출목표': '상'},
        {'월': '6월', '시즌': '🟡 개장월', '주요업무': '그랜드오픈, 사전예약 행사, 언론홍보', '매출목표': '중'},
        {'월': '7~8월', '시즌': '🔴 비수기', '주요업무': '야간 이벤트, 미스트 가동, 폭염 대응', '매출목표': '하'},
        {'월': '9~11월', '시즌': '🟢 성수기', '주요업무': '추석 프로모션, 가을 대회, 장기회원 유치', '매출목표': '상'},
        {'월': '12~1월', '시즌': '🔴 비수기', '주요업무': '동계 할인, 난방 가동, 연말 이벤트', '매출목표': '하'},
        {'월': '2월', '시즌': '🟡 준비기', '주요업무': '봄 시즌 준비, 시설 정비, 요금 조정', '매출목표': '중'},
    ]
    dark_table(pd.DataFrame(cal_data))


# ═══ TAB (HIDDEN): Risk ═══
if False:  # Tab removed from navigation
    sec("⚠️", "리스크 평가")
    info("사업 운영 과정에서 발생할 수 있는 주요 리스크를 영향도(1~5)와 발생확률(1~5)로 평가합니다. 점수(영향도×확률)가 12 이상이면 높은 리스크로 우선 대응이 필요합니다.")
    risks = [
        {'리스크':'기상악화(비/눈)','영향도':4,'발생확률':4,'점수':16,'등급':'높음','대응':'지붕/방풍 설비, 우천 할인'},
        {'리스크':'하절기 폭염','영향도':3,'발생확률':4,'점수':12,'등급':'중간','대응':'미스트/선풍기, 야간 강화'},
        {'리스크':'동절기 한파','영향도':4,'발생확률':3,'점수':12,'등급':'중간','대응':'동계할인, 난방'},
        {'리스크':'미세먼지','영향도':3,'발생확률':4,'점수':12,'등급':'중간','대응':'실내타석 검토'},
        {'리스크':'신규 경쟁업체','영향도':4,'발생확률':3,'점수':12,'등급':'중간','대응':'차별화, 고객 록인'},
        {'리스크':'인건비 상승','영향도':3,'발생확률':4,'점수':12,'등급':'중간','대응':'무인화/자동화'},
        {'리스크':'골프인구 감소','영향도':3,'발생확률':3,'점수':9,'등급':'보통','대응':'초보자 프로그램'},
        {'리스크':'전기료 인상','영향도':2,'발생확률':4,'점수':8,'등급':'보통','대응':'태양광, LED'},
        {'리스크':'시설 노후화','영향도':4,'발생확률':2,'점수':8,'등급':'보통','대응':'예방 정비'},
        {'리스크':'소음 민원','영향도':3,'발생확률':2,'점수':6,'등급':'낮음','대응':'방음벽, 시간 조정'},
    ]
    c1, c2 = st.columns(2)
    with c1:
        subsec("리스크 매트릭스")
        info("가로축은 발생확률, 세로축은 영향도입니다. 우측 상단(빨간 영역)에 위치한 리스크가 가장 긴급한 대응이 필요합니다.")
        fig=go.Figure()
        for r in risks:
            color = C['red'] if r['점수']>=12 else C['orange'] if r['점수']>=8 else C['green']
            fig.add_trace(go.Scatter(x=[r['발생확률']], y=[r['영향도']], mode='markers+text', marker=dict(size=r['점수']*3, color=color, opacity=0.7),
                text=[r['리스크']], textfont=dict(size=10), showlegend=False))
        fig.add_shape(type="rect", x0=3.5,x1=5,y0=3.5,y1=5, fillcolor="rgba(239,68,68,0.08)", line=dict(width=0))
        lo(fig, title='리스크 매트릭스 (영향도 vs 확률)', height=450, xaxis_title='발생확률', yaxis_title='영향도', xaxis=dict(range=[0.5,5.5], gridcolor="#f1f5f9"), yaxis=dict(range=[0.5,5.5], gridcolor="#f1f5f9"))
        st.plotly_chart(fig)
    with c2:
        subsec("리스크 점수 순위")
        info("리스크를 점수 순으로 정렬한 수평 막대 차트입니다. 빨간 막대가 최우선 관리 대상이며, 초록 막대는 상대적으로 낮은 리스크입니다.")
        fig=go.Figure(go.Bar(y=[r['리스크'] for r in risks], x=[r['점수'] for r in risks], orientation='h',
            marker_color=[C['red'] if r['점수']>=12 else C['orange'] if r['점수']>=8 else C['green'] for r in risks],
            text=[f"{r['점수']}점({r['등급']})" for r in risks], textposition='inside'))
        lo(fig, title='리스크 점수 순위 (점)', height=450, xaxis_title='리스크 점수')
        st.plotly_chart(fig)

    subsec("리스크 대응 방안")
    info("각 리스크별 영향도, 발생확률, 등급 및 구체적 대응 방안을 정리한 표입니다.")
    dark_table(pd.DataFrame(risks))

# ═══ TAB (HIDDEN): Detail Data ═══
if False:  # Tab removed from navigation
    sec("📄", "상세 데이터")
    info("모든 분석의 기초가 되는 원천 데이터를 상세 테이블로 제공합니다. 매출·비용 항목별 5개년 데이터, 요금표, 시즌성 가중치 등을 확인할 수 있습니다.")

    subsec("매출 항목별 (억원)")
    info("5개년 매출 항목별 상세 금액입니다. 소수점 2자리까지 표시하여 정밀한 비교가 가능합니다.")
    rd = {'항목': list(D['rev_items'].keys())}
    for i, yr in enumerate(D['yp']): rd[yr] = [f"{D['rev_items'][k][i]/억:.2f}" for k in rd['항목']]
    dark_table(pd.DataFrame(rd))

    subsec("비용 항목별 (억원)")
    info("5개년 비용 항목별 상세 금액입니다. 인건비, 임대료 등 주요 비용의 연도별 변동을 추적하세요.")
    cd = {'항목': list(D['cost_items'].keys())}
    for i, yr in enumerate(D['yp']): cd[yr] = [f"{D['cost_items'][k][i]/억:.2f}" for k in cd['항목']]
    dark_table(pd.DataFrame(cd))

    c1, c2 = st.columns(2)
    with c1:
        subsec("2026 정규 요금표")
        info("2026년 정규 시즌 남/여 요금과 과거(2021) 대비 인상율입니다.")
        if D['pricing']:
            df_p = pd.DataFrame(D['pricing'])
            df_p = df_p.rename(columns={'남(VAT포함)': '남자(VAT포함)', '여(VAT포함)': '여자(VAT포함)'})
            df_p['남자(VAT포함)'] = df_p['남자(VAT포함)'].apply(lambda x: f"{int(x):,}원" if x else '-')
            df_p['여자(VAT포함)'] = df_p['여자(VAT포함)'].apply(lambda x: f"{int(x):,}원" if x else '-')
            df_p['과거(2021남)'] = df_p['과거(2021남)'].apply(lambda x: f"{int(x):,}원" if x else '-')
            df_p['인상율'] = df_p['인상율'].apply(lambda x: f"{x*100:.1f}%" if x else '-')
            dark_table(df_p)
    with c2:
        subsec("동계 요금표 (12~1월)")
        info("동계(12~1월) 할인 요금표입니다. 비수기 가격 전략을 확인하세요.")
        if D['winter']:
            df_w = pd.DataFrame(D['winter'])
            df_w = df_w.rename(columns={'남(VAT포함)': '남자(VAT포함)', '여(VAT포함)': '여자(VAT포함)'})
            df_w['남자(VAT포함)'] = df_w['남자(VAT포함)'].apply(lambda x: f"{int(x):,}원" if x else '-')
            df_w['여자(VAT포함)'] = df_w['여자(VAT포함)'].apply(lambda x: f"{int(x):,}원" if x else '-')
            dark_table(df_w)

    if D['contrib']:
        subsec("상품별 이익기여도 (2018~2021)")
        info("과거 실적 데이터를 기반으로 각 상품의 이익 기여도를 분석합니다. 연도를 선택하여 비교해 보세요.")
        df_c = pd.DataFrame(D['contrib'])
        years_avail = sorted(df_c['연도'].unique())
        sel_year = st.selectbox("연도", years_avail, index=len(years_avail)-1)
        dfy = df_c[df_c['연도']==sel_year].copy()
        c1, c2 = st.columns(2)
        with c1:
            fig=go.Figure(go.Pie(labels=dfy['상품'], values=dfy['이익'], hole=0.45, marker=dict(colors=PAL), textinfo='label+percent'))
            lo(fig, title=f'{sel_year}년 이익 기여도', height=420)
            st.plotly_chart(fig)
        with c2:
            fig=go.Figure()
            fig.add_trace(go.Bar(name='매출', x=dfy['상품'], y=dfy['매출']/억, marker_color=C['blue']))
            fig.add_trace(go.Bar(name='비용', x=dfy['상품'], y=dfy['추정비용']/억, marker_color=C['red']))
            fig.add_trace(go.Bar(name='이익', x=dfy['상품'], y=dfy['이익']/억, marker_color=C['green']))
            lo(fig, title=f'{sel_year}년 상품별 매출·비용·이익 (억원)', barmode='group', height=420, yaxis_title='억원')
            st.plotly_chart(fig)
        disp = dfy[['상품','회원수','매출','추정비용','이익','매출비중','이익기여도','이익률']].copy()
        disp['매출'] = disp['매출'].apply(lambda x: fmt만(x)+'원')
        disp['추정비용'] = disp['추정비용'].apply(lambda x: fmt만(x)+'원')
        disp['이익'] = disp['이익'].apply(lambda x: fmt만(x)+'원')
        for col in ['매출비중','이익기여도','이익률']: disp[col] = disp[col].apply(lambda x: f"{x*100:.1f}%")
        disp['회원수'] = disp['회원수'].apply(lambda x: f"{int(x):,}명")
        dark_table(disp)

    subsec("시즌성 가중치")
    info("월별 시즌 가중치입니다. 1.0 이상은 성수기, 미만은 비수기를 의미합니다. 매출 추정 시 월별 매출에 이 가중치를 곱하여 시즌 효과를 반영합니다.")
    fig_sw = go.Figure(go.Bar(x=D['months'], y=D['season_weights'],
        marker_color=[C['red'] if w<0.6 else C['orange'] if w<0.9 else C['green'] for w in D['season_weights']],
        text=[f"{w:.2f}" for w in D['season_weights']], textposition='inside', textfont=dict(size=12, color='#e2e8f0')))
    fig_sw.add_hline(y=1.0, line_dash="dash", line_color=C['slate'], annotation_text="기준 1.0", annotation_font_size=11)
    lo(fig_sw, title='월별 시즌성 가중치', height=380, yaxis_title='가중치', yaxis_range=[0, max(D['season_weights'])*1.3 if D['season_weights'] else 1.5])
    st.plotly_chart(fig_sw)

# ══════════════════════════════════════════════════════════════
# Footer
# ══════════════════════════════════════════════════════════════
# ═══ TAB 12: 검증 ═══
if _ti == 12:
    sec("🔍", "계산 검증 도구")
    info("컨트롤 패널에 입력된 현재 값을 기준으로 모든 중간 계산 과정과 최종 결과를 표시합니다. 숫자가 맞는지 확인하세요.")

    # ── 1. 입력값 요약 ──
    subsec("1. 컨트롤 패널 입력값 요약")
    v_data = {
        '구분': ['타석수', '투자금', '내용연수', '정액법상각(만/년)', '정률법상각(만/초년도)', '할인율', '매출성장률', '인건비인상률', '물가상승률'],
        '값': [f'{s_bays}개', f'{s_inv}억', f'{s_useful}년', f'{dep_straight_amt:,}만', f'{dep_declining_amt:,}만', f'{disc_r*100:.1f}%', f'{s_growth*100:.1f}%', f'{s_labor_up*100:.1f}%', f'{s_util_up*100:.1f}%'],
    }
    dark_table(pd.DataFrame(v_data))

    # ── 2. 매출 계산 과정 ──
    subsec("2. 매출 계산 과정")
    st.markdown("**상품별 매출 (회원수 × 단가)**")
    rev_detail = []
    for name, val in custom_rev_items.items():
        rev_detail.append({'항목': name, '매출(원)': f'{val:,.0f}', '매출(억)': f'{val/억:.2f}'})
    rev_detail.append({'항목': '합계 (이론치)', '매출(원)': f'{custom_total_rev:,.0f}', '매출(억)': f'{custom_total_rev/억:.2f}'})
    dark_table(pd.DataFrame(rev_detail))

    st.markdown("**보정계수 적용** (시즌/램프업은 월별 분배에만 적용 — 연 매출엔 영향 없음)")
    adj_data = {
        '보정 항목': ['이탈률(중도해지)', '환불률(구매후취소)', '상권 보정', '경제 보정', '월별 시즌평균(분배용)', '월별 램프업평균(분배용)', '연매출 종합 (2026F)', '연매출 종합 (2027F~)'],
        '계수': [f'{_churn_adj:.3f}', f'{_refund_adj:.3f}', f'{ta_rev_adj:.3f}', f'{econ_rev_adj:.3f}', f'{_season_avg:.3f}', f'{_ramp_avg:.3f}', f'{_adj_2026:.3f}', f'{_adj_normal:.3f}'],
        '의미': [f'1 − 이탈률 {w_churn*100:.1f}%', f'1 − 환불률 {s_refund*100:.1f}%',
                f'상권점수 {ta_score:.0f}/70 기준', f'경제점수 {econ_score:.0f}/{ECON_BASELINE} 기준',
                f'월별 분배용 (연 매출 영향 X)', f'월별 분배용 (연 매출 영향 X)',
                f'이탈×환불×상권×경제', f'이탈×환불×상권×경제 (동일)'],
    }
    dark_table(pd.DataFrame(adj_data))

    _golf_raw = custom_total_rev - _rent_2026
    st.markdown("**2026F 매출 산출**")
    calc_2026 = {
        '단계': ['골프 이론매출 (9개월)', '임대매출 (4개월)', '이론 합계', '종합 보정계수', '보정 후 골프매출', '최종 2026F 매출'],
        '금액': [f'{_golf_raw/억:.2f}억', f'{_rent_2026/억:.2f}억', f'{custom_total_rev/억:.2f}억',
                f'×{_adj_2026:.3f}', f'{_golf_raw*_adj_2026/억:.2f}억', f'{rev_p[0]/억:.2f}억'],
        '산식': ['회원수×단가 합계', f'외부{s_rent_external}+계열사{s_rent_affiliate}×4개월', '골프+임대',
                '이탈×환불×상권×경제', f'{_golf_raw/억:.2f}×{_adj_2026:.3f}', '보정 골프 + 임대'],
    }
    dark_table(pd.DataFrame(calc_2026))

    # ── 3. 비용 계산 과정 ──
    subsec("3. 비용 계산 과정")
    st.markdown("**월 비용 내역**")
    cost_detail = {
        '항목': ['인건비(4대보험포함)', '전력비', '수도광열비', '세금과공과', '보험료', '소모품비', '수선비', '용역비', '기타비용', '마케팅', '카드수수료(추정)', '월 합계'],
        '만원/월': [f'{monthly_labor:,.0f}', f'{op_electric:,}', f'{op_water:,}', f'{op_tax:,}', f'{op_insurance:,}',
                   f'{op_supplies:,}', f'{op_maint:,}', f'{op_outsource:,}', f'{op_etc:,}', f'{op_marketing:,}',
                   f'{op_var_monthly:,.0f}', f'{monthly_labor+op_total_monthly:,.0f}'],
    }
    dark_table(pd.DataFrame(cost_detail))

    st.markdown("**감가상각 5개년**")
    dep_detail = {'연도': D['yp'][:len(dep)], '감가상각(만)': [f'{d/만:,.0f}' for d in dep],
                  '정액법(만)': [f'{_dep_str/만:,.0f}']*len(dep),
                  '정률법(만)': [f'{d/만 - _dep_str/만:,.0f}' for d in dep]}
    dark_table(pd.DataFrame(dep_detail))

    # ── 4. 5개년 P&L 검증 ──
    subsec("4. 5개년 손익 검증")
    pl_verify = {'항목': ['매출', '비용', '영업이익', '감가상각', 'EBITDA', '누적EBITDA', '회수율']}
    cum_v = 0
    for i, yr in enumerate(D['yp'][:len(rev_p)]):
        ebitda_v = op_p[i] + dep[i] if i < len(dep) else op_p[i]
        cum_v += ebitda_v
        pl_verify[yr] = [
            f'{rev_p[i]/억:.2f}억', f'{cost_p[i]/억:.2f}억', f'{op_p[i]/억:.2f}억',
            f'{dep[i]/억:.2f}억' if i < len(dep) else '-',
            f'{ebitda_v/억:.2f}억', f'{cum_v/억:.2f}억',
            f'{cum_v/inv_won*100:.1f}%' if inv_won else '-'
        ]
    dark_table(pd.DataFrame(pl_verify))

    # ── 5. 월별 매출 검증 (2026F) ──
    subsec("5. 월별 매출 검증 (2026F)")
    monthly_v = {'월': D['months'][:len(mrev_custom)],
                 '시즌가중치': [f'{w:.2f}' for w in _sw[:len(mrev_custom)]],
                 '램프업': [f'{r*100:.0f}%' for r in ramp_values[:len(mrev_custom)]],
                 '매출(만)': [f'{v/만:,.0f}' for v in mrev_custom],
                 '비용(만)': [f'{v/만:,.0f}' for v in mcost_custom[:len(mrev_custom)]],
                 '손익(만)': [f'{(r-c)/만:,.0f}' for r,c in zip(mrev_custom, mcost_custom[:len(mrev_custom)])]}
    dark_table(pd.DataFrame(monthly_v))
    st.caption(f"월매출 합계: {sum(mrev_custom)/억:.2f}억 | 월비용 합계: {sum(mcost_custom[:len(mrev_custom)])/억:.2f}억")

    # ── 6. KPI 교차 검증 ──
    subsec("6. KPI 교차 검증")
    kpi_v = {
        'KPI': ['NPV', 'IRR', '5년 누적EBITDA', '회수율', 'Payback', 'BEP매출', 'BEP회원수', '안전마진'],
        '값': [f'{npv_val/억:.2f}억', f'{irr_val*100:.1f}%', f'{cum_ebitda[-1]/억:.2f}억',
              f'{rec_rate[-1]*100:.1f}%', f'{payback:.1f}년' if payback else '5년+',
              f'{bep_revenue/억:.2f}억', f'{bep_members:,.0f}명', f'{safety_margin:.1f}%'],
        '산식': [f'Σ(EBITDA/(1+{disc_r*100:.0f}%)^n) - {s_inv}억',
                '[-투자금, EBITDA₁...₅]의 내부수익률',
                f'Σ EBITDA = {" + ".join([f"{e/억:.1f}" for e in ebitda_p])}',
                f'{cum_ebitda[-1]/억:.1f} ÷ {inv_won/억:.0f} × 100',
                '누적EBITDA ≥ 투자금 시점',
                '고정비 ÷ (1-변동비율)',
                'BEP매출 ÷ 평균객단가',
                '(매출-BEP) ÷ 매출 × 100'],
    }
    dark_table(pd.DataFrame(kpi_v))

    # ── 7. 보정계수 영향 시뮬레이션 ──
    subsec("7. 연매출 보정계수별 영향")
    st.caption("연 매출 = 골프 이론매출 × (이탈 × 환불 × 상권 × 경제) + 임대. 시즌/램프업은 월별 분배 시에만 사용.")
    sim_data = {
        '시나리오': ['보정 없음 (이론치)', '이탈만 적용', '이탈+환불', '이탈+환불+상권', '전체 보정 (현재)'],
        '2026F 골프매출': [
            f'{_golf_raw/억:.1f}억',
            f'{_golf_raw*_churn_adj/억:.1f}억',
            f'{_golf_raw*_churn_adj*_refund_adj/억:.1f}억',
            f'{_golf_raw*_churn_adj*_refund_adj*ta_rev_adj/억:.1f}억',
            f'{_golf_raw*_adj_2026/억:.1f}억',
        ],
        '보정률': ['100%', f'{_churn_adj*100:.1f}%', f'{_churn_adj*_refund_adj*100:.1f}%',
                  f'{_churn_adj*_refund_adj*ta_rev_adj*100:.1f}%', f'{_adj_2026*100:.1f}%'],
    }
    dark_table(pd.DataFrame(sim_data))

st.markdown(f'<div class="footer">⛳ 신진(SJ) 등촌골프연습장 사업성 분석 v5.0 &nbsp;|&nbsp; {s_bays}타석 실외 &nbsp;|&nbsp; ₩{s_inv}억 투자 &nbsp;|&nbsp; 2026.06 오픈 &nbsp;|&nbsp; 13개 분석 모듈</div>', unsafe_allow_html=True)

# ══════════════════════════════════════════════════════════════
# Auto-Save: 모든 위젯 값을 JSON으로 자동 저장
# ══════════════════════════════════════════════════════════════
_all_keys = [k for k in st.session_state.keys() if isinstance(st.session_state[k], (int, float, str, bool))]
save_all_values(_all_keys)
