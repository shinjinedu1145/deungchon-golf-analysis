# -*- coding: utf-8 -*-
"""
등촌골프연습장 사업성 분석 — 경영진 보고용 5장 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ──────────────────────────────────────────────────────────────────
# Design Tokens (Midnight Executive Palette)
# ──────────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1E, 0x27, 0x61)  # primary
DEEP_NAVY = RGBColor(0x14, 0x1B, 0x44)  # darker
ICE_BLUE  = RGBColor(0xCA, 0xDC, 0xFC)  # secondary
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GOLD      = RGBColor(0xFB, 0xBF, 0x24)  # accent
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
RED       = RGBColor(0xDC, 0x26, 0x26)
SLATE_700 = RGBColor(0x33, 0x4E, 0x68)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
SLATE_BG  = RGBColor(0xF8, 0xFA, 0xFC)

# ──────────────────────────────────────────────────────────────────
# Setup
# ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9 WIDE
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

# ──────────────────────────────────────────────────────────────────
# Helper functions
# ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh

def add_text(slide, x, y, w, h, text, *, font='Malgun Gothic', size=14, bold=False,
             color=SLATE_700, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_multitext(slide, x, y, w, h, runs, *, valign=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
                  line_spacing=1.2):
    """runs = list of dicts: {text, size, bold, color, italic, break_after}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for i, run in enumerate(runs):
        if i > 0 and runs[i-1].get('break_after'):
            p = tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = run['text']
        r.font.name = run.get('font', 'Malgun Gothic')
        r.font.size = Pt(run.get('size', 14))
        r.font.bold = run.get('bold', False)
        r.font.italic = run.get('italic', False)
        r.font.color.rgb = run.get('color', SLATE_700)
    return tb

def add_table(slide, x, y, w, h, data, *, col_widths=None,
              header_fill=NAVY, header_color=WHITE,
              row_fills=(WHITE, SLATE_BG), border_color=SLATE_300,
              size=11, header_size=11, bold_first_col=False, highlight_last_col=False):
    """data: list of lists. First row is header."""
    rows = len(data); cols = len(data[0])
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for ri, row in enumerate(data):
        is_header = (ri == 0)
        for ci, cell_val in enumerate(row):
            cell = tbl.cell(ri, ci)
            if is_header:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                fill = row_fills[(ri-1) % len(row_fills)]
                cell.fill.solid(); cell.fill.fore_color.rgb = fill
            tf = cell.text_frame
            tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            # First col left-align bold, others center
            if ci == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = str(cell_val)
            r.font.name = 'Malgun Gothic'
            r.font.size = Pt(header_size if is_header else size)
            r.font.bold = is_header or (ci == 0 and bold_first_col)
            if is_header:
                r.font.color.rgb = header_color
            else:
                if highlight_last_col and ci == cols - 1:
                    r.font.color.rgb = NAVY
                    r.font.bold = True
                else:
                    r.font.color.rgb = SLATE_700
    return tbl

def add_chip(slide, x, y, text, *, fill=GOLD, fg=DEEP_NAVY, size=10):
    """small inline chip — auto sized to text"""
    w = max(0.6, len(text) * 0.10) * 914400
    h = Inches(0.28)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, int(w), h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    tf = sh.text_frame
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = 'Malgun Gothic'
    r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = fg
    return sh

def add_footer(slide, page_num, total=5):
    """slide footer with page number"""
    add_rect(slide, Emu(0), SH - Inches(0.35), SW, Inches(0.35), DEEP_NAVY)
    add_text(slide, Inches(0.5), SH - Inches(0.32), Inches(8), Inches(0.30),
             '등촌골프연습장 사업성 분석 · 경영진 보고용 · 신진(SJ)',
             size=9, color=ICE_BLUE, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, SW - Inches(1.5), SH - Inches(0.32), Inches(1.2), Inches(0.30),
             f'{page_num} / {total}',
             size=9, color=ICE_BLUE, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, bold=True)

def add_page_title(slide, num, title, subtitle=None):
    """Standard page header for content slides"""
    # Top navy strip
    add_rect(slide, Emu(0), Emu(0), SW, Inches(0.85), NAVY)
    # Page number badge
    add_rect(slide, Inches(0.45), Inches(0.22), Inches(0.4), Inches(0.4), GOLD)
    add_text(slide, Inches(0.45), Inches(0.22), Inches(0.4), Inches(0.4),
             str(num), size=18, bold=True, color=DEEP_NAVY,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    # Title
    add_text(slide, Inches(1.05), Inches(0.18), Inches(11), Inches(0.42),
             title, size=22, bold=True, color=WHITE, valign=MSO_ANCHOR.MIDDLE)
    # Subtitle
    if subtitle:
        add_text(slide, Inches(1.05), Inches(0.55), Inches(11), Inches(0.28),
                 subtitle, size=11, color=ICE_BLUE, valign=MSO_ANCHOR.TOP)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# Dark navy background
add_rect(s1, Emu(0), Emu(0), SW, SH, DEEP_NAVY)
# Decorative accent diagonal lines (top-right)
add_rect(s1, SW - Inches(0.08), Inches(0.4), Inches(0.08), Inches(2.5), GOLD)
add_rect(s1, SW - Inches(0.24), Inches(0.4), Inches(0.08), Inches(1.5), ICE_BLUE)

# Brand / category
add_text(s1, Inches(0.8), Inches(0.7), Inches(6), Inches(0.4),
         '경영진 투자 검토 보고서', size=12, bold=True, color=GOLD, italic=True)

# Big title
add_text(s1, Inches(0.8), Inches(1.3), Inches(11.5), Inches(1.0),
         '등촌골프연습장', size=44, bold=True, color=WHITE)
add_text(s1, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.7),
         '사업성 분석 — 5개년·10개년 재무 전망', size=24, color=ICE_BLUE)

# Divider
add_rect(s1, Inches(0.8), Inches(3.2), Inches(2), Inches(0.04), GOLD)

# Key facts row
fact_y = Inches(3.7)
fact_w = Inches(2.8)
fact_gap = Inches(0.2)
fact_h = Inches(1.4)
facts = [
    ('88', '타석', '실외 대형 88타석'),
    ('17억', '투자금', '자가소유 부지'),
    ('2026.06', '재오픈', '5년 운영 중단 후 재가동'),
    ('8%', 'WACC', '재오픈·자가소유 기준'),
]
for i, (val, lbl, sub) in enumerate(facts):
    fx = Inches(0.8) + i * (fact_w + fact_gap)
    # Card
    add_rect(s1, fx, fact_y, fact_w, fact_h, NAVY)
    add_rect(s1, fx, fact_y, Inches(0.08), fact_h, GOLD)  # left accent
    add_text(s1, fx + Inches(0.3), fact_y + Inches(0.18), fact_w - Inches(0.3), Inches(0.4),
             lbl, size=11, color=ICE_BLUE, bold=True)
    add_text(s1, fx + Inches(0.3), fact_y + Inches(0.45), fact_w - Inches(0.3), Inches(0.6),
             val, size=28, color=WHITE, bold=True)
    add_text(s1, fx + Inches(0.3), fact_y + Inches(1.0), fact_w - Inches(0.3), Inches(0.35),
             sub, size=10, color=ICE_BLUE)

# Footer (date + author)
add_text(s1, Inches(0.8), SH - Inches(0.8), Inches(8), Inches(0.3),
         '작성일: 2026년 5월 · 작성: 신진(SJ) 경영기획팀',
         size=11, color=SLATE_300)
add_text(s1, SW - Inches(4), SH - Inches(0.8), Inches(3.2), Inches(0.3),
         'CONFIDENTIAL · 대외비',
         size=10, bold=True, color=GOLD, align=PP_ALIGN.RIGHT, italic=True)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2: 사업 개요 & 핵심 가정
# ══════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, Emu(0), Emu(0), SW, SH, SLATE_BG)
add_page_title(s2, '01', '사업 개요 및 핵심 가정',
               '재무 전망의 근거가 되는 5대 가정 — 매출·비용·자본구조·시장·자산')

# Left column: 사업 개요 (KPI grid)
left_x = Inches(0.5); right_x = Inches(7.0)
cw = Inches(6.3)
top_y = Inches(1.15)

# Section: 사업 개요
add_text(s2, left_x, top_y, cw, Inches(0.35),
         '■  사업 개요', size=15, bold=True, color=NAVY)

biz_y = top_y + Inches(0.45)
biz_data = [
    ['항목', '내용'],
    ['소재지', '서울 강서구 등촌동 (자가소유 부지)'],
    ['시설 규모', '88타석 실외 + 사우나·헬스·락카'],
    ['오픈 일자', '2026년 6월 (9개월 운영) / 2027년 정상가동'],
    ['투자 규모', '17억원 (시설·장비·인테리어·운영자금 포함)'],
    ['감가상각', '5년 내용연수 (정액 2,000만 + 정률 6.36억 첫해)'],
]
add_table(s2, left_x, biz_y, cw, Inches(2.2), biz_data,
          col_widths=[Inches(1.6), Inches(4.7)], size=11, header_size=11)

# Section: 회원수 회복률
biz2_y = biz_y + Inches(2.4)
add_text(s2, left_x, biz2_y, cw, Inches(0.35),
         '■  회원수 회복률 (2018 등촌 실적 대비)', size=15, bold=True, color=NAVY)

rec_y = biz2_y + Inches(0.45)
rec_data = [
    ['연도', '2026', '2027', '2028', '2029', '2030', '2031~35'],
    ['회복률', '60%×9/12', '80%', '80%', '82%', '82%', '매년 -3%'],
    ['근거', '오픈 9개월', '시장조사', 'Excel Base', 'Excel Base', 'Excel Base', '시장 자연 축소'],
]
add_table(s2, left_x, rec_y, cw, Inches(1.5), rec_data,
          col_widths=[Inches(1.0), Inches(0.85), Inches(0.8), Inches(0.85), Inches(0.85), Inches(0.85), Inches(1.1)],
          size=10, header_size=10)

# Right column
# Section: 가격 정책
add_text(s2, right_x, top_y, cw, Inches(0.35),
         '■  가격 정책', size=15, bold=True, color=NAVY)
price_y = top_y + Inches(0.45)
price_data = [
    ['연도', '가격 기준', '비고'],
    ['2026', '사이드바 입력가', '오픈 첫해, 제니스 대비 저렴'],
    ['2027', '사이드바 입력가', '동결'],
    ['2028~31', '제니스 동일가', '시장가 따라잡기'],
    ['2032', '제니스 ×1.05', '+5% 인상'],
    ['2033~34', '동결 (×1.05)', '인상 효과 유지'],
    ['2035', '제니스 ×1.1025', '+5% 추가 인상'],
]
add_table(s2, right_x, price_y, cw, Inches(2.4), price_data,
          col_widths=[Inches(1.1), Inches(2.1), Inches(3.1)],
          size=10, header_size=10)

# Section: 5대 핵심 가정 (자본·시장)
mkt_y = price_y + Inches(2.6)
add_text(s2, right_x, mkt_y, cw, Inches(0.35),
         '■  자본 구조 · 시장 가정', size=15, bold=True, color=NAVY)
mkt_y2 = mkt_y + Inches(0.45)
mkt_data = [
    ['항목', '값'],
    ['자본 조달', '100% 자기자본 (차입 0)'],
    ['WACC (할인율)', '8% — 무위험 3% + 사업리스크 3% + 유동성 2%'],
    ['법인세율', '22% (적자 시 0)'],
    ['보정계수 (2026)', '이탈×환불×상권×경제×시즌 = 약 0.60'],
    ['보정계수 (2027~)', '이탈×환불×상권×경제 = 약 0.80'],
]
add_table(s2, right_x, mkt_y2, cw, Inches(1.7), mkt_data,
          col_widths=[Inches(1.8), Inches(4.5)], size=10, header_size=10)

add_footer(s2, 2)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3: 5개년 손익 전망 (2026~2030)
# ══════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, Emu(0), Emu(0), SW, SH, SLATE_BG)
add_page_title(s3, '02', '5개년 손익 전망 (2026~2030)',
               '한국 회계기준 손익계산서 · 임대료수익은 영업외수익으로 분리')

# KPI cards row (4 cards)
kpi_y = Inches(1.15)
kpi_w = Inches(2.95); kpi_h = Inches(1.3); kpi_gap = Inches(0.2)
kpis = [
    ('5년 NPV', '+0.5억', 'WACC 8% 기준', GREEN),
    ('5년 IRR', '8.9%', 'WACC 8% 소폭 상회', GOLD),
    ('Payback', '3.5년', '5년 내 회수', GREEN),
    ('5년 누적 EBITDA', '28.3억', '투자금 17억의 166%', GOLD),
]
for i, (lbl, val, sub, accent) in enumerate(kpis):
    kx = Inches(0.5) + i * (kpi_w + kpi_gap)
    add_rect(s3, kx, kpi_y, kpi_w, kpi_h, WHITE)
    add_rect(s3, kx, kpi_y, Inches(0.08), kpi_h, accent)  # left accent
    add_text(s3, kx + Inches(0.25), kpi_y + Inches(0.15), kpi_w - Inches(0.3), Inches(0.3),
             lbl, size=11, bold=True, color=SLATE_500)
    add_text(s3, kx + Inches(0.25), kpi_y + Inches(0.42), kpi_w - Inches(0.3), Inches(0.6),
             val, size=26, bold=True, color=NAVY)
    add_text(s3, kx + Inches(0.25), kpi_y + Inches(0.97), kpi_w - Inches(0.3), Inches(0.3),
             sub, size=10, color=SLATE_500, italic=True)

# 5년 손익계산서 표
pl_y = kpi_y + Inches(1.55)
add_text(s3, Inches(0.5), pl_y, Inches(8), Inches(0.32),
         '■  5개년 추정 손익계산서 (단위: 억)', size=14, bold=True, color=NAVY)

pl_data = [
    ['항목',         '2026',  '2027',  '2028',  '2029',  '2030'],
    ['매출(영업)',   '10.79', '24.85', '26.82', '27.45', '27.45'],
    ['영업비용',     '21.43', '24.84', '24.13', '24.04', '24.33'],
    ['영업이익',     '-10.64', '0.02',  '2.69',  '3.41',  '3.11'],
    ['영업외수익',   '1.23',  '3.19',  '3.19',  '3.19',  '3.19'],
    ['경상이익',     '-9.41', '3.20',  '5.88',  '6.60',  '6.30'],
    ['당기순이익',   '-9.41', '2.50',  '4.58',  '5.15',  '4.92'],
    ['EBITDA',       '-2.8',  '7.2',   '8.4',   '8.2',   '7.3'],
    ['누적 FCF',     '-2.8',  '3.7',   '10.7',  '17.5',  '23.4'],
]
add_table(s3, Inches(0.5), pl_y + Inches(0.4), Inches(8.5), Inches(3.5), pl_data,
          col_widths=[Inches(1.7), Inches(1.36), Inches(1.36), Inches(1.36), Inches(1.36), Inches(1.36)],
          size=11, header_size=11, bold_first_col=True)

# Right side: 출구전략 시나리오 + 시사점
exit_x = Inches(9.2); exit_w = Inches(3.8)
# Header
add_rect(s3, exit_x, pl_y, exit_w, Inches(0.42), GOLD)
add_text(s3, exit_x + Inches(0.15), pl_y, exit_w - Inches(0.3), Inches(0.42),
         '🏆  출구전략 시나리오', size=12, bold=True, color=DEEP_NAVY,
         valign=MSO_ANCHOR.MIDDLE)

# Card body
add_rect(s3, exit_x, pl_y + Inches(0.42), exit_w, Inches(3.45), WHITE,
         line=GOLD)

add_multitext(s3, exit_x + Inches(0.2), pl_y + Inches(0.6), exit_w - Inches(0.4), Inches(3.25), [
    {'text': '2030년말 매각 가정', 'size': 12, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '매각가:  ', 'size': 11, 'color': SLATE_700},
    {'text': '3,000억', 'size': 11, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '할인 기간:  5년', 'size': 11, 'color': SLATE_700, 'break_after': True},
    {'text': 'WACC:  8% (할인계수 0.6806)', 'size': 11, 'color': SLATE_700, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},
    {'text': '매각가 PV', 'size': 11, 'color': SLATE_500, 'break_after': True},
    {'text': '= 3,000억 × 0.6806 = ', 'size': 11, 'color': SLATE_700},
    {'text': '2,041.7억', 'size': 11, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},
    {'text': '기본 NPV: ', 'size': 11, 'color': SLATE_700},
    {'text': '+0.55억', 'size': 11, 'bold': True, 'color': GREEN, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},
    {'text': '매각 반영 NPV', 'size': 12, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '2,042.3억', 'size': 26, 'bold': True, 'color': GOLD},
])

add_footer(s3, 3)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4: 10개년 확장 전망 (2026~2035)
# ══════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, Emu(0), Emu(0), SW, SH, SLATE_BG)
add_page_title(s4, '03', '10개년 확장 전망 (2026~2035)',
               '2031~ 회원 -3%/년 감소 · 2032·2035 요금 +5% 인상 · CAPEX 2.5억(2033)')

# KPI row
kpi_y = Inches(1.15)
kpis2 = [
    ('10년 NPV', '+8.3억', 'WACC 8% 기준', GREEN),
    ('10년 IRR', '17.6%', '5년 IRR 8.9% 대비 향상', GREEN),
    ('Payback', '3.5년', '동일', GREEN),
    ('10년 누적 FCF', '34.4억', 'CAPEX 2.5억 반영', GOLD),
]
for i, (lbl, val, sub, accent) in enumerate(kpis2):
    kx = Inches(0.5) + i * (kpi_w + kpi_gap)
    add_rect(s4, kx, kpi_y, kpi_w, kpi_h, WHITE)
    add_rect(s4, kx, kpi_y, Inches(0.08), kpi_h, accent)
    add_text(s4, kx + Inches(0.25), kpi_y + Inches(0.15), kpi_w - Inches(0.3), Inches(0.3),
             lbl, size=11, bold=True, color=SLATE_500)
    add_text(s4, kx + Inches(0.25), kpi_y + Inches(0.42), kpi_w - Inches(0.3), Inches(0.6),
             val, size=26, bold=True, color=NAVY)
    add_text(s4, kx + Inches(0.25), kpi_y + Inches(0.97), kpi_w - Inches(0.3), Inches(0.3),
             sub, size=10, color=SLATE_500, italic=True)

# 10년 손익 추이 표
pl_y = kpi_y + Inches(1.55)
add_text(s4, Inches(0.5), pl_y, Inches(8), Inches(0.32),
         '■  10개년 매출·EBITDA·FCF 추이 (단위: 억)', size=14, bold=True, color=NAVY)

pl10_data = [
    ['연도',     '2026',  '2027', '2028', '2029', '2030', '2031', '2032', '2033', '2034', '2035'],
    ['매출',     '12.0',  '28.0', '30.0', '30.6', '30.6', '29.8', '30.3', '29.5', '28.7', '29.2'],
    ['비용',     '21.4',  '24.8', '24.1', '24.0', '24.3', '24.2', '25.1', '26.0', '27.0', '28.1'],
    ['영업이익', '-9.4',  '3.2',  '5.9',  '6.6',  '6.3',  '5.6',  '5.2',  '3.5',  '1.7',  '1.2'],
    ['EBITDA',   '-2.8',  '7.2',  '8.4',  '8.2',  '7.3',  '5.6',  '5.2',  '3.5',  '1.7',  '1.2'],
    ['CAPEX',    '-',     '-',    '-',    '-',    '-',    '-',    '-',    '2.5',  '-',    '-'],
    ['누적 FCF', '-2.8',  '3.7',  '10.7', '17.5', '23.4', '27.8', '31.9', '32.1', '33.5', '34.4'],
]
add_table(s4, Inches(0.5), pl_y + Inches(0.4), Inches(8.5), Inches(2.8), pl10_data,
          col_widths=[Inches(0.85)] + [Inches(0.765)]*10,
          size=10, header_size=10, bold_first_col=True)

# Right side: 출구전략 시나리오
exit_x = Inches(9.2); exit_w = Inches(3.8)
add_rect(s4, exit_x, pl_y, exit_w, Inches(0.42), GOLD)
add_text(s4, exit_x + Inches(0.15), pl_y, exit_w - Inches(0.3), Inches(0.42),
         '🏆  출구전략 시나리오', size=12, bold=True, color=DEEP_NAVY,
         valign=MSO_ANCHOR.MIDDLE)

add_rect(s4, exit_x, pl_y + Inches(0.42), exit_w, Inches(2.8), WHITE,
         line=GOLD)

add_multitext(s4, exit_x + Inches(0.2), pl_y + Inches(0.6), exit_w - Inches(0.4), Inches(2.6), [
    {'text': '2035년말 매각 가정', 'size': 12, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '매각가:  ', 'size': 11, 'color': SLATE_700},
    {'text': '3,500억', 'size': 11, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '할인 기간:  10년', 'size': 11, 'color': SLATE_700, 'break_after': True},
    {'text': 'WACC:  8% (할인계수 0.4632)', 'size': 11, 'color': SLATE_700, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},
    {'text': '매각가 PV = ', 'size': 11, 'color': SLATE_700},
    {'text': '1,621.2억', 'size': 11, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '기본 10년 NPV: ', 'size': 11, 'color': SLATE_700},
    {'text': '+8.27억', 'size': 11, 'bold': True, 'color': GREEN, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},
    {'text': '매각 반영 NPV', 'size': 12, 'bold': True, 'color': NAVY, 'break_after': True},
    {'text': '1,629.4억', 'size': 26, 'bold': True, 'color': GOLD},
])

# Bottom: CAPEX 안내
capex_y = pl_y + Inches(3.4)
add_rect(s4, Inches(0.5), capex_y, Inches(12.3), Inches(0.6), NAVY)
add_multitext(s4, Inches(0.75), capex_y + Inches(0.1), Inches(12), Inches(0.4), [
    {'text': '📦  CAPEX 2.5억 (2033년)  ', 'size': 12, 'bold': True, 'color': GOLD},
    {'text': '— 실외 시설 노후 주기 7~10년 → 8년차 권장. 구성: 철골 프레임 보수·도장 1.5억 + 그물망 전면 교체 1.0억',
     'size': 11, 'color': WHITE},
], valign=MSO_ANCHOR.MIDDLE)

add_footer(s4, 4)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5: 결론 및 권고사항
# ══════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, Emu(0), Emu(0), SW, SH, DEEP_NAVY)

# Top: gold accent strip
add_rect(s5, Emu(0), Emu(0), SW, Inches(0.15), GOLD)

# Page heading
add_text(s5, Inches(0.5), Inches(0.45), Inches(0.4), Inches(0.4),
         '04', size=18, bold=True, color=GOLD,
         align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
add_text(s5, Inches(1.05), Inches(0.4), Inches(11), Inches(0.5),
         '결론 및 권고사항', size=24, bold=True, color=WHITE,
         valign=MSO_ANCHOR.MIDDLE)
add_text(s5, Inches(1.05), Inches(0.85), Inches(11), Inches(0.3),
         '5년 vs 10년 시나리오 비교 + 출구전략 분석을 통한 투자 판단',
         size=11, color=ICE_BLUE)

# 시나리오 비교표
comp_y = Inches(1.5)
add_text(s5, Inches(0.5), comp_y, Inches(8), Inches(0.32),
         '■  시나리오 비교 (단위: 억 / %)', size=14, bold=True, color=GOLD)

comp_data = [
    ['시나리오',         '5년 운영',  '10년 운영', '5년+매각',  '10년+매각'],
    ['기본 NPV',         '+0.55',     '+8.27',     '+0.55',     '+8.27'],
    ['매각가',           '-',         '-',         '3,000억',   '3,500억'],
    ['매각 PV',          '-',         '-',         '2,041.7',   '1,621.2'],
    ['최종 NPV',         '+0.55',     '+8.27',     '2,042.3',   '1,629.4'],
    ['IRR',              '8.9%',      '17.6%',     '— ',        '— '],
    ['Payback',          '3.5년',     '3.5년',     '3.5년',     '3.5년'],
    ['투자 적합성',      '적합',      '적합',      '매우 적합', '매우 적합'],
]
add_table(s5, Inches(0.5), comp_y + Inches(0.4), Inches(8), Inches(3.6), comp_data,
          col_widths=[Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6)],
          size=11, header_size=11, bold_first_col=True,
          header_fill=NAVY, header_color=WHITE,
          row_fills=(WHITE, SLATE_BG),
          highlight_last_col=False)

# Right: 권고사항 박스
rec_x = Inches(8.9); rec_w = Inches(4.1)
add_rect(s5, rec_x, comp_y, rec_w, Inches(4.0), NAVY,
         line=GOLD)
# Title bar
add_rect(s5, rec_x, comp_y, rec_w, Inches(0.45), GOLD)
add_text(s5, rec_x + Inches(0.2), comp_y, rec_w - Inches(0.4), Inches(0.45),
         '💡  경영진 권고사항', size=13, bold=True, color=DEEP_NAVY,
         valign=MSO_ANCHOR.MIDDLE)

add_multitext(s5, rec_x + Inches(0.25), comp_y + Inches(0.6), rec_w - Inches(0.5), Inches(3.3), [
    {'text': '1.  ', 'size': 12, 'bold': True, 'color': GOLD},
    {'text': '투자 적합 판정 — ', 'size': 12, 'bold': True, 'color': WHITE},
    {'text': 'NPV 양수(+0.5억/+8.3억), IRR 8.9%(5년)·17.6%(10년) 모두 WACC 8% 상회.',
     'size': 11, 'color': ICE_BLUE, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},

    {'text': '2.  ', 'size': 12, 'bold': True, 'color': GOLD},
    {'text': '5년 매각 시나리오 우위 — ', 'size': 12, 'bold': True, 'color': WHITE},
    {'text': '시간가치 효과로 5년차 매각이 10년차 대비 PV +420억 유리.',
     'size': 11, 'color': ICE_BLUE, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},

    {'text': '3.  ', 'size': 12, 'bold': True, 'color': GOLD},
    {'text': '운영 리스크 — ', 'size': 12, 'bold': True, 'color': WHITE},
    {'text': '2026년 적자(-9.4억) 1회 발생, 2027부터 흑자 전환. 2031~ 회원 감소·비용 인상 이중 압박으로 영업이익 점진 감소.',
     'size': 11, 'color': ICE_BLUE, 'break_after': True},
    {'text': ' ', 'size': 6, 'break_after': True},

    {'text': '4.  ', 'size': 12, 'bold': True, 'color': GOLD},
    {'text': '핵심 액션 — ', 'size': 12, 'bold': True, 'color': WHITE},
    {'text': '오픈 첫해 운전자금 5억 별도 확보, 2032 가격 인상 시점 정확한 시장 모니터링, 2033 CAPEX 2.5억 사전 적립.',
     'size': 11, 'color': ICE_BLUE},
])

# Bottom: 최종 한 줄
bot_y = Inches(5.7)
add_rect(s5, Inches(0.5), bot_y, Inches(12.3), Inches(1.3), NAVY,
         line=GOLD)
add_rect(s5, Inches(0.5), bot_y, Inches(0.15), Inches(1.3), GOLD)

add_multitext(s5, Inches(0.85), bot_y + Inches(0.15), Inches(12), Inches(1.1), [
    {'text': '최종 결론', 'size': 14, 'bold': True, 'color': GOLD, 'break_after': True},
    {'text': '본 사업은 ', 'size': 14, 'color': WHITE},
    {'text': '5년 운영 시 NPV +0.55억 / IRR 8.9%', 'size': 14, 'bold': True, 'color': GOLD},
    {'text': '으로 자본비용(WACC 8%)을 소폭 상회하며, ', 'size': 14, 'color': WHITE},
    {'text': '2030년말 3,000억 매각', 'size': 14, 'bold': True, 'color': GOLD},
    {'text': ' 시 NPV ', 'size': 14, 'color': WHITE},
    {'text': '+2,042억', 'size': 14, 'bold': True, 'color': GOLD},
    {'text': '으로 매우 우수한 투자수익이 확보됨.', 'size': 14, 'color': WHITE, 'break_after': True},
    {'text': '재오픈 후 5년 안정 운영 → 최적 시점 매각 전략 권장.', 'size': 13, 'italic': True, 'color': ICE_BLUE},
], line_spacing=1.3)

# Footer
add_text(s5, Inches(0.5), SH - Inches(0.35), Inches(8), Inches(0.3),
         '등촌골프연습장 사업성 분석 · 경영진 보고용 · 2026.05',
         size=9, color=SLATE_300, valign=MSO_ANCHOR.MIDDLE)
add_text(s5, SW - Inches(1.5), SH - Inches(0.35), Inches(1.2), Inches(0.3),
         '5 / 5', size=9, color=ICE_BLUE, bold=True,
         align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)

# ──────────────────────────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────────────────────────
import shutil
output_path = '등촌골프연습장_경영진보고_5장.pptx'
prs.save(output_path)
print(f'✓ Created: {output_path}')

# 데스크탑에도 복사 (사용자 접근 편의)
desktop_korean = r'C:\Users\win\Desktop\등촌골프연습장_경영진보고.pptx'
desktop_english = r'C:\Users\win\Desktop\Deungchon_Executive_Report.pptx'
try:
    shutil.copyfile(output_path, desktop_korean)
    print(f'✓ Desktop (한글): {desktop_korean}')
except Exception as e:
    print(f'Desktop 한글 복사 실패: {e}')
try:
    shutil.copyfile(output_path, desktop_english)
    print(f'✓ Desktop (영문): {desktop_english}')
except Exception as e:
    print(f'Desktop 영문 복사 실패: {e}')
